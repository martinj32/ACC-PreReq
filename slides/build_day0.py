"""
ACC Day 0 — AI / Agentics Basics Slide Deck Builder — Army Cyber Dark Theme
Source: docs/day0/ai-agentics-basics.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_summary_table(prs, title, subtitle, headers, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, subtitle)
    _footer(slide)

    col_x = [L, Inches(3.5), Inches(8.1)]
    col_w = [Inches(2.9), Inches(4.4), Inches(5.0)]
    ht = Inches(1.52)
    rh = Inches(0.42)

    hdr = _rect(slide, L, ht, R, rh, ELECTRIC_BLUE)
    _glow_blue(hdr, 4)
    for i, h in enumerate(headers):
        _tb(slide, h, col_x[i] + Inches(0.1), ht + Inches(0.06), col_w[i], rh,
            sz=11, bold=True, color=CYBER_BLACK)

    for ri, row in enumerate(rows):
        top = ht + rh + ri * Inches(0.83)
        _rect(slide, L, top, R, Inches(0.83), DARK_NAVY if ri % 2 == 0 else NAVY_MID)
        bar = _rect(slide, L, top, Pt(3), Inches(0.83), ELECTRIC_BLUE)
        _glow_blue(bar, 3)
        colors = [ELECTRIC_BLUE, WHITE, LIGHT_GRAY]
        for ci, (val, color) in enumerate(zip(row, colors)):
            _tb(slide, val, col_x[ci] + Inches(0.12), top + Inches(0.08),
                col_w[ci] - Inches(0.15), Inches(0.8),
                sz=11, bold=(ci == 0), color=color)


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # -----------------------------------------------------------------------
    # TITLE
    # -----------------------------------------------------------------------
    add_title_slide(prs,
        "AI / Agentics Basics\nDay 0 — Pre-Course",
        "PREREQUISITE BLOCK  |  MULTI-WEEK DAILY REPS",
        "[IMAGE: Soldier at a terminal workstation reviewing AI tool output — zero to capable in four phases]")

    # -----------------------------------------------------------------------
    # OVERVIEW
    # -----------------------------------------------------------------------
    add_overview_slide(prs,
        "A student who has never touched a terminal walks out with a working mental model of the LLM, "
        "deliberate prompting skills, command-line fluency, and a supervisor mindset — ready to run an agent without freezing.",
        "Multi-week, short daily reps (four phases)",
        [
            "Understand what an LLM actually is and how it generates output",
            "Know the four failure modes — and produce one with your own hands",
            "Prompt deliberately: role, context, example, output spec",
            "Navigate and act in a real filesystem from the command line",
            "Understand the chatbot-vs-agent distinction and what access it implies",
            "Apply the supervisor mindset: delegate clearly, verify the work, own the outcome",
        ])

    # -----------------------------------------------------------------------
    # PHASE 1 HEADER
    # -----------------------------------------------------------------------
    add_section_header(prs, "01", "Know What You\nAre Talking To",
        "Start in the tool they already use. Convert intuition into method — mental model, failure modes, deliberate prompting.",
        "[IMAGE: Soldier studying a technical field manual — converting experience into doctrine]")

    # -----------------------------------------------------------------------
    # MODULE 0.1 — What an LLM Actually Is
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.1", "What an LLM\nActually Is",
        "An LLM predicts text from learned patterns — not rules, not memory, not live data — and that single fact explains both its power and its failure modes.",
        "[IMAGE: Soldier at a radio set — receiving signals, interpreting patterns, generating responses]")

    add_concept_slide(prs, "How an LLM Works", [
        ("The model reads a sequence of text and predicts what comes next — one chunk at a time.", False),
        ("No lookup table. No truth check. No memory between sessions.", False),
        ("Given everything written so far: what is the statistically most likely next piece?", False),
        ("Trained, not programmed.", False),
        ("No one wrote rules. The model was trained on enough text that patterns learned into its weights.", True),
        ("Upside: generalizes to almost anything. Failure mode: bad training data means bad patterns.", True),
        ("The LLM is the engine — a brain in a jar. Fluent output is not the same as correct output.", False),
    ], note="Swap the verb: the model does not 'know,' 'want,' 'think,' or 'lie.' It predicts. That is the only accurate verb. Catch yourself using the others and swap it in.")

    add_check_on_learning(prs,
        "The model just gave you a confident-sounding answer.\n\n"
        "What would it take to verify it? Where would you check?\n\n"
        "A language model was never given a rule that says 'Paris is the capital of France.' "
        "How does it produce that answer?")

    add_hands_on(prs, "What an LLM Actually Is", [
        "Ask your chatbot something you already know well enough to spot an error. Read the output.",
        "Ask yourself: 'How would I verify this?' Where would you check?",
        "Ask it to continue a sentence you leave unfinished. Watch the completion shift with small word changes.",
        "You are watching prediction happen in real time — not lookup, not rules, not memory.",
    ], "[IMAGE: Soldier testing equipment against known specifications — spot-checking output against ground truth]")

    add_readiness_check(prs, [
        "I can state in one sentence what an LLM does",
        "I can explain 'trained not programmed' in plain terms",
        "I know that confident output and correct output are not the same thing",
        "I have watched the model predict — not just heard about it",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.2 — Tokens and Context
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.2", "Tokens\nand Context",
        "The model reads and writes in tokens — roughly 0.75 of an English word each — and can only hold a fixed number at once; fill that limit and earlier content starts to fall out.",
        "[IMAGE: Tactical whiteboard running out of space — overlapping notes crowding out earlier entries]")

    add_concept_slide(prs, "Tokens and the Context Window", [
        ("The model does not read words — it reads chunks of text called tokens.", False),
        ("In English prose, 1 token is roughly three-quarters of a word.", False),
        ("Code and non-English text tokenize less efficiently — more tokens per word.", False),
        ("Every model has a hard ceiling: the context window.", False),
        ("Think of it as a whiteboard with finite space — new writing crowds out old writing once it fills.", True),
        ("Current approximate sizes: Claude ~200K tokens | GPT-4o ~128K | Gemini 2.0 ~1M+", True),
        ("Bigger is not always better — a full window degrades quality in the middle.", False),
    ], note="Verify context window sizes before teaching — they change with releases. The instinct 'shorter and focused beats long and sprawling' is the deliverable, not the numbers.")

    add_concept_slide(prs, "What a Full Window Looks Like", [
        ("The model starts hedging.", False),
        ("It contradicts something it stated confidently ten turns ago.", False),
        ("It forgets a constraint you set at the start.", False),
        ("That is not a bug — that is the whiteboard running out of space.", False),
        ("Two situations that warrant starting a fresh chat:", False),
        ("The task has changed direction and earlier turns are dead weight.", True),
        ("The model is drifting — hedging on things it stated confidently before.", True),
    ], note="Critical instructions buried mid-conversation may be partially ignored. Models attend better to the beginning and end of their context window than to the middle.")

    add_check_on_learning(prs,
        "You set an important constraint at the start of a long conversation. "
        "Ten exchanges later, the model seems to have forgotten it.\n\n"
        "What happened, and what would you do differently next time?\n\n"
        "You are in a long chat and the model starts contradicting instructions you gave at the start. "
        "What is the most likely cause?")

    add_hands_on(prs, "Tokens and Context", [
        "Open a fresh chat and paste several paragraphs of dense text.",
        "Ask the model to summarize just the first paragraph.",
        "In the same chat, ask it about something from the end of what you pasted.",
        "Start a new chat and ask the same question. Compare the quality.",
        "You are probing the window's edges — watching it degrade in real time.",
    ], "[IMAGE: Soldier reviewing a long operation log — noting what has dropped off the whiteboard]")

    add_readiness_check(prs, [
        "I can define a token in plain terms",
        "I can explain what a context window is and what happens when it fills",
        "I know the two situations that warrant starting a fresh chat",
        "I have observed the window effect — not just read about it",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.3 — How LLMs Fail
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.3", "How LLMs\nFail",
        "The model will state false things with total confidence, give different answers to the same question, and know nothing past its training cutoff — none of this is a malfunction.",
        "[IMAGE: Soldier double-checking map coordinates before calling in a grid — verifying before acting]")

    add_concept_slide(prs, "Four Failure Modes", [
        ("Hallucination — the model has no truth-checking step.", False),
        ("Generates tokens statistically likely to follow context. Cannot distinguish accurate from plausible.", True),
        ("Confident output and correct output are entirely unrelated. This is expected behavior.", True),
        ("Confident-Wrong — hallucination is not always dramatic.", False),
        ("Subtly wrong dates, statistics, citations — equally confident tone regardless of accuracy.", True),
        ("Never use the model's confidence as a signal of accuracy.", True),
    ], note="The goal is calibrated trust, not distrust. The model is extraordinarily capable and it hallucinates. Both are true simultaneously. If students leave deciding it is useless, the lesson failed.")

    add_concept_slide(prs, "Failure Modes 3 and 4", [
        ("Nondeterminism — the model uses a temperature parameter that introduces variation.", False),
        ("Same prompt, different run, different result — by design.", True),
        ("Do not treat one output as 'the answer' for high-stakes work. Run important prompts more than once.", True),
        ("Knowledge Cutoff — the model was trained on data up to a point in time.", False),
        ("It does not know what happened after that cutoff.", True),
        ("Ask about a recent event and it may refuse, guess, or confidently describe something that did not happen.", True),
        ("Build the verification reflex now: anything that matters gets checked before you act on it.", False),
    ], note="All four failure modes look identical from the outside. Verify anything time-sensitive regardless of how confidently the model presents it.")

    add_check_on_learning(prs,
        "You just watched the model invent a source.\n\n"
        "What does that mean for the next time it gives you a fact you have not heard before?\n\n"
        "You run the same prompt twice in two separate chats and get different answers. "
        "What is the most accurate explanation?")

    add_hands_on(prs, "How LLMs Fail", [
        "Ask the model for 5 peer-reviewed sources on a specific narrow topic — author names, journal names, publication years.",
        "Pick one citation and try to verify it exists. Check the DOI if it provides one.",
        "Ask the same question in a new chat. Compare the answers.",
        "Ask about something you know happened recently. Watch how it responds.",
        "You have just produced hallucination, confident-wrong, nondeterminism, and cutoff — all four modes.",
    ], "[IMAGE: Analyst comparing two documents side by side — source document vs. model output]")

    add_readiness_check(prs, [
        "I can name the four core failure modes",
        "I have produced a hallucination with my own hands",
        "I understand that confidence and correctness are unrelated in model output",
        "I know why this module is the reason a human stays in the loop",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.4 — Deliberate Prompting
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.4", "Deliberate\nPrompting",
        "Converting prompting-by-feel into method means telling the model who to be, giving it context, showing it an example, and saying what good output looks like.",
        "[IMAGE: Soldier issuing a clear five-paragraph OPORD — structured brief, no ambiguity]")

    add_concept_slide(prs, "The Four Elements of a Deliberate Prompt", [
        ("Without structure, the model must guess: who you are, what you need, what 'good' looks like.", False),
        ("Guessing produces generic output. Four elements close that gap:", False),
        ("Role — tell the model who to be. Sets the frame for every response.", True),
        ("Context — what it needs to know that it cannot infer: task, audience, constraints.", True),
        ("Example — show it what good output looks like before asking. One strong example beats three paragraphs.", True),
        ("Output spec — format, length, and tone you want.", True),
        ("The iterative habit matters more than the perfect prompt.", False),
    ], note="The template trap: the four elements are a structure, not a script to recite. The durable skill is clarity — giving the model what it needs to not guess. How you deliver that varies by task.")

    add_example_slide(prs, "Before and After: Prompting in Practice",
        "Same model. Different brief. Different result.",
        [
            "WEAK PROMPT:",
            '  "Write me a summary."',
            "  Model must guess: summary of what? For whom? What length? What format?",
            "  Result: generic, too long, misses the point.",
            "",
            "STRONG PROMPT:",
            '  "You are summarizing this for a senior leader who has two minutes.',
            '   Pull out the three most important points. Use bullet points.',
            '   No jargon. Max 100 words."',
            "",
            "If you are stuck: ask the model to interview you.",
            '  Type: "I need help with [task]. Ask me the questions you need answered before you start."',
        ],
        "[IMAGE: Two mission briefs side by side — one a wall of text, one a clean OPORD format]")

    add_check_on_learning(prs,
        "What specifically changed between your first output and your second?\n\n"
        "Which of the four elements made the biggest difference for your task?\n\n"
        "You ask the model to 'write a report.' It produces something generic and too long. "
        "What is the most effective fix?")

    add_hands_on(prs, "Deliberate Prompting", [
        "Pick a real task you need help with.",
        "Write a one-line version of the request. Submit it. Save the output.",
        "Now add role, context, an example of what good looks like, and an output spec. Submit again.",
        "Compare the two outputs side by side.",
        "Which of the four elements made the biggest difference for your task?",
    ], "[IMAGE: Soldier editing a draft operations order — tightening the mission statement line by line]")

    add_readiness_check(prs, [
        "I can name the four elements of a deliberate prompt",
        "I have run a before/after comparison with my own real task",
        "I treat prompting as a conversation, not a one-shot command",
        "I know the difference between prompting by feel and prompting by structure",
    ])

    add_section_summary(prs, "Phase 1 — Know What You Are Talking To", [
        "The LLM predicts — it does not know, think, or lie. Swap the verb.",
        "Four failure modes: hallucination, confident-wrong, nondeterminism, knowledge cutoff.",
        "Context windows are finite — fill them and the model loses earlier content.",
        "Deliberate prompting: role, context, example, output spec. Structure beats feel every time.",
    ])

    # -----------------------------------------------------------------------
    # PHASE 2 HEADER
    # -----------------------------------------------------------------------
    add_section_header(prs, "02", "The Computer\nUnderneath",
        "Genuinely new material. No AI here — plain computer literacy. The terminal modules in Phase 3 all depend on being able to think about where things are.",
        "[IMAGE: Soldier studying a topographic map — building the mental terrain model before movement]")

    # -----------------------------------------------------------------------
    # MODULE 0.5 — Files, Folders, and the Tree
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.5", "Files, Folders,\nand the Tree",
        "Everything on a computer lives in a file, files live in folders, folders nest into a tree, and the address of any file is the path you walk to reach it.",
        "[IMAGE: Soldier navigating a sectored map — understanding the terrain before moving through it]")

    add_concept_slide(prs, "The File Tree", [
        ("A file holds data: a document, an image, a script.", False),
        ("A folder (directory) holds files and other folders. Folders nest into a tree.", False),
        ("Every file on the machine has exactly one location in that tree.", False),
        ("The path is the route through the tree to reach a file.", False),
        ("Windows: C:\\Users\\YourName\\Documents\\report.txt", True),
        ("Mac/Linux: /home/yourname/documents/report.txt", True),
        ("Most apps hide this. The tree exists whether you see it or not.", False),
    ], note="Many casual users save files and find them through search or recent items — they have genuinely never navigated the folder structure. Verify by asking: 'Can you find a file you saved last week without using search?'")

    add_concept_slide(prs, "File Extensions — Turn Them On Now", [
        ("A file extension is the suffix after the dot: .docx, .txt, .pdf, .py", False),
        ("By default, Windows hides extensions. Turn them on before anything else.", False),
        ("File Explorer → View → Show → File name extensions.", False),
        ("You will need to see extensions throughout this course.", False),
        ("Why it matters: the difference between .txt and .docx is the whole next module.", False),
    ], note="Do not skip the extension step. It prevents confusion in Module 0.6 where the difference between .txt and .docx is the entire lesson.")

    add_check_on_learning(prs,
        "You just moved a file by dragging.\n\n"
        "Write out the path to where it now lives.\n\n"
        "Could you give someone else those directions and have them find it?\n\n"
        "What is a file path?")

    add_hands_on(prs, "Files, Folders, and the Tree", [
        "Open File Explorer (Windows) or Finder (Mac).",
        "Navigate from your home folder down to Downloads by clicking — do not use search.",
        "Find a file. Read its full path in the address bar.",
        "Create a new folder on your Desktop. Name it 'practice'.",
        "Move any file into that folder by dragging. You are walking the tree — Phase 3 is the same moves, typed.",
    ], "[IMAGE: Soldier navigating a route on a physical map — same route, different medium than digital]")

    add_readiness_check(prs, [
        "I can describe a file, a folder, and how they nest into a tree",
        "I can read a path as a route through that tree",
        "I have navigated to a specific file without using search",
        "File extensions are now visible on my machine",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.6 — Plaintext vs Rich Text
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.6", "Plaintext vs\nRich Text",
        "A plaintext file is just characters — no hidden formatting, no metadata — and that is exactly what models and tools want.",
        "[IMAGE: Soldier reading an unformatted signal message — raw, unambiguous, machine-readable]")

    add_concept_slide(prs, "Plaintext vs Rich Text", [
        ("Plaintext: text, nothing else. What you see is every character the file contains.", False),
        ("No hidden formatting. No metadata. What you see is what the machine reads.", False),
        ("Common plaintext: .txt  .md  .csv  .json  .py", True),
        ("Rich text files (.docx, .xlsx, .pdf) store formatting instructions alongside the text.", False),
        ("Bold, font size, margins, revision history — all packed in alongside the words.", True),
        ("When you paste from Word into a terminal or Markdown, hidden characters come with it.", True),
        ("Smart quotes, em-dashes, non-breaking spaces — invisible and often breaking.", False),
    ], note="Students paste content from Word documents into AI tools and wonder why things break. The answer is hidden characters. This is the diagnosis for a category of problems they will hit repeatedly.")

    add_check_on_learning(prs,
        "You paste a block of text from a Word document into your AI chatbot and the formatting looks off.\n\n"
        "What is the most likely cause, and what would you do instead?\n\n"
        "Why do models and command-line tools prefer plaintext over rich-text formats like .docx?")

    add_hands_on(prs, "Plaintext vs Rich Text", [
        "Open Notepad (Windows) and type a few sentences. Save as test.txt.",
        "Open the same file in VS Code. Compare how it looks.",
        "Now open any .docx file in VS Code (not Word). Look at what is actually inside it.",
        "You are seeing the difference between 'what a human reads' and 'what a machine reads.'",
    ], "[IMAGE: Two documents side by side — one clean plaintext, one Word file opened raw showing markup]")

    add_readiness_check(prs, [
        "I can distinguish plaintext from rich text and give an example of each",
        "I understand why pasting from Word can silently break things",
        "I know the common plaintext extensions",
        "I have seen what a .docx file actually looks like underneath",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.7 — Looking at Files in a Real Editor
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.7", "Looking at Files\nin a Real Editor",
        "A code editor is a better text viewer — install one, open a folder, read a file, close it without fear; that is the whole job today.",
        "[IMAGE: Soldier setting up a communications suite — installing equipment before the mission]")

    add_concept_slide(prs, "VS Code — Three Things to Know Today", [
        ("A code editor is not an IDE. It is a text viewer with good defaults.", False),
        ("VS Code: free, runs on Windows and Mac, handles every file type in this course.", False),
        ("File tree panel (left side) — the same folder structure from Module 0.5, now labeled.", False),
        ("Editor panel (right side) — file content.", False),
        ("Markdown preview — click the split-square icon to render a .md file.", False),
        ("Scope today: open, view, close.", False),
        ("Do not install extensions, change settings, or start editing files.", False),
    ], note="Editors are intimidating by appearance. The goal is comfort at the front door, not a full tour. A broken install on one machine should not stall the room — Notepad++ is a lighter backup.")

    add_check_on_learning(prs,
        "You are looking at the same Markdown file in two panes — raw on the left, rendered on the right.\n\n"
        "Which version does the machine work with? Which version does the human read?\n\n"
        "What is the main reason to use a code editor instead of Word for working with plaintext files?")

    add_hands_on(prs, "Looking at Files in a Real Editor", [
        "Download and install VS Code if you have not already.",
        "Go to File → Open Folder. Open the 'practice' folder from Module 0.5.",
        "Click the .txt file you have in there. Read it in the editor panel.",
        "Find a .md file anywhere on your machine. Open it. Toggle the preview (Ctrl+Shift+V).",
    ], "[IMAGE: Soldier using two viewing tools on the same map — one physical, one digital]")

    add_readiness_check(prs, [
        "VS Code (or equivalent) is installed and opens without errors",
        "I can open a folder and navigate its tree in the editor",
        "I have opened a Markdown file and toggled the rendered preview",
        "I did not change any settings or install any extensions",
    ])

    add_section_summary(prs, "Phase 2 — The Computer Underneath", [
        "Files nest in folders that nest in a tree — every file has one path.",
        "Plaintext is honest; rich text hides things. Tools want plaintext.",
        "The code editor is the plaintext viewer — open, view, close without fear.",
        "File extensions are visible now. They matter for every module ahead.",
    ])

    # -----------------------------------------------------------------------
    # PHASE 3 HEADER
    # -----------------------------------------------------------------------
    add_section_header(prs, "03", "The Terminal",
        "The spine of the course. Maximum scaffolding, maximum reps. Every rep is framed as: this is how the agent will move around your computer for you.",
        "[IMAGE: Soldier at a communications terminal — direct interface, typed commands, no menu in between]")

    # -----------------------------------------------------------------------
    # MODULE 0.8 — What the Terminal Is (and Why It Is Safe)
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.8", "What the Terminal Is\n(and Why It Is Safe)",
        "The terminal is a typed way to give the computer the same instructions you used to give by clicking — it is not hacking, it feels alien for about an hour before it feels normal.",
        "[IMAGE: Soldier first day at a new communications system — unfamiliar but not dangerous]")

    add_concept_slide(prs, "The Terminal — Direct Comms with the Machine", [
        ("A command is a typed instruction. Same as a click — it tells the computer to do something.", False),
        ("The difference: a click is constrained by what a menu shows you.", False),
        ("A command can express anything the machine understands.", False),
        ("Clicking is talking through an interpreter. The command line is speaking the language directly.", False),
        ("The prompt line is where your typing goes. $ (Mac/Linux) or > (Windows PowerShell) means ready.", False),
        ("The commands in this module and the next two are read-only and harmless.", False),
    ], note="Address the bait-and-switch objection out loud before a student says it: 'Some of you are wondering why an AI course spent two weeks on the command line.' The agent acts through the terminal. You cannot supervise something you cannot read.")

    add_check_on_learning(prs,
        "You typed a command and saw a response.\n\n"
        "What is the difference between that and clicking a button in an app? What stayed the same?\n\n"
        "What does a terminal actually do?")

    add_hands_on(prs, "What the Terminal Is", [
        "Open the terminal on your machine. (Windows: search PowerShell. Mac: search Terminal.)",
        "Sit with the blinking cursor for ten seconds. Name the feeling. Then continue.",
        "Type the command to print the current date. (Mac/Linux: date  |  Windows: Get-Date). Press Enter.",
        "Close the terminal and open it again.",
        "You opened it, ran something, and closed it. The anxiety has a smaller surface area now.",
    ], "[IMAGE: Soldier completing first radio check — awkward at first, routine after a few reps]")

    add_readiness_check(prs, [
        "I can open the terminal on my machine",
        "I ran one command and read its output",
        "I can find the prompt line and know where my typing goes",
        "The anxiety is smaller than it was at the start of this module",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.9 — Navigating: Where Am I, What Is Here
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.9", "Navigating:\nWhere Am I, What Is Here",
        "Three commands cover most of navigation — where am I, what is here, move — and mastering them replaces every click from Module 0.5.",
        "[IMAGE: Soldier performing land navigation — map, compass, grid coordinates]")

    add_concept_slide(prs, "The Three Navigation Commands", [
        ("This is land navigation: where am I, what is around me, how do I move.", False),
        ("pwd — Print Working Directory. Your current location in the tree.", False),
        ("      Works on: Mac, Linux, WSL, PowerShell", True),
        ("ls — List. Everything in the current folder.", False),
        ("      Add -la (Mac/WSL) for details including hidden files.", True),
        ("cd foldername — Change Directory. Moves down into a folder.", False),
        ("cd .. — Move up one level.", False),
    ], note="Same three commands, many short repetitions, spread across several days. Do not introduce file creation yet. One capability at a time. Reps build muscle memory; coverage builds confusion.")

    add_example_slide(prs, "Walking the Tree — Typed Instead of Clicked",
        "Same moves as Module 0.5, typed instead of clicked",
        [
            "pwd              # See where you are",
            "ls               # See what is here",
            "cd Documents     # Move into Documents",
            "pwd              # Confirm you moved",
            "ls               # See what is in Documents",
            "cd ..            # Move back up one level",
            "pwd              # Confirm you are back",
            "",
            "pwd is your reset button.",
            "  Any time you feel lost: type pwd.",
            "  It tells you exactly where you are.",
            "  You can never be permanently lost if you can always ask 'where am I.'",
        ],
        "[IMAGE: Soldier confirming grid position on a map — reset to known location before continuing]")

    add_check_on_learning(prs,
        "You typed 'cd Documents' and then 'pwd'.\n\n"
        "What does pwd tell you, and how does it prove you moved?\n\n"
        "You are in the terminal and have no idea where you are in the filesystem. "
        "What is the first command you run?")

    add_hands_on(prs, "Navigating the Filesystem", [
        "Open the terminal. Run pwd. Read the path. Find that same location in File Explorer.",
        "Run ls. Identify three items by name.",
        "Run cd to move into one of them.",
        "Run pwd and ls again.",
        "Run cd .. to go back up. Repeat this cycle until it feels routine.",
    ], "[IMAGE: Soldier moving through a sectored area — confirming position at each checkpoint]")

    add_readiness_check(prs, [
        "I can print my current location with pwd",
        "I can list folder contents with ls",
        "I can move into and back out of folders with cd and cd ..",
        "I have walked the same tree by clicking (0.5) and by typing (this module)",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.10 — Acting: Make, Move, Copy
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.10", "Acting: Make,\nMove, Copy",
        "Four commands cover most day-to-day file work — make a folder, make a file, copy, move — the same actions from Module 0.5, now typed.",
        "[IMAGE: Soldier preparing equipment — creating, organizing, moving gear to the right location]")

    add_concept_slide(prs, "The Four File Action Commands", [
        ("Every action you take, you confirm. You do not assume the round landed — you check the target.", False),
        ("mkdir foldername — Make a new folder.", False),
        ("touch filename.txt (Mac/WSL) / New-Item filename.txt (PowerShell) — Make a new empty file.", False),
        ("cp source destination — Copy a file.", False),
        ("mv source destination — Move or rename a file.", False),
        ("No deletion yet. Deletion is irreversible at the command line — no Recycle Bin.", False),
        ("Typos create new files. Run ls after every action and read the output.", False),
    ], note="'Trust the action, not the narration' is a principle that runs through all agentic work. Build it here, at the command line, before the agent is involved.")

    add_check_on_learning(prs,
        "You ran 'mv notes.txt notes-final.txt'. What does ls show you?\n\n"
        "What would it look like if you had made a typo in the destination name?\n\n"
        "You run 'mv report.txt final-report.txt' and then ls. What are you checking for?")

    add_hands_on(prs, "Make, Move, Copy", [
        "Make a folder called 'phase3-practice'. Navigate into it.",
        "Create a file called day1.txt.",
        "Copy it to day1-backup.txt.",
        "Rename day1.txt to day1-v1.txt.",
        "Run ls after each step to verify the change happened.",
        "Build the verify-after-acting habit here — before any agent is involved.",
    ], "[IMAGE: Soldier verifying equipment load after staging — checking each item before mission launch]")

    add_readiness_check(prs, [
        "I can create a folder and a file from the command line",
        "I can copy and move/rename files",
        "I run ls after every action to verify it happened",
        "I have not attempted to delete anything yet",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.11 — Paths, Tab-Completion, and the Up-Arrow
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.11", "Paths, Tab-Completion,\nand the Up-Arrow",
        "A path is the full address of a file; tab-completion finishes names so you stop fat-fingering them; the up-arrow recalls your last command.",
        "[IMAGE: Soldier using grid coordinates precisely — exactness prevents the round from landing in the wrong place]")

    add_concept_slide(prs, "Absolute vs Relative Paths", [
        ("An absolute path starts from the root and gives the full route — works from anywhere.", False),
        ("Windows: C:\\Users\\YourName\\Documents\\report.txt", True),
        ("Mac/Linux: /home/yourname/documents/report.txt", True),
        ("A relative path starts from wherever you currently are — only works in the right starting position.", False),
        ("documents/report.txt  or  ../otherfolder/file.txt", True),
        ("A path is a grid coordinate: precise, unambiguous, and unforgiving of a wrong character.", False),
        ("A space or capital letter in the wrong place either fails or hits the wrong target.", False),
    ], note="The terminal's unforgiving exactness is both the lesson and the frustration. A student annoyed that a capital letter breaks a path is learning the right thing. Tab-completion makes this structurally impossible for the completed portion.")

    add_concept_slide(prs, "Tab-Completion and the Up-Arrow", [
        ("Tab-completion: start typing a folder or file name and press Tab — the terminal finishes it.", False),
        ("If there is more than one match, press Tab twice to see the options.", False),
        ("This is not a shortcut — it is the standard. Use it on every path, every time.", False),
        ("A typo in a path is invisible until something breaks. Tab-completion makes typos impossible for the completed portion.", False),
        ("The up-arrow: press to cycle through previous commands.", False),
        ("Edit a recalled command and re-run it instead of retyping from scratch.", False),
    ], note="Never finish typing a long path by hand if tab-completion can do it. This is the single biggest accuracy and morale win in the terminal.")

    add_check_on_learning(prs,
        "You are deep in a nested folder and need to copy a file to a location three levels up.\n\n"
        "Would you use an absolute or a relative path? Why?\n\n"
        "Why does tab-completion reduce errors in terminal commands?")

    add_hands_on(prs, "Paths, Tab-Completion, and the Up-Arrow", [
        "Navigate to a folder using the full absolute path — type it out by hand once to feel the length.",
        "Navigate to the same folder using tab-completion. Press Tab after each segment.",
        "Run ls in a folder with several items. Start typing one item's name, press Tab, watch it complete.",
        "Run a command. Press the up-arrow. Edit one character and re-run it.",
    ], "[IMAGE: Soldier using precise grid coordinates on a map — every digit matters]")

    add_readiness_check(prs, [
        "I can write both an absolute and a relative path to the same file",
        "I use tab-completion on every path — it is now a habit, not an option",
        "I can recall and edit previous commands with the up-arrow",
        "I understand why exactness (spaces, capitalization) matters in a path",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.12 — Commands, Flags, and Knowing When to Stop
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.12", "Commands, Flags,\nand Knowing When to Stop",
        "Most commands take flags that change what they do; --help shows all available flags without memorization; Ctrl+C stops anything that will not quit.",
        "[IMAGE: Soldier adjusting radio settings with flags/options — same radio, different mission profile]")

    add_concept_slide(prs, "Flags and How to Find Them", [
        ("Flags are options you add to a command to change its behavior.", False),
        ("Short form: -l  |  Long form: --help", False),
        ("ls          — list files, default view", True),
        ("ls -l       — list files, detailed view", True),
        ("ls -la      — list files, detailed, including hidden files", True),
        ("Same command. Different flags. Different output.", False),
        ("You do not memorize flags. You look them up: ls --help  |  Get-Help ls", False),
    ], note="The concept of a flag and how to look one up is the lesson. Memorizing specific flags is the wrong target. The skill is knowing flags exist and knowing how to find the one you need.")

    add_concept_slide(prs, "The Prompt and Ctrl+C", [
        ("The prompt is the signal. When a command finishes, the prompt returns.", False),
        ("If the cursor blinks with no prompt and nothing is happening, the command is still running.", False),
        ("Ctrl+C stops a running command and returns the prompt.", False),
        ("It does not delete anything — it stops execution and hands control back to you.", False),
        ("Students who know they can stop something are less afraid to start it.", False),
        ("This is a control mechanism, not an emergency procedure.", False),
    ], note="Flags prepare students for CLI tools in agentic environments. Ctrl+C is a control mechanism. Name the forward connection: these are the same tools the agent uses.")

    add_check_on_learning(prs,
        "You run a command and nothing happens for thirty seconds — the cursor just blinks.\n\n"
        "What are the two possibilities, and how do you handle each one?\n\n"
        "You need to run ls but want to see hidden files. You do not know the flag. What do you do?")

    add_hands_on(prs, "Commands, Flags, and Knowing When to Stop", [
        "Run ls --help (Mac/Linux) or Get-Help ls (PowerShell). Read the first ten lines without panic.",
        "Find the flag that shows hidden files. Run ls with that flag.",
        "Run ping google.com (ping -t google.com on Windows). Let it run for five seconds.",
        "Press Ctrl+C. Confirm the prompt returned.",
        "You just demonstrated: flags change behavior, and you can stop anything you start.",
    ], "[IMAGE: Soldier stopping a radio transmission — controlled interruption, returning to standby]")

    add_readiness_check(prs, [
        "I can add a flag to a command and observe the difference in output",
        "I can use --help to find a flag I do not know",
        "I know the prompt returning means a command finished",
        "I can stop a running command with Ctrl+C without fear",
    ])

    add_section_summary(prs, "Phase 3 — The Terminal", [
        "pwd, ls, cd — where am I, what is here, move. Know them cold.",
        "mkdir, touch/New-Item, cp, mv — make, copy, move. Run ls after every action.",
        "Tab-completion is not a shortcut — it is standard operating procedure.",
        "Flags change behavior. --help finds them. Ctrl+C stops anything.",
    ])

    # -----------------------------------------------------------------------
    # PHASE 4 HEADER
    # -----------------------------------------------------------------------
    add_section_header(prs, "04", "From Chatbot\nto Agent",
        "The bridge that makes Phases 2 and 3 pay off. Every command you learned in Phase 3 is what the agent runs when it works in your filesystem.",
        "[IMAGE: Soldier handing off a mission to a capable junior — clear brief, defined boundaries, active oversight]")

    # -----------------------------------------------------------------------
    # MODULE 0.13 — Chatbot vs Agent
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.13", "Chatbot\nvs Agent",
        "A chatbot gives advice; an agent takes action — the difference is read, write, and execute access to the real files and terminal the student just learned.",
        "[IMAGE: Soldier comparing advisory role to action role — one briefs, one executes]")

    add_concept_slide(prs, "The One-Line Difference", [
        ("Ask a chatbot to rename a folder: it tells you the command.", False),
        ("Ask an agent: it renames the folder.", False),
        ("That is the entire difference. The agent has access.", False),
        ("Three levels of access:", False),
        ("Read — the agent can look at your files.", True),
        ("Write — the agent can create and modify your files.", True),
        ("Execute — the agent can run commands in your terminal.", True),
        ("Engine + Harness + Operator. Everything in advanced agentic work stacks on this primitive.", False),
    ], note="Make the connection explicit: every command you learned in Phase 3 — pwd, ls, cd, mkdir, mv — that is what the agent runs when it works in your filesystem. You learned to navigate the terrain so you can supervise someone else navigating it.")

    add_concept_slide(prs, "Read, Write, Execute Is a Lot of Trust", [
        ("This is not abstract.", False),
        ("An agent with write and execute access can create, modify, or delete files on your real machine.", False),
        ("The LLM is the engine. A harness with tools gives it hands.", False),
        ("You are the operator who points it at the right problem and pulls the plug when it goes wrong.", False),
        ("Module 0.16 sets the bright line for what that access means.", False),
        ("Do not skip it.", False),
    ], note="The student who does not understand the terminal cannot supervise an agent that is using it. That is why Phases 2 and 3 existed.")

    add_check_on_learning(prs,
        "You just asked a chatbot to rename a folder.\n\n"
        "What would an agent do differently — and what are you trusting it with when it does?\n\n"
        "What is the key difference between a chatbot and an agent?")

    add_hands_on(prs, "Chatbot vs Agent", [
        "Open your chatbot. Ask it: 'Rename the folder project to project-v1.' Read the response.",
        "Identify: did it rename the folder, or tell you how? That gap is the distinction.",
        "Think about what would need to change for it to perform the action instead of describe it.",
        "If you have Claude Code available, ask it the same thing in a project folder. Compare.",
    ], "[IMAGE: Soldier at terminal comparing two different outputs — advisory vs. action]")

    add_readiness_check(prs, [
        "I can state the one-line difference between a chatbot and an agent",
        "I can explain what read, write, and execute access means in practice",
        "I understand why Phases 2 and 3 were prerequisites for this module",
        "I can name all three parts of the engine-harness-operator model",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.14 — Why Version Control Exists
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.14", "Why Version\nControl Exists",
        "When an agent can change your files, you want a logbook of every change — what changed, when, and why, with the ability to go back.",
        "[IMAGE: NCO reviewing the duty logbook — every entry recorded, every action accountable]")

    add_concept_slide(prs, "Version Control as the Duty Logbook", [
        ("Every time a file changes, version control records a snapshot: what changed, when, and why.", False),
        ("Stored locally on your machine. Every snapshot can be rewound.", False),
        ("The duty logbook: every change recorded — what, when, why, who.", False),
        ("The agent is about to start making entries in your files.", False),
        ("The logbook is how you keep accountability over a teammate who works fast and never sleeps.", False),
        ("Two layers: local logbook (your machine) and remote copy (cloud backup for teams).", False),
        ("No commands today. You are loading the reason. The commands come in the git module.", False),
    ], note="For a military audience, the duty logbook framing lands immediately. For civilian: 'undo button for your whole project.' Read the room. Resist teaching git commands here — the concept is the prerequisite.")

    add_check_on_learning(prs,
        "An agent just ran a batch edit on 40 files. You look at one and it is not right.\n\n"
        "Without version control, what are your options?\n\n"
        "With version control, what changes?\n\n"
        "Why does version control matter specifically once an agent can edit your files?")

    add_hands_on(prs, "Why Version Control Exists", [
        "Think about a file you have edited multiple times over the past month.",
        "Ask yourself: 'If I needed to see what it looked like three weeks ago, could I?'",
        "Ask yourself: 'If something went wrong, what would I lose?'",
        "That gap — between what you want to recover and what you can currently recover — is what version control fills.",
    ], "[IMAGE: Soldier reviewing a duty logbook — every action recorded, every change traceable]")

    add_readiness_check(prs, [
        "I can explain version control as a logbook of file changes I can rewind",
        "I understand why it matters once an agent has write access",
        "I can distinguish the local logbook from the remote copy",
        "I am not yet expected to know any git commands",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.15 — How AI Is Delivered and Paid For
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.15", "How AI Is Delivered\nand Paid For",
        "The same model reaches you through different doors — flat-rate app, pay-per-token API, or your organization's cloud account — and tokens cost money.",
        "[IMAGE: Logistics NCO tracking supply costs against a budget — knowing what each item costs before ordering]")

    add_concept_slide(prs, "Three Delivery Models", [
        ("Subscription app — flat monthly fee, model access included.", False),
        ("Examples: ChatGPT Plus, Claude.ai Pro", True),
        ("Pay-per-token API — you pay per token sent and received.", False),
        ("Examples: Anthropic API, OpenAI API", True),
        ("Bring-your-own-key — your organization's API key foots the bill.", False),
        ("Examples: enterprise deployments, Claude Code", True),
        ("Cloud vs. local: most models run in the cloud. Local models are smaller, less capable, data stays on machine.", False),
    ], note="Do not teach current prices — they change frequently. Teach the structure: tokens cost money, bigger models cost more, match the model to the job. The behavior is the deliverable.")

    add_concept_slide(prs, "Match the Model to the Job", [
        ("Token cost scales with model size and reasoning effort.", False),
        ("'Always use the biggest model' is wrong.", False),
        ("Drafting a quick message: small, fast model.", True),
        ("Analyzing a complex document: larger model.", True),
        ("Writing and running code in an agentic environment: frontier model with tool use.", True),
        ("Using the biggest model for everything is like running a diesel generator to charge a phone.", False),
        ("On a pay-per-token plan: a long conversation with a large model costs more than a short one.", False),
    ], note="Know your billing model before you start a long task. Cost is real — it is not an abstraction.")

    add_check_on_learning(prs,
        "You need to run an analysis task with a very long document and several back-and-forth exchanges.\n\n"
        "What factors determine the cost of that task, and how would you reduce it?\n\n"
        "Why is 'always use the most powerful model' not good advice?")

    add_hands_on(prs, "How AI Is Delivered and Paid For", [
        "Open your chatbot's settings or account page. Find which model you are using.",
        "If you have API access, open the pricing page. Read it for structure — not memorization.",
        "Note: input vs. output tokens, model tiers, cost scaling.",
        "Identify which delivery model you are currently on.",
    ], "[IMAGE: Soldier checking supply manifest — knowing what is available, what it costs, and what is authorized]")

    add_readiness_check(prs, [
        "I can name the three ways AI is paid for",
        "I understand that tokens cost money and model choice affects cost",
        "I know the difference between cloud and local delivery",
        "I can identify which delivery model I am currently using",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.16 — Data Handling: What Never to Paste
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.16", "Data Handling:\nWhat Never to Paste",
        "Once you are pasting real work into a cloud-connected tool, what you paste matters — this rule does not expire.",
        "[IMAGE: Soldier at a classification review station — checking authorization before any information leaves the building]")

    add_concept_slide(prs, "The Bright Line", [
        ("Authorization is a property of the system, not the impressiveness of the tool.", False),
        ("A highly capable model on an unauthorized system is still unauthorized.", False),
        ("What must never go into an unauthorized system:", False),
        ("Personally identifiable information (PII): names + identifiers, SSNs, DOBs, home addresses.", True),
        ("Sensitive, controlled, or classified material of any kind.", True),
        ("Anything above the system's authorization ceiling.", True),
        ("Default posture: when in doubt, do not paste. Ask someone who can authorize it before you proceed.", False),
    ], note="This is one of two modules where the tone is deliberately direct and does not retreat to nuance. Soften it and students hear 'be careful sometimes.' The message is: this is a hard rule.")

    add_concept_slide(prs, "The Classify-Before-You-Paste Habit", [
        ("Before pasting anything into an AI tool, take two seconds:", False),
        ("What is this, and is this system authorized for it?", False),
        ("Build that pause until it is automatic.", False),
        ("This bright line does not expire at the end of this course.", False),
        ("It does not expire under time pressure.", False),
        ("It does not move because the tool is impressive.", False),
        ("One careless paste in an unauthorized system is the kind of mistake with real and lasting consequences.", False),
    ], note="The paraphrase-and-summarize loophole does not exist. If the content is above the authorization ceiling, a summary of it is also above the ceiling.")

    add_check_on_learning(prs,
        "You are under a deadline and want to paste a document into your AI tool to summarize it quickly.\n\n"
        "You are not sure whether the system is authorized for that content.\n\n"
        "What do you do?\n\n"
        "You have a borderline-sensitive document and an unauthorized but highly capable AI tool "
        "that would save hours of work. What is the correct call?")

    add_hands_on(prs, "Data Handling — What Never to Paste", [
        "Think about the last three things you pasted into an AI tool.",
        "For each one: was the system authorized for that type of content?",
        "Identify one category of content you work with regularly that you will never paste into an unauthorized tool.",
        "Write that category down. It is your personal bright line.",
    ], "[IMAGE: Soldier at a classification checkpoint — pausing before passing information through]")

    add_readiness_check(prs, [
        "I can identify the categories that must never go into an unauthorized system",
        "I know the default posture: when in doubt, do not paste",
        "I understand that authorization is a property of the system, not the tool's capability",
        "I have identified my personal bright line for content I handle regularly",
    ])

    # -----------------------------------------------------------------------
    # MODULE 0.17 — The Supervisor Mindset
    # -----------------------------------------------------------------------
    add_section_header(prs, "0.17", "The Supervisor\nMindset",
        "You delegate to the agent, verify its work, and own the outcome — the agent is a capable, fast, assumption-prone junior who will confidently fill gaps you did not address.",
        "[IMAGE: NCO supervising a junior soldier completing a task — clear brief, active oversight, final accountability]")

    add_concept_slide(prs, "The Motivated-Junior Model", [
        ("The agent is a junior teammate with file-system access.", False),
        ("Capable. Fast. Willing to fill ambiguity with plausible-sounding assumptions.", False),
        ("It will never push back or say 'I'm not sure' unless you build that into the prompt.", False),
        ("It executes your intent — including the parts you left implicit.", False),
        ("Three duties of the supervisor:", False),
        ("Delegate clearly — vague intent produces confident but wrong execution.", True),
        ("Verify the work — check what the agent actually did, not just what it said it did.", True),
    ], note="Two failure modes: Blind trust ('it sounds right') — confident output is not verified output. Learned helplessness ('I can't check this, it's too technical') — you don't replicate the work, you check whether it makes sense.")

    add_concept_slide(prs, "Own the Outcome", [
        ("Own the outcome — the capability does not transfer the accountability.", False),
        ("You are still the one who signs for the result.", False),
        ("The delegate-verify-own loop in practice:", False),
        ("Delegate: 'Move files by date into year subfolders. Show me what you plan before you act.'", True),
        ("Verify: check that files landed where they should. Spot-check three. Confirm nothing deleted.", True),
        ("Own: if one landed wrong, you catch it, fix it, and note what to specify more precisely next time.", True),
        ("This is the through-line of the entire course — not a capstone topic.", False),
    ], note="Reinforce throughout, not just at the end. Every agentic action in later modules should be framed with the supervisor loop. Students who only hear it once will not carry it forward.")

    add_check_on_learning(prs,
        "The agent completed a task and the output looks correct at first glance.\n\n"
        "What would you check to actually verify it, and how would you know if an unauthorized assumption had been made?\n\n"
        "An agent completes a task and narrates what it did in clear, confident language. What do you do next?")

    add_hands_on(prs, "The Supervisor Mindset", [
        "Give the agent (or chatbot) a small real task.",
        "Write a clear brief: who you are, what you need, what good output looks like, and what is off-limits.",
        "Submit it. Read the output.",
        "Verify: does it do what you asked? Did it make any assumptions you did not authorize? Is anything wrong?",
        "If something is off, correct it. Identify what you should have specified more precisely.",
    ], "[IMAGE: NCO reviewing completed task before signing off — verifying the work, not just accepting the narration]")

    add_readiness_check(prs, [
        "I can state the three duties of the supervisor: delegate, verify, own",
        "I know the difference between a capable junior and a trustworthy one",
        "I have completed at least one delegate-verify-correct loop with a real task",
        "I am the commander, not the typist",
    ])

    add_section_summary(prs, "Phase 4 — From Chatbot to Agent", [
        "Chatbot gives advice. Agent takes action — read, write, execute on your real machine.",
        "Version control is the duty logbook — rewindable record of every agent change.",
        "Data handling: authorization is a property of the system, not the tool's capability.",
        "Supervisor mindset: delegate clearly, verify the work, own the outcome.",
    ])

    # -----------------------------------------------------------------------
    # SUMMARY TABLE
    # -----------------------------------------------------------------------
    add_summary_table(prs,
        "DAY 0 MODULE SUMMARY",
        "Four Phases — Zero to Ready",
        ["PHASE", "WHAT IT BUILDS", "KEY HABIT"],
        [
            ("Phase 1: Know What You Are Talking To",
             "LLM mental model, failure modes, deliberate prompting",
             "Verify before you act on model output"),
            ("Phase 2: The Computer Underneath",
             "File/folder literacy, plaintext vs rich text, code editor",
             "See the tree — apps hide it, tools need it"),
            ("Phase 3: The Terminal",
             "pwd, ls, cd, mkdir, cp, mv, flags, tab-completion, Ctrl+C",
             "Run ls after every action — verify the round landed"),
            ("Phase 4: From Chatbot to Agent",
             "Chatbot vs agent, version control, data handling, supervisor mindset",
             "Delegate clearly, verify the work, own the outcome"),
        ])

    # -----------------------------------------------------------------------
    # FINAL READINESS CHECK
    # -----------------------------------------------------------------------
    add_readiness_check(prs, [
        "I can explain what an LLM does and name the four failure modes",
        "I have produced a hallucination with my own hands",
        "I can write a deliberate prompt: role, context, example, output spec",
        "I can navigate the filesystem from the command line — pwd, ls, cd",
        "I can create, copy, and move files from the command line",
        "I use tab-completion on every path without thinking about it",
        "I can state the one-line difference between a chatbot and an agent",
        "I know what data must never go into an unauthorized system",
        "I can name the three duties of the supervisor: delegate, verify, own",
    ])

    # -----------------------------------------------------------------------
    # END SLIDE
    # -----------------------------------------------------------------------
    add_end_slide(prs, "AI / Agentics Basics\nDay 0 — Pre-Course", [
        "Phases 2 and 3 paid off in Module 0.13. The terminal is how the agent moves — now you can read the terrain.",
        "Complete the supervisor loop at least once before Day 1: give a task, verify the output, own the result.",
        "Identify your personal data bright line and write it down — it applies from this day forward.",
        "You are not a typist. You are a commander. Everything that follows stacks on that.",
    ])

    # -----------------------------------------------------------------------
    # Save
    # -----------------------------------------------------------------------
    out = os.path.join(os.path.dirname(__file__), "day0.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
