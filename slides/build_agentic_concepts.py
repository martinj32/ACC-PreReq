"""
ACC Agentic Concepts Slide Deck Builder — Army Cyber Dark Theme
Source: docs/agentic-ai/agent-concepts.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_summary_table(prs, title, subtitle, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, subtitle)
    _footer(slide)

    col_x = [L, Inches(3.5), Inches(8.1)]
    col_w = [Inches(2.9), Inches(4.4), Inches(5.0)]
    ht = Inches(1.52)
    rh = Inches(0.42)

    hdr = _rect(slide, L, ht, R, rh, ELECTRIC_BLUE)
    _glow_blue(hdr, 4)
    for i, h in enumerate(["CONCEPT", "CORE IDEA", "WHY IT MATTERS"]):
        _tb(slide, h, col_x[i] + Inches(0.1), ht + Inches(0.06), col_w[i], rh,
            sz=11, bold=True, color=CYBER_BLACK)

    for ri, (concept, idea, why) in enumerate(rows):
        top = ht + rh + ri * Inches(0.83)
        _rect(slide, L, top, R, Inches(0.83), DARK_NAVY if ri % 2 == 0 else NAVY_MID)
        bar = _rect(slide, L, top, Pt(3), Inches(0.83), ELECTRIC_BLUE)
        _glow_blue(bar, 3)
        for ci, (val, color) in enumerate([(concept, ELECTRIC_BLUE), (idea, WHITE), (why, LIGHT_GRAY)]):
            _tb(slide, val, col_x[ci] + Inches(0.12), top + Inches(0.08),
                col_w[ci] - Inches(0.15), Inches(0.8),
                sz=11, bold=(ci == 0), color=color)


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- TITLE ---
    add_title_slide(prs,
        "Agentic AI —\nCore Concepts",
        "MODULE  |  AGENTIC AI FUNDAMENTALS",
        "[IMAGE: Soldier at a command post terminal with a supervisory posture, overseeing automated system activity on screen]")

    # --- OVERVIEW ---
    add_overview_slide(prs,
        "A chatbot gives advice; an agent takes action. This module explains the difference, why version control becomes critical the moment an agent can edit your files, and the supervisor mindset that keeps agentic work responsible.",
        "Self-contained module (prerequisites: Bedrock, Terminal Basics)",
        [
            "State the one-line difference between a chatbot and an agent",
            "Explain what read, write, and execute access means in practice",
            "Identify all three parts of the engine-harness-operator model",
            "Explain why version control matters once an agent has write access",
            "State the three duties of the supervisor: delegate, verify, own",
            "Recognize and counter the two main supervision failure modes",
        ])

    # ==========================================================================
    # SECTION 1 — Chatbot vs Agent
    # ==========================================================================
    add_section_header(prs, "01", "Chatbot\nvs Agent",
        "A chatbot gives advice. An agent takes action. The difference is read, write, and execute access to your real files and terminal.",
        "[IMAGE: Soldier asking for directions versus a soldier physically leading the patrol — advisor vs. actor]")

    add_concept_slide(prs, "What Is an Agent?", [
        ("Ask a chatbot to rename a folder — it tells you the command.", False),
        ("Ask an agent — it renames the folder.", False),
        ("That is the entire difference. The agent has access.", False),
        ("Three levels of access:", False),
        ("Read — the agent can look at your files", True),
        ("Write — the agent can create and modify your files", True),
        ("Execute — the agent can run commands in your terminal", True),
        ("The LLM is the engine. A harness with tools gives it hands. You are the operator.", False),
    ], note="This is the payoff for Terminal Basics. Every command students learned — pwd, ls, cd, mkdir, mv — is what the agent runs in their filesystem. They learned the terrain to supervise someone else navigating it.")

    add_concept_slide(prs, "Read, Write, Execute Is a Lot of Trust", [
        ("This is not abstract.", False),
        ("An agent with write and execute access can create, modify, or delete files on your real machine.", False),
        ("You are extending significant trust when you give an agent these permissions.", False),
        ("The engine-harness-operator model:", False),
        ("Engine — the LLM: reasons, generates, plans", True),
        ("Harness — the tool layer: gives the engine reach into your filesystem and terminal", True),
        ("Operator — you: directs the mission and pulls the plug when it goes wrong", True),
        ("Everything in advanced agentic work stacks on this primitive.", False),
    ], note="WARNING: An agent with write and execute access can create, modify, or delete files on your real machine. The supervisor mindset in Section 3 is what keeps that from becoming a problem.")

    add_check_on_learning(prs,
        "You asked a chatbot to rename a folder. It gave you instructions.\n\n"
        "You asked an agent the same question. It renamed the folder.\n\n"
        "What are you trusting the agent with when it does that — and what would happen if it made an assumption about which folder to rename?")

    add_hands_on(prs, "Chatbot vs Agent", [
        "Open your AI chatbot (Claude web, ChatGPT, or equivalent).",
        "Ask it: 'Rename the folder project to project-v1.' Read the response. Did it rename anything or tell you how?",
        "If Claude Code is available — open it in a project folder. Ask the same question. Watch what it does.",
        "Does it call a tool? Ask for confirmation? Does the folder change on disk?",
        "Write one sentence: what was different between the two responses? Which one did the work?",
    ], "[IMAGE: Soldier comparing two radio responses — one with instructions, one with a completed action report]")

    add_section_summary(prs, "Chatbot vs Agent", [
        "A chatbot gives advice. An agent takes action — that is the entire difference.",
        "Agents have three levels of access: read, write, and execute.",
        "Engine + harness + operator = the complete agentic system.",
        "You are extending real trust. The supervisor mindset is what keeps it responsible.",
    ])

    # ==========================================================================
    # SECTION 2 — Why Version Control Exists
    # ==========================================================================
    add_section_header(prs, "02", "Why Version\nControl Exists",
        "When an agent can change your files, you want a logbook of every change — what changed, when, and why — with the ability to go back.",
        "[IMAGE: Army duty NCO maintaining a logbook at a battalion operations desk]")

    add_concept_slide(prs, "The Logbook for Your Files", [
        ("Every time a file changes, version control records a snapshot.", False),
        ("What changed. When. A note about why. Stored locally on your machine.", False),
        ("Every snapshot can be rewound.", False),
        ("Version control is the duty logbook.", False),
        ("Every change recorded: what, when, why, who.", True),
        ("The agent is about to start making entries in your files.", True),
        ("The logbook is how you keep accountability over a teammate who works fast and never sleeps.", True),
    ], note="Why software people seem to remember everything: they do not. They have the log. 'Who changed this and when?' is a two-second lookup. 'What did this file look like last Tuesday?' is one command.")

    add_concept_slide(prs, "Two Layers of Version Control", [
        ("Local logbook — your machine keeps the record.", False),
        ("You can always rewind to any previous snapshot without needing the internet.", True),
        ("Every intermediate state is saved.", True),
        ("Remote copy — the same history synchronized to a cloud server.", False),
        ("Backup — protection against local machine failure.", True),
        ("The mechanism for teammates working on the same files.", True),
        ("No commands today. You are loading the reason. Commands come in the git session.", False),
    ], note="Instructor note: Resist teaching git commands here. The concept is the prerequisite. If a student asks 'how do I actually do this' — tell them it comes in the next session. Today is the why.")

    add_check_on_learning(prs,
        "An agent just ran a batch edit on 40 files.\n\n"
        "You look at one file and it is not right.\n\n"
        "Without version control — what are your options?\n\n"
        "With version control — what changes?")

    add_hands_on(prs, "Why Version Control Exists", [
        "Think about a file you have edited multiple times over the past month.",
        "Ask yourself: 'If I needed to see what it looked like three weeks ago, could I?'",
        "Ask yourself: 'If something went wrong in that file right now, what would I lose?'",
        "That gap — between what you want to recover and what you can currently recover — is what version control fills.",
        "No commands today. Load the reason. The commands come next.",
    ], "[IMAGE: Soldier reviewing a handwritten duty log — cross-referencing entries from the previous week]")

    add_section_summary(prs, "Why Version Control Exists", [
        "Version control is a logbook: every file change recorded with what, when, and why.",
        "Every snapshot can be rewound — local recovery without the internet.",
        "Remote copy provides backup and enables team coordination.",
        "Once an agent has write access, you need this logbook — it works fast and never sleeps.",
    ])

    # ==========================================================================
    # SECTION 3 — The Supervisor Mindset
    # ==========================================================================
    add_section_header(prs, "03", "The Supervisor\nMindset",
        "Delegate clearly, verify the work, and own the outcome. The agent fills ambiguity with plausible assumptions — your job is to command well and check.",
        "[IMAGE: NCO supervising a junior soldier completing a task — reviewing the work before signing off]")

    add_concept_slide(prs, "The Motivated-Junior Model", [
        ("The agent is a junior teammate with file-system access.", False),
        ("Capable. Fast. Willing to fill ambiguity with plausible-sounding assumptions.", False),
        ("It will never push back or say 'I'm not sure' unless you build that into the prompt.", False),
        ("It executes your intent — including the parts you left implicit.", False),
        ("Three duties of the supervisor:", False),
        ("Delegate clearly — vague intent produces confident but wrong execution", True),
        ("Verify the work — check what the agent actually did, not just what it said it did", True),
        ("Own the outcome — the capability does not transfer the accountability", True),
    ], note="The capability does not transfer the accountability. You are still the one who signs for the result.")

    add_concept_slide(prs, "Two Failure Modes to Name and Counter", [
        ("Failure Mode 1: Blind Trust", False),
        ("'It sounds right.' Confident output from an agent is not verified output.", True),
        ("Check the work before you act on it — fluency is not evidence of accuracy.", True),
        ("Failure Mode 2: Learned Helplessness", False),
        ("'I can't check this, it's too technical.'", True),
        ("You do not need to replicate the agent's work.", True),
        ("You need to check: does output make sense? Were constraints honored? Does anything look wrong?", True),
        ("That is a human judgment call — not a technical skill.", True),
    ], note="Both failure modes leave you exposed. Blind trust means you sign for something you did not verify. Learned helplessness means you hand over accountability without a check. Neither is acceptable.")

    add_concept_slide(prs, "The Delegate-Verify-Own Loop", [
        ("Step 1: Delegate — write a clear brief.", False),
        ("Who you are. What you need. What good output looks like. What is off-limits.", True),
        ("Step 2: Verify — check what the agent actually did.", False),
        ("Not what it narrated. The actual output.", True),
        ("Step 3: Own — if something is wrong, you catch it, fix it, and note what to specify more precisely.", False),
        ("Example — reorganize a folder of files by date:", False),
        ("Delegate: 'Move files by date. Do not delete anything. Show me the plan first.'", True),
        ("Verify: spot-check three files. Confirm nothing was deleted.", True),
    ], note="Optimize for the loop, not the polish. A rough artifact that was checked is a better outcome than a finished artifact that was not.")

    add_check_on_learning(prs,
        "An agent completes a task and narrates what it did in clear, confident language.\n\n"
        "The output looks correct at first glance.\n\n"
        "What would you check to actually verify it — and how would you know if an unauthorized assumption had been made?")

    add_hands_on(prs, "The Supervisor Mindset", [
        "Give the agent (or chatbot) a small real task.",
        "Write a clear brief: who you are, what you need, what good output looks like, and what is off-limits.",
        "Submit it. Read the output.",
        "Verify: does it do what you asked? Did it make assumptions you did not authorize? Is anything wrong?",
        "If something is off — correct it. Identify what you should have specified more precisely.",
        "That loop — delegate, verify, correct — is the whole skill.",
    ], "[IMAGE: NCO with clipboard reviewing completed work by a junior soldier before signing off at a command post]")

    add_section_summary(prs, "The Supervisor Mindset", [
        "The agent is a motivated junior with file-system access — capable, fast, and assumption-prone.",
        "Three duties: delegate clearly, verify the work, own the outcome.",
        "Blind trust and learned helplessness are both failure modes — check the work.",
        "The loop: delegate → verify → correct. That is the entire skill.",
    ])

    # ==========================================================================
    # WORKED EXAMPLE
    # ==========================================================================
    add_section_header(prs, "04", "Worked\nExample",
        "All three concepts — agent access, version control, and supervisor mindset — in a single scenario.",
        "[IMAGE: After-action review board with soldiers analyzing an operation log at a conference table]")

    add_example_slide(prs, "Example: The Supervisor Loop in Action",
        "Scenario: Agent reorganizes a project folder of dated field reports",
        [
            "BRIEF (Delegate clearly):",
            "  'Move all files with a date in the filename into subfolders by year.",
            "   Do not delete anything. Show me the plan before you execute.'",
            "",
            "AGENT RESPONSE:",
            "  Lists 47 files. Proposes folder structure. Waits for approval.",
            "",
            "VERIFY:",
            "  You review the plan. Spot-check three files after execution.",
            "  Confirm nothing was deleted. Version control log shows every move.",
            "",
            "OWN:",
            "  One file landed in wrong year folder — you catch it, fix it, update the brief.",
        ],
        "[IMAGE: Soldier reviewing a printed file manifest against what appeared on screen — line by line verification]")

    # ==========================================================================
    # SUMMARY TABLE
    # ==========================================================================
    add_summary_table(prs,
        "MODULE SUMMARY",
        "Three Foundations of Agentic Work",
        [
            ("Chatbot vs Agent",      "Agents read, write, and execute — they act",      "Trust is real. You need a supervisor mindset before granting access."),
            ("Version Control",       "Logbook of every file change you can rewind",      "Once an agent edits files, this is not optional — it is accountability."),
            ("Supervisor Mindset",    "Delegate clearly, verify the work, own the outcome", "The loop keeps capable-and-fast from becoming a liability."),
        ])

    # ==========================================================================
    # READINESS CHECK
    # ==========================================================================
    add_readiness_check(prs, [
        "I can state the one-line difference between a chatbot and an agent",
        "I can explain what read, write, and execute access means in practice",
        "I understand why Terminal Basics was the prerequisite for this module",
        "I can name all three parts of the engine-harness-operator model",
        "I can explain version control as a logbook of file changes I can rewind",
        "I understand why version control matters once an agent has write access",
        "I can distinguish the local logbook from the remote copy",
        "I can state the three duties of the supervisor: delegate, verify, own",
        "I know the difference between blind trust and verified output",
        "I have completed at least one delegate-verify-correct loop with a real task",
    ])

    # ==========================================================================
    # END SLIDE
    # ==========================================================================
    add_end_slide(prs, "Agentic AI —\nCore Concepts", [
        "Complete the hands-on exercises in all three sections before continuing.",
        "In your first Claude Code session, name out loud which access level you are granting: read, write, or execute.",
        "Before any agentic task: write the brief, verify the output, own the result.",
        "Next module: working with git — the commands that make the logbook real.",
    ])

    out = os.path.join(os.path.dirname(__file__), "agentic-concepts.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
