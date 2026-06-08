"""
ACC Mental Models Slide Deck Builder — Army Cyber Dark Theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_summary_table(prs, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE SUMMARY", "The Six Mental Models")
    _footer(slide)

    col_x = [L, Inches(3.5), Inches(8.1)]
    col_w = [Inches(2.9), Inches(4.4), Inches(5.0)]
    ht = Inches(1.52)
    rh = Inches(0.42)

    hdr = _rect(slide, L, ht, R, rh, ELECTRIC_BLUE)
    _glow_blue(hdr, 4)
    for i, h in enumerate(["MODEL", "CORE IDEA", "WHY IT MATTERS"]):
        _tb(slide, h, col_x[i] + Inches(0.1), ht + Inches(0.06), col_w[i], rh,
            sz=11, bold=True, color=CYBER_BLACK)

    for ri, (model, idea, why) in enumerate(rows):
        top = ht + rh + ri * Inches(0.83)
        _rect(slide, L, top, R, Inches(0.83), DARK_NAVY if ri % 2 == 0 else NAVY_MID)
        bar = _rect(slide, L, top, Pt(3), Inches(0.83), ELECTRIC_BLUE)
        _glow_blue(bar, 3)
        for ci, (val, color) in enumerate([(model, ELECTRIC_BLUE), (idea, WHITE), (why, LIGHT_GRAY)]):
            _tb(slide, val, col_x[ci] + Inches(0.12), top + Inches(0.08),
                col_w[ci] - Inches(0.15), Inches(0.8),
                sz=11, bold=(ci == 0), color=color)


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    add_title_slide(prs,
        "Mental Models for\nAI-Assisted Operations",
        "PREREQUISITE MODULE  |  35 MINUTES",
        "[IMAGE: Soldier reviewing data at a forward command post terminal]")

    add_overview_slide(prs,
        "Six mental models separate operators who use AI effectively from those who waste tokens, break things, and overpay.",
        "35 minutes (self-contained)",
        ["Understand the 'harness': an LLM becomes agentic when paired with tools",
         "Grasp why context windows matter and when they become a constraint",
         "Treat tokens as a resource with real cost",
         "Understand what tool calls are and why they enable verification",
         "Adopt active supervision as your default posture",
         "Build cost-consciousness as a core habit"])

    # --- SECTION 1 ---
    add_section_header(prs, "01", "The Harness\nMental Model",
        "LLM + Tools = Agency. Without tools, Claude is a chatbot. With tools, Claude is a collaborator.",
        "[IMAGE: Soldier with radio harness — the gear that connects and enables action]")

    add_concept_slide(prs, "What Is the Harness?", [
        ("An LLM by itself is text-in, text-out. It cannot act on the world.", False),
        ("The harness is the tool layer that makes an LLM into an agent:", False),
        ("Eyes — tools to read files, check systems, fetch URLs", True),
        ("Hands — tools to write files, edit code, create directories", True),
        ("Body — tools to run commands, call APIs, deploy outputs", True),
        ("The engine cannot act without a harness. The harness needs an operator.", False),
    ], note="Working vocabulary: Engine = the LLM (reasons, generates). Harness = tool layer (gives the engine reach). Operator = you (directs the mission, owns the result).")

    add_concept_slide(prs, "The Engine-Harness-Operator Stack", [
        ("Engine — the LLM. Reasoning brain. Generates, plans, decides.", False),
        ("Harness — the tool layer. Gives the engine access to files, commands, and external systems.", False),
        ("Operator — you. Directs the mission, approves consequential actions, carries accountability.", False),
        ("Formula: LLM (knowledge + reasoning) + Tools (sensors + actuators) = Agency", False),
        ("Without tools: Claude writes code but cannot run it or read your file system.", False),
        ("With tools: Claude reads, modifies, tests, commits, and iterates on your actual work.", False),
    ])

    add_concept_slide(prs, "Your Job in the Harness", [
        ("When you invoke Claude with tools enabled, you are building a feedback loop.", False),
        ("Claude sees the world via tools, takes action, observes the result, and adjusts.", False),
        ("Your four responsibilities as operator:", False),
        ("Set up the right tools and permissions", True),
        ("Give clear instructions — vague orders get vague results", True),
        ("Supervise the feedback loop", True),
        ("Interrupt and correct when it goes wrong", True),
    ], note="Warning: Your job is not to give the LLM a task and walk away. You are not a spectator. You are the operator — you own the output.")

    add_check_on_learning(prs,
        "You give Claude a mission: search 200 field reports for grid 38SMB4521 and return every matching excerpt.\n\n"
        "Without a harness, Claude cannot read the files. It can only generate plausible-sounding excerpts from memory — "
        "and you have no way to know which it did.\n\n"
        "With a harness, Claude calls tools against the actual files and returns real, traceable results.\n\n"
        "What does that difference mean for how you verify the output?")

    add_hands_on(prs, "The Harness Mental Model", [
        "Open your AI chatbot. Ask it: 'Rename the folder project to project-v1.' Read the response.",
        "Identify: did it rename the folder, or tell you how? That gap is the harness.",
        "If you have Claude Code available, ask it the same thing in a project folder. Compare.",
        "You are seeing the difference between a text engine and an agent.",
    ], "[IMAGE: Soldier at terminal comparing two different outputs on dual monitors]")

    add_section_summary(prs, "The Harness Mental Model", [
        "An LLM alone is text-in, text-out — it cannot act without tools.",
        "The harness (tool layer) gives the engine eyes, hands, and a body.",
        "Engine + Harness + Operator = the complete agentic system.",
        "Your role: set up tools, give clear orders, supervise the loop, own the result.",
    ])

    # --- SECTION 1.5 ---
    add_section_header(prs, "1.5", "The Model\nLandscape",
        "Three families dominate agentic work. Each offers tiers — fast/cheap to slow/powerful. Match the model to the mission.",
        "[IMAGE: Army equipment lineup — different tools for different mission sets]")

    add_concept_slide(prs, "Model Families and Tiers", [
        ("Three major families: Claude (Anthropic), GPT (OpenAI), Gemini (Google).", False),
        ("Each offers multiple tiers — fast/cheap through slow/powerful.", False),
        ("Claude's tiers — the standard working toolchain:", False),
        ("Haiku — fast and cheap. High-volume tasks, evaluation loops, simple formatting.", True),
        ("Sonnet — balanced. The standard working model for most missions.", True),
        ("Opus — powerful and expensive. Reserve for the hardest reasoning tasks.", True),
        ("'Always use the biggest model' is the wrong instinct.", False),
    ], note="Decision rule: Match the model to the task. Summarization: fast tier. Complex analysis: capable tier. Tight loops: cheapest tier that meets quality requirements.")

    add_concept_slide(prs, "Name Drift Is Real", [
        ("Model names and pricing change with every major release.", False),
        ("The tier concept is durable. The specific names are not.", False),
        ("Verify current model IDs and costs at anthropic.com before any course run.", False),
        ("Bookmarking the docs is part of the job — not optional.", False),
    ], note="Do not teach current prices or model names as facts. They will be wrong within a quarter. Teach the structure: tiers, cost scaling, task matching.")

    add_hands_on(prs, "The Model Landscape", [
        "Go to anthropic.com/claude and find the current model tiers.",
        "Note the current names for fast, balanced, and powerful models. Do they match this section?",
        "Pick one task you might do with AI (summarize, write code, analyze options).",
        "Decide which tier is appropriate and write one sentence explaining why.",
    ], "[IMAGE: Soldier cross-referencing field manual with digital screen]")

    add_section_summary(prs, "The Model Landscape", [
        "Three families: Claude, GPT, Gemini. Each has fast, balanced, and powerful tiers.",
        "Match the model to the task — bigger is not always better.",
        "Model names drift with releases. Verify before you teach or build.",
    ])

    # --- SECTION 1.6 ---
    add_section_header(prs, "1.6", "How the\nModel Fails",
        "Confident does not mean correct. The model has no truth-checking step — that step is yours.",
        "[IMAGE: Soldier double-checking map coordinates before calling in a grid]")

    add_concept_slide(prs, "Failure Mode 1: Hallucination", [
        ("The model generates tokens that are statistically likely to follow prior context.", False),
        ("It has no separate truth-checking step.", False),
        ("It cannot distinguish accurate training data from plausible completion.", False),
        ("Confident output and correct output are entirely unrelated.", False),
        ("This is expected system behavior — not a rare bug.", False),
    ], note="The sharp junior analyst who never says 'I don't know.' Confident, fluent, and occasionally making it up. You check the work — fluency is not evidence of accuracy.")

    add_concept_slide(prs, "Failure Modes 2, 3, and 4", [
        ("Confident-Wrong — hallucination is not always dramatic.", False),
        ("Subtly wrong dates, statistics, citations. Equally confident tone regardless of accuracy.", True),
        ("Never use the model's confidence as a signal of truth.", True),
        ("Nondeterminism — the model uses temperature, introducing randomness by design.", False),
        ("Same prompt, different run, different result. One output is not 'the answer.'", True),
        ("Knowledge Cutoff — the model was trained to a point in time.", False),
        ("It may refuse, guess, or confidently describe events that did not happen.", True),
    ], note="All three failure modes look identical from the outside. Verify anything time-sensitive regardless of how confidently the model presents it.")

    add_check_on_learning(prs,
        "You receive an intelligence summary from Claude with three specific claims about adversary activity. "
        "Claude presents all three confidently.\n\n"
        "You have no way to tell whether they came from source documents Claude read, or from training data pattern-matching.\n\n"
        "What do you do before you include this summary in a product?")

    add_hands_on(prs, "How the Model Fails", [
        "Ask Claude: 'What is the current price per million tokens for Claude Sonnet?' Read the answer.",
        "Go to anthropic.com/pricing and compare. Are they the same?",
        "Ask Claude: 'When is your knowledge cutoff?' Note whether it flags any uncertainty.",
        "Ask: 'What major Claude models were released after your cutoff?' Observe how it handles the gap.",
        "You have just demonstrated all three failure modes in a single exercise.",
    ], "[IMAGE: Analyst comparing two documents side by side — verifying source vs. output]")

    add_section_summary(prs, "How the Model Fails", [
        "Hallucination: the model completes plausibly, not accurately. No truth-checking step.",
        "Nondeterminism: same prompt, different run, different result — by design.",
        "Knowledge cutoff: confident answers about events past training may be wrong.",
        "Verification is your job — a property of how the system works, not a sign of distrust.",
    ])

    # --- SECTION 2 ---
    add_section_header(prs, "02", "Context\nWindows",
        "The context window is the LLM's working memory. Fill it and it starts forgetting — silently.",
        "[IMAGE: Tactical whiteboard filling up with overlapping mission graphics and notes]")

    add_concept_slide(prs, "What Is a Context Window?", [
        ("A context window is the total text (in tokens) the LLM can 'see' at one time.", False),
        ("Think of it as working memory — a finite whiteboard.", False),
        ("Current Claude models: 200,000 tokens.", False),
        ("One page of text ≈ 300-500 tokens. A 50-page document ≈ 15,000-25,000 tokens.", False),
        ("Tool call results accumulate in context — a long session can fill the window through results alone.", False),
        ("Watch for: contradictions, missed instructions, degrading response quality.", False),
    ], note="Verify current context window sizes before any course run — they change with model releases.")

    add_concept_slide(prs, "The Three Rules of Context", [
        ("Rule 1: Context is not infinite.", False),
        ("At 200,000 tokens you have real limits. Long conversations and verbose reasoning add up fast.", True),
        ("Rule 2: Tool results accumulate in context.", False),
        ("The file read itself doesn't consume tokens. But every result returned does — and it builds up.", True),
        ("Rule 3: You control what's in context.", False),
        ("You decide what to include, how verbose to be, how much history to keep.", True),
        ("Strategic use of context is a skill — not a setting.", False),
    ])

    add_concept_slide(prs, "Context Windows in the Harness", [
        ("Claude does not need the entire codebase in context to work on large projects.", False),
        ("The tool call loop keeps context lean:", False),
        ("Claude generates a request: 'Read /path/to/file.js'", True),
        ("The tool executes and returns the file contents into context", True),
        ("Claude reasons about what it read, then calls the next tool or generates output", True),
        ("Claude reads what it needs, when it needs it. Smart sampling, not cramming.", False),
    ], note="Solution when context fills: start a fresh session, paste only relevant context, or ask Claude to read from files instead of pasting everything.")

    add_check_on_learning(prs,
        "You are in a Claude Code session that has been running for two hours. "
        "You have read 30 files, run 15 commands, and had a long back-and-forth.\n\n"
        "Claude starts giving you answers that contradict what it said an hour ago.\n\n"
        "What is the most likely cause, and what do you do?")

    add_hands_on(prs, "Context Windows", [
        "Open Claude (web). Start a fresh conversation. Type: 'My name is [your name]. Remember this.'",
        "Have a long, unrelated conversation — 10 to 20 back-and-forth exchanges.",
        "Ask: 'What was my name?' Note whether Claude still knows.",
        "If you have Claude Code, run a long session with several file reads. Ask Claude to summarize what it has done. Check accuracy.",
    ], "[IMAGE: Soldier reviewing a long operation log — noting what has been forgotten or dropped]")

    add_section_summary(prs, "Context Windows", [
        "The context window is finite working memory — fill it and the model loses earlier content.",
        "Tool call results accumulate in context and can silently fill the window.",
        "You control what goes in. Strategic use of context is a skill.",
        "Start fresh sessions proactively before context degradation causes problems.",
    ])

    # --- SECTION 3 ---
    add_section_header(prs, "03", "Tokens as\nCurrency",
        "Every token costs time and money. Concise, structured prompts save iterations. Verbose prompts waste both.",
        "[IMAGE: Supply sergeant accounting for limited ammunition and fuel — finite resources]")

    add_concept_slide(prs, "What Is a Token?", [
        ("A token is the unit of text the LLM processes — not a word, but a chunk.", False),
        ("1 token ≈ 0.75 English words (varies by language and content type).", False),
        ("1 page of text ≈ 300-500 tokens. 1 line of code ≈ 5-15 tokens (code is denser).", False),
        ("Every input token costs money. Every output token costs more.", False),
        ("Total cost = (input tokens × input price) + (output tokens × output price).", False),
        ("More tokens = slower response time and higher dollar cost.", False),
    ], note="Verify current pricing at anthropic.com/pricing — do not rely on token costs in this deck. The structure is durable; the numbers are not.")

    add_concept_slide(prs, "Four Token Principles", [
        ("Principle 1: Concision is a feature.", False),
        ("A 100-word structured prompt often beats a 500-word rambling one. Less noise, lower cost.", True),
        ("Principle 2: Repetition is waste.", False),
        ("Asking the same question 3 times = 3x the tokens. Iterate once.", True),
        ("Principle 3: Cheap vs. expensive operations matter.", False),
        ("Cheap: ask Claude to read a file via tool call. Expensive: paste 30,000 tokens into the prompt.", True),
        ("Principle 4: Plan before you prompt.", False),
        ("What context does Claude need? What should it read via tools vs. receive pasted?", True),
    ])

    add_check_on_learning(prs,
        "You need Claude to analyze a 200-page field manual.\n\n"
        "Option A: Paste the entire manual into the chat and ask your question.\n\n"
        "Option B: Ask Claude to search the manual using tools and read only the relevant sections.\n\n"
        "Which is more cost-effective? Which produces more reliable results? Why?")

    add_hands_on(prs, "Tokens as Currency", [
        "Write a vague prompt asking Claude to help you with a project. Count the words.",
        "Rewrite it as a structured prompt: task, context, constraints, output format. Count the words again.",
        "Submit both and compare the quality of the first response each time.",
        "Which required fewer follow-up iterations? Which approach was cheaper?",
        "Optional: paste each prompt into a tokenizer and count the actual token cost.",
    ], "[IMAGE: Two mission briefs side by side — one a wall of text, one a clean OPORD format]")

    add_section_summary(prs, "Tokens as Currency", [
        "Every token has a cost in time and money. Treat it as a finite resource.",
        "Structured prompts outperform verbose ones — less noise, fewer iterations.",
        "Use tool calls to read files instead of pasting large documents into context.",
        "Plan before you prompt. One clear ask beats five vague ones every time.",
    ])

    # --- SECTION 4 ---
    add_section_header(prs, "04", "Tool Calls and\nDeterminism",
        "Tool calls let Claude verify its reasoning against actual system state — not generate from memory.",
        "[IMAGE: Soldier physically verifying equipment readiness against a checklist]")

    add_concept_slide(prs, "What Is a Tool Call?", [
        ("A tool call is a structured request from the LLM to an external system.", False),
        ("It says: 'Execute this specific action and return the result.'", False),
        ("Claude does not execute tool calls itself — it requests that the harness execute them.", False),
        ("Types of tool calls:", False),
        ("Reading: read files, list directories, check git status, fetch URLs", True),
        ("Writing: create, edit, delete, move files", True),
        ("Executing: run commands, run tests, deploy code, call APIs", True),
    ], note="Example: read(\"/src/auth.js\") → result → write(\"/src/auth.js\", fix) → run(\"npm test\") → all tests pass. Every step is traceable.")

    add_concept_slide(prs, "Why Tool Calls Enable Verification", [
        ("Without tools: Claude generates text. You cannot verify whether it hallucinated.", False),
        ("With tools: Claude can verify its own reasoning against actual system state.", False),
        ("Every tool call is a verifiable, traceable action — your audit trail.", False),
        ("Claude read that specific file. Claude ran that specific test.", True),
        ("Claude wrote that specific content. You can trace every step.", True),
        ("That is not possible with pure text generation.", False),
    ])

    add_check_on_learning(prs,
        "Claude tells you it has verified that your export script produces correctly formatted output.\n\n"
        "Scenario A: Claude reviewed the output description you pasted and said it looks correct.\n\n"
        "Scenario B: Claude called a read tool on the actual output file and checked it against the template.\n\n"
        "Which scenario gives you a defensible product? Why?")

    add_hands_on(prs, "Tool Calls and Determinism", [
        "In Claude Code, ask: 'What is in my current directory?' Watch it call the list tool and return real results.",
        "Ask: 'Read [a specific file in your project].' Watch it read the actual file — not guess at contents.",
        "Ask Claude to make a small edit. Before it writes, ask: 'What exactly are you going to change and why?'",
        "Verify before approving. You are tracing the tool call loop live. Every action is visible.",
    ], "[IMAGE: Soldier reading back a radio transmission to confirm accuracy before relaying]")

    add_section_summary(prs, "Tool Calls and Determinism", [
        "A tool call is a structured request: execute this action and return the result.",
        "Tool calls let Claude verify its reasoning against real system state — not generate from memory.",
        "Every tool call is a verifiable, traceable action — your audit trail.",
        "Without tool calls, you cannot distinguish correct output from confident hallucination.",
    ])

    # --- SECTION 5 ---
    add_section_header(prs, "05", "Operator Posture\nand Supervision",
        "The agent is a motivated PFC with file-system access. Over-brief, supervise closely, verify before you sign.",
        "[IMAGE: NCO supervising a junior soldier completing a task at a terminal]")

    add_concept_slide(prs, "What Is Operator Posture?", [
        ("Your posture is your stance toward the AI:", False),
        ("Delegating — asking it to do everything while you step away (wrong)", True),
        ("Supervising — watching closely, intervening as needed (correct)", True),
        ("Collaborating — both thinking, both deciding (correct, context-dependent)", True),
        ("The agent is a motivated PFC with file-system access and a mission.", False),
        ("Capable, fast, eager — and will execute confidently on an ambiguous brief.", False),
        ("The capability does not transfer the accountability. You sign for the result.", False),
    ], note="Human in the loop means you are the decision point — not a passive observer. Review before committing. Intervene early.")

    add_concept_slide(prs, "Three Ways Supervision Breaks Down", [
        ("Mistake 1: Over-trust", False),
        ("You let Claude run for hours without reviewing. You get 20 breaking changes you didn't expect.", True),
        ("Fix: review after each major step. Ask Claude to show changes before committing.", True),
        ("Mistake 2: Under-involvement", False),
        ("You give a vague order ('Make this faster'). Claude spends 5,000 tokens on guesses.", True),
        ("Fix: specific orders. Name the file, the goal, the constraint, what success looks like.", True),
        ("Mistake 3: Automation Fallacy", False),
        ("You chain build, test, and deploy. One step fails silently. Broken output ships.", True),
    ], note="Fix for automation fallacy: verify intermediate results. Chain actions with verification between each step — never run the full pipeline blind.")

    add_concept_slide(prs, "The Golden Rule", [
        ("If you don't understand what Claude is about to do — don't let it do it.", False),
        ("Ask Claude to explain before approving any consequential action.", False),
        ("You are the operator. You are responsible for the output.", False),
        ("The supervision loop:", False),
        ("You specify → Claude clarifies → You confirm → Claude acts → You review → Claude adjusts → You approve", True),
        ("You are at every decision point. Claude executes and reports back.", False),
    ], note="If Claude says 'I'll clean up the database' and you don't understand the change — ask first. Review before execution. Always.")

    add_check_on_learning(prs,
        "Claude has been running a long session and proposes to 'clean up the database by removing duplicate records.' "
        "It sounds reasonable.\n\n"
        "What questions do you ask before approving this action?\n\n"
        "What could go wrong if you approve without reviewing?")

    add_hands_on(prs, "Operator Posture and Supervision", [
        "Give Claude a task with a vague brief: 'Improve my code.' Read what it produces.",
        "Give the same task with a specific brief: name the file, the goal, the constraint, what success looks like.",
        "Compare. Which result could you review and approve without asking follow-up questions?",
        "Label your vague brief: which supervision failure mode does it most resemble?",
    ], "[IMAGE: Two mission briefs — one handwritten with no detail, one full 5-paragraph OPORD]")

    add_section_summary(prs, "Operator Posture and Supervision", [
        "Active supervision: you review at each step and intervene before consequential actions execute.",
        "Three failure modes: over-trust, under-involvement, automation fallacy.",
        "The agent is a motivated PFC — capable and eager, but needs clear orders and oversight.",
        "You own the output. The capability does not transfer the accountability.",
    ])

    # --- SECTION 6 ---
    add_section_header(prs, "06", "Cost-Consciousness\nas a Core Skill",
        "Efficiency is not optional — it is the difference between a sustainable workflow and one that burns the budget on noise.",
        "[IMAGE: Logistics NCO tracking supply expenditures against a limited budget line]")

    add_concept_slide(prs, "The Three Cost Dimensions", [
        ("Dimension 1: Token Cost", False),
        ("Every input token costs money. Every output token costs more. 100k session = 10x a 10k session.", True),
        ("Dimension 2: Time Cost", False),
        ("Longer conversations take longer. More tool calls take longer. Dead time adds up.", True),
        ("Dimension 3: Iteration Cost", False),
        ("10 iterations to get it right = 10x the tokens. One clear prompt beats five vague ones.", True),
        ("Efficient communication is the gate to everything else.", False),
    ])

    add_concept_slide(prs, "Five Cost-Conscious Behaviors", [
        ("1. Write clear, structured prompts first — not verbose ones.", False),
        ("2. Use tools to read files instead of pasting large documents into the prompt.", False),
        ("3. Iterate tightly — review immediately, not a day later.", False),
        ("4. Know when to ask Claude vs. when to solve it yourself.", False),
        ("Use Claude: architectural decisions, complex logic, unfamiliar code.", True),
        ("Use search: syntax lookups, simple API questions, quick references.", True),
        ("5. Plan before you build. 5 minutes planning often saves 50 minutes of iteration.", False),
    ], note="One clear prompt beats five vague ones every time. Iteration is where cost hides — not in the prompt itself.")

    add_check_on_learning(prs,
        "You need to debug a slow API endpoint.\n\n"
        "You have the endpoint's source file, the database schema, and the server logs.\n\n"
        "What is the most cost-efficient way to give Claude what it needs?\n\n"
        "What would the expensive version look like?")

    add_hands_on(prs, "Cost-Consciousness", [
        "Write a prompt you would normally send to Claude for help with a task.",
        "Apply cost-conscious principles: cut repetition, add structure, narrow the scope, specify output format.",
        "Count words before and after. Is the revised version clearer and shorter?",
        "Submit the revised prompt. Did you get a usable first response, or still need multiple follow-ups?",
    ], "[IMAGE: Soldier editing a draft operations order — cutting lines, tightening the mission statement]")

    add_section_summary(prs, "Cost-Consciousness", [
        "Three cost dimensions: token cost, time cost, iteration cost.",
        "Structured prompts reduce iterations — which is where most cost actually hides.",
        "Use tools to read files instead of pasting large documents.",
        "Plan before you build. Five minutes planning saves fifty minutes of iteration.",
    ])

    # --- WORKED EXAMPLES ---
    add_section_header(prs, "07", "Worked\nExamples",
        "All six mental models in action on scenarios you will actually face.",
        "[IMAGE: After-action review board — soldiers analyzing what went right and wrong]")

    add_example_slide(prs, "Example 1: The Harness in Action",
        "Scenario: Intelligence summary from raw source documents",
        ["WITHOUT THE HARNESS:",
         '  You: "Write an intel summary on adversary logistics activity."',
         "  Claude generates from training data. Two claims unverifiable. Brief gets pulled.",
         "",
         "WITH THE HARNESS:",
         '  You: "Write an intel summary. Read source docs in /reports/AO-North/."',
         "  Claude lists 14 files, reads 11 with relevant content, excludes 3.",
         "  Returns summary with citations — every claim is traceable.",
         "",
         "First approach — you cannot defend the output.",
         "Second approach — every claim traces to a source file."],
        "[IMAGE: Analyst comparing raw source documents to a final intelligence product]")

    add_example_slide(prs, "Example 2: Context Windows and Smart Sampling",
        "Scenario: Find all mentions of a grid coordinate across 500 field reports",
        ["BAD APPROACH:",
         "  Paste all 500 reports (100,000 tokens) into the chat.",
         "  Claude overwhelmed. Results incomplete. Coverage unverifiable.",
         "",
         "GOOD APPROACH:",
         '  "Search /reports/field/ for grid 38SMB4521. Return excerpts with filenames."',
         "  Claude: lists 500 files → grep finds 7 matches → reads those 7 files.",
         "  Returns 9 mentions across 7 reports. Complete and traceable.",
         "",
         "Context cost — Bad: 100,000 tokens.  Good: ~20,000 tokens."],
        "[IMAGE: Targeted database search vs. reading an entire stack of paper reports]")

    add_example_slide(prs, "Example 4: Operator Supervision Preventing a Bad Product",
        "Scenario: Drafting a SITREP from unstructured notes",
        ["WITHOUT SUPERVISION (over-trust):",
         '  You: "Draft a SITREP from my notes."',
         "  Claude infers a casualty figure not in your notes. You don't review. You submit.",
         "  Correction required. Trust degraded.",
         "",
         "WITH SUPERVISION (active oversight):",
         '  You: "Draft SITREP from /notes/0600-update.txt. Use only that file. Flag gaps."',
         "  Claude drafts, flags two required fields, traces every claim to source lines.",
         "  You fill the gaps with verified data. Defensible product."],
        "[IMAGE: NCO reviewing a SITREP line-by-line before initialing and passing it up]")

    # --- SUMMARY TABLE ---
    add_summary_table(prs, [
        ("The Harness",        "LLM + tools = agency",                  "Without tools, Claude is a chatbot. With tools, a collaborator."),
        ("Context Windows",    "Limited working memory",                 "Use tools to read on-demand. Don't cram everything into context."),
        ("Tokens as Currency", "Every token costs time and money",       "Concise, structured prompts save iterations and cost."),
        ("Tool Calls",         "Requests that execute outside context",  "Verification is possible. Hallucinations can be caught."),
        ("Operator Posture",   "You supervise, Claude executes",         "You are responsible. Review before committing. Intervene early."),
        ("Cost-Consciousness", "Efficiency is a core skill",             "Plan before building. Iterate tightly. Match model to mission."),
    ])

    # --- READINESS CHECK ---
    add_readiness_check(prs, [
        "I can explain what the harness is and why tools are required for agency",
        "I know the three parts of the engine-harness-operator stack and what each does",
        "I can describe what a context window is and what happens when tool results accumulate",
        "I understand what tokens cost and can identify the efficient vs. expensive approach",
        "I can explain what a tool call is and why it enables verification",
        "I know what active supervision means and can name the three ways it breaks down",
        "I can identify cost-conscious vs. wasteful prompting behavior",
    ])

    # --- END ---
    add_end_slide(prs, "Mental Models for\nAI-Assisted Operations", [
        "Read this module twice — once to understand, once to internalize.",
        "Complete the Section 8 exercise: identify the mental models in the transcript.",
        "In your first Claude Code session, call out which model is in play as you work.",
        "Build the habit. These are not concepts — they are habits of thought.",
    ])

    out = os.path.join(os.path.dirname(__file__), "mental-models.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
