"""
ACC Terminal Basics — The Terminal Slide Deck Builder — Army Cyber Dark Theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- TITLE ---
    add_title_slide(prs,
        "Terminal Basics:\nThe Terminal",
        "TERMINAL BASICS MODULE  |  PREREQUISITE",
        "[IMAGE: Soldier at a command post terminal typing a command — focused and deliberate]")

    # --- OVERVIEW ---
    add_overview_slide(prs,
        "The terminal is a typed way to give the computer the same instructions you used to give by clicking. "
        "You are learning these moves because the agent acts through the terminal — and you need to supervise.",
        "Self-paced (approximately 60 minutes)",
        [
            "Open the terminal and run basic commands with confidence",
            "Navigate the filesystem using pwd, ls, and cd",
            "Create, copy, and move files from the command line",
            "Use absolute and relative paths correctly",
            "Apply tab-completion and the up-arrow to reduce errors",
            "Use flags and --help to extend any command",
            "Stop a running command with Ctrl+C",
            "Explain why version control matters once an agent has write access",
        ])

    # =========================================================
    # SECTION 01 — What the Terminal Is
    # =========================================================
    add_section_header(prs, "01", "What the\nTerminal Is",
        "A typed instruction is the same as a click — but without the menu. It feels alien for about an hour.",
        "[IMAGE: Soldier sitting at a terminal for the first time — cursor blinking, ready for input]")

    add_concept_slide(prs, "The Terminal Is Not Hacking", [
        ("The terminal is a typed way to give the computer instructions — the same instructions a click gives.", False),
        ("The difference is precision: a click is constrained by what a menu shows you.", False),
        ("A command can express anything the machine understands.", False),
        ("The prompt line is where your typing goes.", False),
        ("Mac/Linux: ends with $    |    Windows PowerShell: ends with >", True),
        ("When you see that character, the terminal is ready for input.", True),
        ("Nothing in this module can hurt the machine — all commands are read-only and harmless.", False),
    ], note="Instructor: Name the frustration before a student says it. 'Some of you are wondering why an AI course spent time on the command line.' Then connect it: the agent acts through the terminal. The student who cannot read the terminal cannot supervise an agent using it.")

    add_concept_slide(prs, "Why the Terminal Matters Here", [
        ("The agent acts through the terminal.", False),
        ("Every file read, every command run, every edit made — the agent does it through a terminal interface.", False),
        ("You need to understand the terrain it operates on.", False),
        ("You do not need to be an expert.", False),
        ("You need enough fluency to read what the agent is doing and intervene when it goes wrong.", True),
        ("These are the exact moves an agent makes on your machine.", False),
        ("You are learning them so you can supervise.", False),
    ])

    add_example_slide(prs, "Your First Command",
        "Scenario: Running date to confirm the terminal is alive",
        [
            "Open the terminal:",
            "  Windows: search 'PowerShell' or 'Terminal' in the Start menu",
            "  Mac: search 'Terminal' in Spotlight",
            "",
            "Mac / Linux:      type  date  and press Enter",
            "Windows PowerShell:  type  Get-Date  and press Enter",
            "",
            "The terminal printed the current date and returned the prompt.",
            "That is a command.",
            "",
            "You gave the machine an instruction. It executed. The prompt returned.",
            "That is the full loop — and everything else is a variation of it.",
        ],
        "[IMAGE: Soldier's screen showing a date output in a terminal window — first successful command]")

    add_check_on_learning(prs,
        "You typed a command and saw a response.\n\n"
        "What is the difference between that and clicking a button in an app?\n\n"
        "What stayed the same?\n\n"
        "What does a terminal actually do — in one sentence?")

    add_hands_on(prs, "What the Terminal Is", [
        "Open the terminal on your machine. Sit with the blinking cursor for ten seconds. Name the feeling.",
        "Mac/WSL: run  date, whoami, pwd, ls, echo \"hello\"  — one at a time, read each output.",
        "Windows PowerShell: run  Get-Date, $env:USERNAME, Get-Location, Get-ChildItem, Write-Output \"hello\"",
        "Close the terminal and open it again. Re-run one command from memory without looking at the list.",
        "You opened it, ran five things, and closed it. The anxiety has a smaller surface area now.",
    ], "[IMAGE: Soldier running five commands in sequence at a terminal — reading each output before continuing]")

    add_section_summary(prs, "What the Terminal Is", [
        "The terminal accepts typed instructions and passes them to the OS — the same as clicking, without the menu.",
        "The prompt ($  or  >) means the terminal is ready. The prompt returning means the command finished.",
        "The agent acts through the terminal. You learn it to supervise.",
        "Nothing in this section can damage the machine — all commands are read-only.",
    ])

    # =========================================================
    # SECTION 02 — Navigating: Where Am I, What Is Here
    # =========================================================
    add_section_header(prs, "02", "Navigating:\nWhere Am I",
        "Three commands cover most navigation — pwd, ls, cd — and knowing them cold means you can never stay lost.",
        "[IMAGE: Soldier using a map and compass to confirm position before moving — deliberate orientation]")

    add_concept_slide(prs, "The Three Navigation Commands", [
        ("pwd — Print Working Directory. Your current location in the tree.", False),
        ("Any time you feel lost, type pwd. You can never be permanently lost if you can always ask where you are.", True),
        ("ls — List. Everything in the current folder.", False),
        ("Add  -la  (Mac/WSL) for details including hidden files.", True),
        ("cd foldername — Change Directory. Move into a folder.", False),
        ("cd .. — Move up one level. Two dots always means 'the folder above this one.'", True),
        ("These three commands replace every click you made in The Machine module.", False),
    ], note="Reps over coverage. Same three commands, many short repetitions, spread across several days. Do not rush to the next section. One capability at a time — reps build muscle memory.")

    add_example_slide(prs, "Walking the Tree",
        "Scenario: Navigate the same filesystem you explored by clicking — now typed",
        [
            "pwd              # See where you are",
            "ls               # See what is here",
            "cd Documents     # Move into Documents",
            "pwd              # Confirm you moved",
            "ls               # See what is in Documents",
            "cd ..            # Move back up one level",
            "pwd              # Confirm you are back",
            "",
            "Walk this sequence on your own machine.",
            "These are the exact same moves from The Machine module — same tree, typed.",
            "",
            "pwd is your reset button. Run it any time you feel lost.",
        ],
        "[IMAGE: Soldier navigating a directory tree on screen — each cd command confirmed with pwd]")

    add_check_on_learning(prs,
        "You are in the terminal and have no idea where you are in the filesystem.\n\n"
        "What is the first command you run — and why that one first?\n\n"
        "You typed  cd Documents  and then  pwd.\n\n"
        "What does pwd tell you, and how does it prove you moved?")

    add_hands_on(prs, "Navigating: Where Am I, What Is Here", [
        "Open the terminal. Run  pwd. Read the path. Find that same location in File Explorer.",
        "Run  ls. Identify three items by name.",
        "Run  cd  to move into one of those items.",
        "Run  pwd  and  ls  again — confirm where you are and what is here.",
        "Run  cd ..  to go back up. Run  pwd  to confirm.",
        "Repeat this cycle until it feels routine.",
    ], "[IMAGE: Soldier alternating between File Explorer and terminal — matching what they see in each view]")

    add_section_summary(prs, "Navigating: Where Am I, What Is Here", [
        "pwd prints your current location — run it any time you feel lost.",
        "ls lists the contents of the current folder.",
        "cd foldername moves in;  cd ..  moves up one level.",
        "These three commands replace every navigation click from The Machine module.",
    ])

    # =========================================================
    # SECTION 03 — Acting: Make, Move, Copy
    # =========================================================
    add_section_header(prs, "03", "Acting: Make,\nMove, Copy",
        "Navigation gets you to the position. Now you act on it — and confirm every action before moving on.",
        "[IMAGE: Soldier creating a field record and then making a copy for the next echelon — deliberate and verified]")

    add_concept_slide(prs, "The Four Action Commands", [
        ("mkdir foldername — Make a new folder.", False),
        ("touch filename.txt (Mac/WSL) or New-Item filename.txt (PowerShell) — Make a new empty file.", False),
        ("cp source destination — Copy a file.", False),
        ("mv source destination — Move or rename a file.", False),
        ("Every action you take, you confirm. Run ls after every command.", False),
        ("Typos create new files — the command does not error, it makes a second file with the wrong name.", False),
        ("No deletion yet. Make, copy, and move only — deletion comes later, behind heavy framing.", False),
    ], note="The verify habit starts here. 'Trust the action, not the narration' runs through all agentic work. A student who verifies every terminal action will verify every agent action later.")

    add_example_slide(prs, "Build Something From Scratch",
        "Scenario: Create a project folder, a file, a copy, and a renamed version — confirming each step",
        [
            "mkdir project",
            "cd project",
            "touch notes.txt          # (Mac/WSL)  or  New-Item notes.txt  (PowerShell)",
            "ls                       # Verify notes.txt is there",
            "mkdir drafts",
            "cp notes.txt drafts/notes-copy.txt",
            "ls drafts                # Verify the copy landed",
            "mv notes.txt notes-v1.txt",
            "ls                       # Verify the rename happened",
            "",
            "After each command, run ls to confirm the change happened.",
            "That is the verify-after-acting habit — build it now.",
        ],
        "[IMAGE: Soldier running a command then immediately checking with ls — deliberate confirmation after each action]")

    add_check_on_learning(prs,
        "You run  mv report.txt final-report.txt  and then  ls.\n\n"
        "What are you checking for?\n\n"
        "You ran  mv notes.txt notes-final.txt  but mistyped the destination. "
        "What does  ls  show you — and what would tip you off to the error?")

    add_hands_on(prs, "Acting: Make, Move, Copy", [
        "Make a folder called  terminal-practice  and navigate into it.",
        "Create a file called  day1.txt.",
        "Copy it to  day1-backup.txt.",
        "Rename  day1.txt  to  day1-v1.txt.",
        "Run  ls  after each step to verify the change.",
        "You should end with two files: day1-v1.txt and day1-backup.txt.",
    ], "[IMAGE: Soldier's terminal showing a clean folder with correctly named files — every step confirmed]")

    add_section_summary(prs, "Acting: Make, Move, Copy", [
        "mkdir, touch/New-Item, cp, mv — four commands cover most day-to-day file work.",
        "Run ls after every action. Confirm the change happened before moving on.",
        "Typos do not error — they create misnamed files. ls is the check.",
        "No deletion yet. That comes later, with heavy framing.",
    ])

    # =========================================================
    # SECTION 04 — Paths, Tab-Completion, and the Up-Arrow
    # =========================================================
    add_section_header(prs, "04", "Paths, Tab,\nand Up-Arrow",
        "A path is a grid coordinate — precise and unforgiving. Tab-completion and the up-arrow are your accuracy tools.",
        "[IMAGE: Soldier calling in a grid coordinate with precision — no rounding, no approximation]")

    add_concept_slide(prs, "Absolute vs. Relative Paths", [
        ("An absolute path starts from the root of the filesystem — it works from anywhere.", False),
        ("Windows: C:\\Users\\YourName\\Documents\\report.txt", True),
        ("Mac/Linux: /home/yourname/documents/report.txt", True),
        ("A relative path starts from wherever you currently are — it only works from the right starting position.", False),
        ("documents/report.txt  or  ../otherfolder/file.txt", True),
        ("A path is a grid coordinate: precise, unambiguous, and unforgiving of a wrong character.", False),
        ("A space or a capital letter in the wrong place either fails or hits the wrong target.", False),
    ], note="A student annoyed that a capital letter breaks a path is learning the right thing. Name it: 'Yes, the space matters. The capital letter matters. That is why tab-completion exists.'")

    add_concept_slide(prs, "Tab-Completion and the Up-Arrow", [
        ("Tab-completion: start typing a name and press Tab — the terminal finishes it.", False),
        ("Two matches? Press Tab twice to see all options.", True),
        ("Use it on every path, every time. Tab-completion makes typos structurally impossible for the completed portion.", True),
        ("The up-arrow: cycles through previous commands.", False),
        ("Edit a recalled command and re-run it instead of retyping from scratch.", True),
        ("Never finish typing a long path by hand if tab-completion can do it.", False),
        ("These are not tricks — they are standard operating procedure.", False),
    ])

    add_check_on_learning(prs,
        "Why does tab-completion reduce errors in terminal commands?\n\n"
        "You are deep in a nested folder and need to copy a file to a location three levels up.\n\n"
        "Would you use an absolute or a relative path — and why?\n\n"
        "What is the risk of typing a long path entirely by hand?")

    add_hands_on(prs, "Paths, Tab-Completion, and the Up-Arrow", [
        "Navigate to a folder using the full absolute path — type it by hand once to feel the length.",
        "Navigate to the same folder using tab-completion. Press Tab after each segment of the path.",
        "Run  ls  in a folder with several items. Start typing one item's name, press Tab, watch it complete.",
        "Run a command. Press the up-arrow. Edit one character and re-run it.",
    ], "[IMAGE: Soldier typing a long path with tab-completion — terminal finishing the name precisely]")

    add_section_summary(prs, "Paths, Tab-Completion, and the Up-Arrow", [
        "Absolute paths work from anywhere. Relative paths only work from the correct starting position.",
        "A path is a grid coordinate — spaces and capitalization are load-bearing.",
        "Tab-completion finishes names from what actually exists — the completed portion cannot contain a typo.",
        "The up-arrow recalls previous commands. Edit and re-run instead of retyping.",
    ])

    # =========================================================
    # SECTION 05 — Commands, Flags, and Knowing When to Stop
    # =========================================================
    add_section_header(prs, "05", "Flags and\nKnowing When to Stop",
        "Most commands take flags that change what they do. --help shows all options. Ctrl+C stops anything.",
        "[IMAGE: Soldier reading a quick-reference card before operating equipment — checking options before acting]")

    add_concept_slide(prs, "Flags Change What a Command Does", [
        ("Flags are options you add to a command to modify its behavior.", False),
        ("Short form: -l    Long form: --long-output", True),
        ("ls           # List files, default view", False),
        ("ls -l        # List files, detailed view", True),
        ("ls -la       # List files, detailed, including hidden files", True),
        ("Same command. Different flags. Different output.", False),
        ("You do not memorize flags. You look them up.", False),
    ], note="The concept of a flag and how to look one up is the lesson. The skill is knowing flags exist and knowing how to find the one you need — not memorizing any specific flag.")

    add_concept_slide(prs, "Finding Flags and Stopping Commands", [
        ("ls --help  (Mac/Linux)  or  Get-Help ls  (PowerShell) — shows all flags and what they do.", False),
        ("Read the --help output top to bottom. Find the flag you need; ignore the rest.", True),
        ("It looks overwhelming the first time. Find one flag. Use it. Done.", True),
        ("The prompt is the signal: when a command finishes, the prompt returns.", False),
        ("If the cursor blinks with no prompt, the command is still running.", True),
        ("Ctrl+C stops a running command and returns the prompt.", False),
        ("It does not delete anything — it stops execution and hands control back to you.", True),
    ])

    add_example_slide(prs, "Stopping a Running Command",
        "Scenario: Demonstrate that you can stop anything you start",
        [
            "Run:   ping google.com",
            "       (Windows: ping -t google.com)",
            "",
            "Watch it run indefinitely.",
            "",
            "Press:  Ctrl+C",
            "",
            "The command stops. The prompt returns.",
            "",
            "You just demonstrated that you can stop anything you start.",
            "",
            "This is not just a terminal skill — it is the same reflex",
            "you use when an agent starts doing something unexpected.",
        ],
        "[IMAGE: Soldier's hand hovering over Ctrl+C — ready to stop a process at any moment]")

    add_check_on_learning(prs,
        "You run a command and nothing happens for thirty seconds — the cursor just blinks.\n\n"
        "What are the two possibilities, and how do you handle each one?\n\n"
        "You need to run  ls  but want to see hidden files. You do not know the flag.\n\n"
        "What do you do — and what do you NOT do?")

    add_hands_on(prs, "Commands, Flags, and Knowing When to Stop", [
        "Run  ls --help  (Mac/Linux) or  Get-Help ls  (PowerShell). Read the first ten lines without panic.",
        "Find the flag that shows hidden files in the --help output. Run  ls  with that flag.",
        "Run  ping google.com  (or  ping -t google.com  on Windows). Let it run for five seconds.",
        "Press  Ctrl+C. Confirm the prompt returned.",
    ], "[IMAGE: Soldier reading --help output — scanning for a specific flag and ignoring the rest]")

    add_section_summary(prs, "Commands, Flags, and Knowing When to Stop", [
        "Flags modify command behavior. Short form: -x    Long form: --option",
        "Use --help (or Get-Help) to find any flag — do not memorize, look it up.",
        "The prompt returning means the command finished. No prompt means still running.",
        "Ctrl+C stops anything. It does not delete — it returns control to you.",
    ])

    # =========================================================
    # SECTION 06 — Why Version Control Exists
    # =========================================================
    add_section_header(prs, "06", "Why Version\nControl Exists",
        "When an agent can change your files, you want a logbook of every change — with the ability to rewind any of them.",
        "[IMAGE: NCO reviewing a detailed change log — tracking every modification made during an operation]")

    add_concept_slide(prs, "The Logbook Mental Model", [
        ("Every time a file changes, version control records a snapshot: what changed, when, and why.", False),
        ("Stored locally on your machine. Every snapshot can be rewound.", False),
        ("Two layers:", False),
        ("Local logbook — your machine keeps the record. Rewind to any snapshot without the internet.", True),
        ("Remote copy — the same history synchronized to a cloud server. Backup and collaboration.", True),
        ("Why software people seem to remember everything: they do not. They have the log.", False),
        ("'What did this file look like last Tuesday?' is one command — not a memory challenge.", False),
    ], note="Resist teaching git commands here. The concept is the prerequisite. If a student asks 'how do I actually do this,' tell them: 'You will do it in the next session. Today you need to know why.'")

    add_concept_slide(prs, "Why This Matters Once an Agent Can Edit", [
        ("You just learned to create, move, and rename files from the command line.", False),
        ("An agent will do those same operations — at speed, on multiple files, sometimes without asking.", False),
        ("Version control is what lets you see every change and undo any of them.", False),
        ("Without version control: batch edit goes wrong, you are guessing which files changed.", False),
        ("With version control: you see exactly what changed, when, and why — and you rewind.", False),
        ("This is not optional once you give an agent write access to your files.", False),
    ], note="No commands today. The concept loads the reason. The commands come when you work with git directly in the next session.")

    add_check_on_learning(prs,
        "An agent just ran a batch edit on 40 files.\n\n"
        "You look at one and it is not right.\n\n"
        "Without version control, what are your options?\n\n"
        "With version control, what changes?\n\n"
        "Why does version control matter specifically once an agent can edit your files?")

    add_hands_on(prs, "Why Version Control Exists", [
        "Think about a file you have edited multiple times over the past month.",
        "Ask yourself: 'If I needed to see what it looked like three weeks ago, could I?'",
        "Ask yourself: 'If something went wrong, what would I lose?'",
        "That gap — between what you want to recover and what you currently can — is what version control fills.",
        "No terminal commands today. Load the reason. The commands come next session.",
    ], "[IMAGE: Soldier looking at a stack of handwritten drafts with no version numbers — unable to tell which is current]")

    add_section_summary(prs, "Why Version Control Exists", [
        "Version control records a snapshot of every change: what, when, and why — and lets you rewind any of them.",
        "Two layers: local logbook (always available) and remote copy (backup + collaboration).",
        "Once an agent has write access, version control is not optional — it is the safety net.",
        "No git commands today. The concept is the prerequisite. Commands come next session.",
    ])

    # =========================================================
    # FINAL READINESS CHECK
    # =========================================================
    add_readiness_check(prs, [
        "I can open the terminal and find the prompt line",
        "I ran at least five commands and read each output",
        "I can navigate the filesystem with pwd, ls, and cd without looking at a reference",
        "I run ls after every action to verify the change happened",
        "I can write both an absolute and a relative path to the same file",
        "I use tab-completion on every path — it is a habit, not an option",
        "I can use --help to find a flag I do not know",
        "I can stop a running command with Ctrl+C without hesitation",
        "I can explain what version control is and why it matters once an agent has write access",
    ])

    # =========================================================
    # END SLIDE
    # =========================================================
    add_end_slide(prs, "Terminal Basics:\nThe Terminal", [
        "Run each hands-on exercise from memory — no peeking at the slide.",
        "Walk the file tree by clicking (The Machine) and then by typing — compare the experience.",
        "In your next Claude Code session, read the terminal output the agent generates and narrate what it is doing.",
        "Move to the next module: Terminal Basics — Git and Version Control.",
    ])

    out = os.path.join(os.path.dirname(__file__), "terminal-terminal.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
