"""
ACC Bedrock — Personalizing Your AI Slide Deck Builder — Army Cyber Dark Theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- TITLE ---
    add_title_slide(prs,
        "Personalizing\nYour AI",
        "BEDROCK MODULE  |  20 MINUTES",
        "[IMAGE: Soldier configuring a workstation at a cyber operations desk]")

    # --- OVERVIEW ---
    add_overview_slide(prs,
        "Most AI tools let you set persistent instructions so the model already knows who you are, how you communicate, "
        "and what to never do — before you type a single word. Personalization is a one-time investment that pays back on every session after.",
        "20 minutes (self-contained)",
        [
            "Locate the custom instruction or project settings in your AI tool of choice",
            "Write at least three persistent instructions covering identity, format, and prohibited behaviors",
            "Understand what to include — and what to keep out — of custom instructions",
            "Distinguish between custom instructions, memory, and projects as persistent context mechanisms",
            "Recognize the context window cost of custom instructions and how to manage it",
        ])

    # =========================================================================
    # SECTION 1 — Default Settings Are a Starting Point, Not a Destination
    # =========================================================================
    add_section_header(prs, "01",
        "Default Settings Are\na Starting Point",
        "Without customization, every chat starts cold. A one-time setup pays back on every conversation after.",
        "[IMAGE: Soldier receiving no briefing before a mission — starting from zero every time]")

    add_concept_slide(prs, "Why Personalization Matters", [
        ("Without customization, every new chat starts cold.", False),
        ("The model does not know who you are, how you communicate, or what context is always relevant.", False),
        ("Every session you either re-explain yourself or accept generic output.", False),
        ("Personalization is a one-time investment that pays back on every single conversation after.", False),
        ("Most major platforms expose persistent customization through custom instructions or project settings.", False),
        ("Instructions load at the start of every conversation — before you type anything.", False),
    ], note="Think of it as a standard SALUTE brief the model gets before every session. You write it once. It briefs itself.")

    add_concept_slide(prs, "Where to Find It: ChatGPT and Claude", [
        ("ChatGPT — Custom Instructions", False),
        ("Settings → Personalization → Custom Instructions", True),
        ("Two fields: 'What should ChatGPT know about you?' and 'How should it respond?'", True),
        ("Applied automatically to every new conversation.", True),
        ("Claude — Projects", False),
        ("Claude.ai → Projects → Create a Project → Set Instructions", True),
        ("One system prompt field — use it the same way.", True),
        ("Upload documents the model can reference across all conversations in that project.", True),
    ])

    add_concept_slide(prs, "What to Put In — and What to Keep Out", [
        ("PUT IN: your role and what you do day to day", False),
        ("PUT IN: how you prefer responses (length, format, tone)", False),
        ("PUT IN: things to always do — 'Lead with a one-sentence summary. Use bullet points.'", False),
        ("PUT IN: things to never do — 'No emojis. No padding. Do not apologize.'", False),
        ("PUT IN: recurring context — tools, team terminology, domain", False),
        ("KEEP OUT: sensitive information, PII, classified or controlled material of any kind", False),
        ("KEEP OUT: anything covered in the Data Handling module — instructions go to the platform's servers on every request", False),
    ], note="Security reminder: custom instructions are sent to the platform on every message. Nothing sensitive, nothing controlled, nothing above the system's authorization ceiling.")

    add_example_slide(prs,
        "Before and After Custom Instructions",
        "Scenario: Analyst asking for a document summary",
        [
            "WITHOUT CUSTOM INSTRUCTIONS:",
            '  You ask: "Summarize this document."',
            "  The model returns a five-paragraph response — neutral tone, hedging,",
            "  introduction, conclusion, multiple caveats. Generic.",
            "",
            "WITH CUSTOM INSTRUCTIONS:",
            '  Instructions set: "I\'m an analyst. Lead with BLUF. Use bullets.',
            '  No padding. Max 150 words."',
            '  You ask: "Summarize this document."',
            "  Same document. Same model. Direct, structured output that matches",
            "  how you actually communicate.",
        ],
        "[IMAGE: Two document summaries side by side — one verbose and generic, one clean and direct]")

    add_concept_slide(prs, "Two Cautions Before You Build", [
        ("Caution 1: Start with one thing.", False),
        ("Tell it one thing it gets wrong every time. Fix that first. Build from there.", True),
        ("A tight 50-word instruction that sticks beats a 1,000-word template the model ignores after three turns.", True),
        ("Caution 2: Custom instructions count against your context window.", False),
        ("They are loaded on every message. A very long instruction reduces the effective window for your actual conversation.", True),
        ("Keep custom instructions under 500 words.", True),
        ("If you need more persistent context, use a Project document instead of cramming it into the instruction field.", True),
    ], note="Instructor Note: Custom instruction UI locations change frequently. Verify the menu path before teaching — the concept is stable, the UI is not.")

    add_check_on_learning(prs,
        "You wrote a custom instruction that says 'always lead with a one-sentence summary.'\n\n"
        "On one response, the model ignores it.\n\n"
        "What are the two most likely causes?")

    add_hands_on(prs, "Setting Up Custom Instructions", [
        "Open your AI tool of choice. Find the custom instructions or project settings.",
        "Write at least three instructions:",
        "    — One about who you are and what you do",
        "    — One about how you like responses formatted",
        "    — One about something to never do",
        "Start a new chat. Ask a question you have asked before without custom instructions.",
        "Compare the two outputs. Did the model follow your instructions?",
    ], "[IMAGE: Soldier filling out a standard form at a workstation — building a reusable configuration]")

    add_section_summary(prs, "Default Settings Are a Starting Point", [
        "Every session starts cold without customization — the model knows nothing about you.",
        "Custom instructions load before every conversation. Write them once; pay back on every session.",
        "Include role, format preferences, always-do and never-do behaviors.",
        "Keep sensitive and controlled information out — instructions reach the platform's servers on every request.",
        "Keep instructions under 500 words to protect your effective context window.",
    ])

    # =========================================================================
    # SECTION 2 — Going Further: Projects, Memory, and Personas
    # =========================================================================
    add_section_header(prs, "02",
        "Projects, Memory,\nand Personas",
        "Beyond custom instructions: structured persistent context that turns a generic tool into one that already knows your work.",
        "[IMAGE: Soldier setting up a dedicated operations cell — organized, pre-configured for a specific mission type]")

    add_concept_slide(prs, "Claude Projects: A Dedicated Workspace", [
        ("A Project is a workspace with its own system prompt and document library.", False),
        ("Use a Project when:", False),
        ("You have a recurring task with the same context requirements every time", True),
        ("You want to upload reference documents the model can consult in every conversation", True),
        ("You need different configurations for different types of work", True),
        ("Create one project for analysis, one for writing, one for research.", False),
        ("Each has its own instructions and document set. Context never bleeds between project types.", False),
    ], note="Example: Your analysis Project has instructions set and your organization's style guide uploaded. You open the project and start working — context already loaded.")

    add_concept_slide(prs, "ChatGPT Memory: Reactive Persistence", [
        ("Memory records facts from your conversations and surfaces them in future sessions.", False),
        ("The model notes your name, role, preferences, and past decisions as they come up.", False),
        ("View, edit, and delete memory entries at any time: Settings → Personalization → Memory.", False),
        ("Memory is reactive — it records what comes up in conversation.", False),
        ("Custom instructions are proactive — you set them deliberately.", False),
        ("Use both: custom instructions for stable preferences, memory for things that evolve.", False),
    ], note="Instructor Note: Memory availability varies by subscription tier and has been rolled out at different times across regions. Verify current availability before teaching.")

    add_concept_slide(prs, "Personas: Named Configurations", [
        ("A persona is a named role, tone, and constraint set you can switch to by name.", False),
        ("Supported on some platforms and all agentic harnesses.", False),
        ("A well-defined persona means you do not re-explain the agent's role every session.", False),
        ("The persona carries the configuration — you start working immediately.", False),
        ("In agentic contexts: define the persona once, invoke it by name, begin the mission.", False),
    ])

    add_example_slide(prs,
        "A Working Project Setup",
        "Scenario: Analyst who does recurring document analysis",
        [
            "PROJECT INSTRUCTIONS:",
            '  "You are a plain-language editor supporting an analyst.',
            "  Lead with the main finding. No jargon.",
            "  Flag anything that requires verification before acting on it.",
            '  Max 200 words per response unless asked for more."',
            "",
            "PROJECT DOCUMENTS:",
            "  — Organization's style guide",
            "  — Glossary of domain terms",
            "  — One-page brief on the current project",
            "",
            "Every analysis conversation starts with that context already loaded.",
            "Open the project. Start working.",
        ],
        "[IMAGE: Analyst at a pre-configured workstation — everything staged before the mission begins]")

    add_check_on_learning(prs,
        "You create a Project for analysis work with detailed instructions.\n\n"
        "Three months later, your role changes.\n\n"
        "What do you need to update, and where?\n\n"
        "What happens if you do not update it?")

    add_hands_on(prs, "Projects, Memory, and Personas", [
        "If you use Claude: create a Project for one type of work you do repeatedly.",
        "    — Add instructions and at least one reference document to the Project.",
        "If you use ChatGPT: go to Settings → Personalization → Memory.",
        "    — Delete any outdated or wrong entries. Add a manual entry for something the model should always know.",
        "Start a new conversation inside the Project or with memory active.",
        "Compare how quickly the model orients to your work versus a cold-start session.",
    ], "[IMAGE: Soldier comparing two mission briefings — one starting from scratch, one pre-staged and ready]")

    add_section_summary(prs, "Projects, Memory, and Personas", [
        "Claude Projects: a dedicated workspace with instructions and uploaded documents for recurring work.",
        "ChatGPT Memory: reactive — records facts from conversations. Supplement custom instructions, not a replacement.",
        "Personas: named configurations that carry role and constraints without re-briefing each session.",
        "Use Projects for structured, recurring work. Use memory for evolving preferences.",
        "Keep project documents free of sensitive or controlled material — same rule as custom instructions.",
    ])

    # =========================================================================
    # READINESS CHECK
    # =========================================================================
    add_readiness_check(prs, [
        "I have found the custom instruction or project settings in my AI tool",
        "I have written at least three instructions: who I am, how I want responses, what to never do",
        "I ran a before/after comparison and noticed a measurable difference in output",
        "I did not include any sensitive or controlled information in my instructions",
        "I understand why custom instructions count against the context window and how to manage that",
        "I know the difference between custom instructions (proactive) and memory (reactive)",
        "I have set up or can describe how to set up a Project or named configuration for recurring work",
    ])

    # =========================================================================
    # END SLIDE
    # =========================================================================
    add_end_slide(prs,
        "Personalizing\nYour AI",
        [
            "Revisit your custom instructions after your first 10 sessions — trim what isn't working.",
            "Create at least one Project or named configuration for your most common AI use case.",
            "Apply the same security discipline here as everywhere: no sensitive data in persistent instructions.",
            "Build the habit: the model is only as well-briefed as the instructions you give it.",
        ])

    out = os.path.join(os.path.dirname(__file__), "bedrock-personalizing.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
