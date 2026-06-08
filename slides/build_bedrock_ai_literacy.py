"""
ACC Bedrock — AI Literacy Slide Deck Builder — Army Cyber Dark Theme
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_summary_table(prs, title, subtitle, rows, col_headers, col_x, col_w):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, subtitle)
    _footer(slide)

    ht = Inches(1.52)
    rh = Inches(0.42)

    hdr = _rect(slide, L, ht, R, rh, ELECTRIC_BLUE)
    _glow_blue(hdr, 4)
    for i, h in enumerate(col_headers):
        _tb(slide, h, col_x[i] + Inches(0.1), ht + Inches(0.06), col_w[i], rh,
            sz=11, bold=True, color=CYBER_BLACK)

    for ri, row in enumerate(rows):
        top = ht + rh + ri * Inches(0.83)
        _rect(slide, L, top, R, Inches(0.83), DARK_NAVY if ri % 2 == 0 else NAVY_MID)
        bar = _rect(slide, L, top, Pt(3), Inches(0.83), ELECTRIC_BLUE)
        _glow_blue(bar, 3)
        colors = [ELECTRIC_BLUE, WHITE, LIGHT_GRAY]
        for ci, val in enumerate(row):
            _tb(slide, val, col_x[ci] + Inches(0.12), top + Inches(0.08),
                col_w[ci] - Inches(0.15), Inches(0.8),
                sz=11, bold=(ci == 0), color=colors[ci])


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- TITLE ---
    add_title_slide(prs,
        "Bedrock:\nAI Literacy",
        "PREREQUISITE MODULE  |  SELF-PACED",
        "[IMAGE: Soldier at a command post terminal reviewing AI-generated analysis on a dual-monitor workstation]")

    # --- OVERVIEW ---
    add_overview_slide(prs,
        "An LLM predicts text from learned patterns — not rules, not memory, not live data. "
        "That single fact explains its extraordinary capability and its exact failure modes.",
        "Self-paced (6 sections)",
        [
            "State what an LLM actually does and why it is not the same as knowing",
            "Explain tokens and context windows in plain terms",
            "Name the four failure modes and produce a hallucination on demand",
            "Write a deliberate, structured prompt using the four-element framework",
            "Identify how AI is delivered and paid for — and why model choice has a cost",
            "Apply the classify-before-you-paste rule to content you handle every day",
        ])

    # =========================================================================
    # SECTION 01 — What an LLM Actually Is
    # =========================================================================
    add_section_header(prs, "01", "What an LLM\nActually Is",
        "The model predicts what comes next — not rules, not memory, not a lookup table. "
        "Fluent output is not the same as correct output.",
        "[IMAGE: Soldier reading a training manual — learning patterns, not memorizing rules]")

    add_concept_slide(prs, "Prediction, Not Lookup", [
        ("The model reads a sequence of text and predicts what comes next — one chunk at a time.", False),
        ("No lookup table. No truth check. No rules written by a programmer.", False),
        ("Given everything written so far: what is the statistically most likely next piece?", False),
        ("Trained, not programmed.", False),
        ("No one wrote 'Paris is the capital of France.' The model trained on enough text that it became the likely completion.", True),
        ("Upside: it generalizes across almost any domain.", True),
        ("Failure mode: if training data was wrong or thin in some area, the model learned those patterns too.", True),
    ], note="The only accurate verb is 'predicts.' The model does not know, want, think, or lie. Catch yourself using those words and swap in 'predicts.'")

    add_concept_slide(prs, "The LLM Is the Engine", [
        ("By itself, the LLM is a brain in a jar.", False),
        ("Capable of reasoning in text — incapable of acting on its own.", False),
        ("No body. No memory between chats. No live connection to the world.", False),
        ("Fluent, confident output is not the same as correct output.", False),
        ("Both are true simultaneously — that is the working mental model.", False),
    ], note="Warning: Anthropomorphism is a wired cognitive shortcut. 'It knew the answer' will feel natural and will mislead you at every later step. The accurate frame is: it predicted a statistically likely completion.")

    add_check_on_learning(prs,
        "The model just gave you a confident-sounding answer.\n\n"
        "What would it take to verify it? Where would you check?\n\n"
        "A language model was never given a rule that says 'Paris is the capital of France.' "
        "How does it produce that answer?\n\n"
        "Think through: training on patterns vs. looked up vs. hard-coded rule — which is it, and what does that mean for accuracy?")

    add_hands_on(prs, "What an LLM Actually Is", [
        "Type an incomplete sentence on a topic you know well and submit it. Read what the model predicts.",
        "Add the word 'definitely' or 'obviously' to the same sentence and submit again. Does the completion change?",
        "Ask the model a factual question you already know the answer to. Check the output for accuracy.",
        "Ask the exact same question in a new chat. Is the answer identical, or slightly different?",
        "Ask yourself: 'How would I verify any of these outputs if I did not already know the answer?'",
    ], "[IMAGE: Soldier testing radio frequencies — observing how small input changes affect output]")

    add_section_summary(prs, "What an LLM Actually Is", [
        "The model predicts what comes next — statistically likely, not looked up or rule-based.",
        "Trained, not programmed: it learned patterns from text, including the wrong ones.",
        "Fluent and confident output is not the same as correct output.",
        "The only accurate verb is 'predicts' — not knows, thinks, or remembers.",
    ])

    # =========================================================================
    # SECTION 02 — Tokens and Context
    # =========================================================================
    add_section_header(prs, "02", "Tokens and\nContext",
        "The model reads in tokens — roughly 0.75 of a word each — and can only hold a fixed number at once. "
        "Fill that limit and earlier content starts to fall out.",
        "[IMAGE: Tactical whiteboard filling up with overlapping mission graphics — running out of space]")

    add_concept_slide(prs, "What Is a Token?", [
        ("The model does not read words — it reads chunks of text called tokens.", False),
        ("In English prose, one token is roughly three-quarters of a word.", False),
        ("Code and non-English text tokenize less efficiently — more tokens per word.", False),
        ("You do not need to count tokens precisely.", False),
        ("The instinct to keep input focused is the deliverable from this section.", False),
    ])

    add_concept_slide(prs, "The Context Window", [
        ("Every model has a hard ceiling on tokens it can hold in a single conversation.", False),
        ("That ceiling is the context window — a whiteboard with finite space.", False),
        ("New writing crowds out old writing once it fills.", False),
        ("Current approximate sizes (verify before teaching — these change with releases):", False),
        ("Claude: 200,000 tokens", True),
        ("GPT-4o: 128,000 tokens", True),
        ("Gemini 2.0: 1,000,000+ tokens (bigger is not always better — quality degrades in the middle)", True),
    ], note="What a full window looks like: the model starts hedging. It contradicts something it said confidently ten turns ago. It forgets a constraint you set at the start. That is not a bug — it is the whiteboard running out of space.")

    add_concept_slide(prs, "The Two Rules for Fresh Chats", [
        ("Start a fresh chat when:", False),
        ("Rule 1: The task has changed direction and earlier turns are dead weight.", True),
        ("Rule 2: The model is drifting — hedging on things it stated confidently before.", True),
        ("Models attend better to the beginning and end of their context window than to the middle.", False),
        ("Critical instructions buried mid-conversation may be partially ignored.", False),
        ("Proactive session management is a skill — not a sign that something is wrong.", False),
    ], note="The practical rule: start fresh before you need to, not after you notice the degradation. A ten-minute-old constraint buried in a long chat is weaker than a constraint in a fresh session.")

    add_check_on_learning(prs,
        "You set an important constraint at the start of a long conversation. "
        "Ten exchanges later, the model seems to have forgotten it.\n\n"
        "What happened, and what would you do differently next time?\n\n"
        "You are in a long chat and the model starts contradicting instructions you gave at the start. "
        "What is the most likely cause?\n\n"
        "A) The model changed its mind\n"
        "B) You used the wrong model for this task\n"
        "C) Earlier content has been crowded out of the context window\n"
        "D) The model is testing whether you notice")

    add_hands_on(prs, "Tokens and Context", [
        "Open a fresh chat and paste several paragraphs of dense text.",
        "Ask the model to summarize just the first paragraph.",
        "In the same chat, ask about something from the end of what you pasted.",
        "Start a new chat and ask the same question about the end. Compare the quality.",
        "Notice: does the model attend differently to content at the top vs. buried in the middle?",
    ], "[IMAGE: Soldier reviewing a long operation log — noting what has been forgotten or dropped off the bottom]")

    add_section_summary(prs, "Tokens and Context", [
        "The model reads in tokens — roughly 0.75 of an English word each.",
        "Every model has a context window — a finite whiteboard that fills up.",
        "When the window fills, earlier content falls out silently. Instructions drift.",
        "Start fresh when the task changes or the model begins to hedge on prior statements.",
    ])

    # =========================================================================
    # SECTION 03 — How LLMs Fail
    # =========================================================================
    add_section_header(prs, "03", "How LLMs\nFail",
        "The model will state false things with total confidence, give different answers to the same question, "
        "and know nothing past its training cutoff. None of this is a malfunction.",
        "[IMAGE: Soldier double-checking map coordinates before calling in a grid — verification before action]")

    add_concept_slide(prs, "Failure Mode 1: Hallucination", [
        ("The model has no truth-checking step.", False),
        ("It generates tokens that are statistically likely to follow prior context.", False),
        ("It cannot distinguish accurate training data from plausible completion.", False),
        ("Confident output and correct output are entirely unrelated.", False),
        ("This is expected behavior of the system — not a rare bug.", False),
    ], note="The goal is calibrated trust, not distrust. The model is extraordinarily capable AND it hallucinates. Both are true simultaneously.")

    add_concept_slide(prs, "Failure Modes 2, 3, and 4", [
        ("Confident-Wrong — hallucination is not always dramatic.", False),
        ("Subtly wrong dates, wrong statistics, plausible-sounding citations that do not exist.", True),
        ("Tone stays equally confident whether the claim is right or wrong.", True),
        ("Nondeterminism — temperature introduces variation by design.", False),
        ("Same prompt, different run, different result. Do not treat one output as 'the answer.'", True),
        ("Knowledge Cutoff — the model was trained to a point in time.", False),
        ("Ask about a recent event: it may refuse, guess, or confidently describe something that did not happen.", True),
    ], note="Build the verification reflex now. Anything that matters gets checked before you act on it — not because the model is bad at its job, but because this is how the system works.")

    add_check_on_learning(prs,
        "You just watched the model invent a source.\n\n"
        "What does that mean for the next time it gives you a fact you have not heard before?\n\n"
        "You run the same prompt twice in two separate chats and get different answers. "
        "What is the most accurate explanation?\n\n"
        "A) One of the chats had a longer context window\n"
        "B) The model updated itself between the two runs\n"
        "C) The model uses randomness by design — the same input does not guarantee the same output\n"
        "D) You phrased the prompt slightly differently without noticing")

    add_hands_on(prs, "How LLMs Fail", [
        "Ask the model for 5 peer-reviewed sources on a specific narrow topic — with author names, journal names, and publication years.",
        "Pick one citation and try to verify it exists. Search the journal's website or a database.",
        "Ask the same question in a new chat. Compare the two citation lists — are they identical?",
        "Ask about something you know happened recently. Watch how it responds to events past its cutoff.",
        "You have just demonstrated all four failure modes. Name which step showed which mode.",
    ], "[IMAGE: Analyst comparing two documents side by side — verifying source vs. model output]")

    add_section_summary(prs, "How LLMs Fail", [
        "Hallucination: the model completes plausibly, not accurately. No truth-checking step.",
        "Confident-Wrong: tone stays equally confident for right and wrong claims — never use confidence as a signal.",
        "Nondeterminism: same prompt, different run, different result — by design.",
        "Knowledge cutoff: confident answers about events past training may be wrong.",
    ])

    # =========================================================================
    # SECTION 04 — Prompt Engineering
    # =========================================================================
    add_section_header(prs, "04", "Prompt\nEngineering",
        "A vague prompt is a vague order — and a vague order gets a vague result from a junior who fills the gaps with assumptions.",
        "[IMAGE: NCO handing a clearly written OPORD to a soldier — specific orders, specific results]")

    add_concept_slide(prs, "The Four Elements of a Deliberate Prompt", [
        ("The model meets you halfway the moment you give it structure.", False),
        ("Without context, it guesses: who you are, what you need, what good looks like.", False),
        ("Role — tell the model who to be.", False),
        ("'You are a plain-language editor' produces different output than no framing at all.", True),
        ("Context — tell it what it needs to know that it cannot infer.", False),
        ("What is the task, who is the audience, what constraints apply.", True),
        ("Example — show it what good output looks like before asking for it.", False),
        ("One strong example beats three paragraphs of description.", True),
        ("Output Spec — tell it the format, length, and tone you want.", False),
    ])

    add_concept_slide(prs, "Before and After: The Same Model, Different Brief", [
        ("Weak prompt:", False),
        ('"Write me a summary."', True),
        ("Strong prompt:", False),
        ('"You are summarizing this for a senior leader who has two minutes. Pull out the three most important points. Use bullet points. No jargon. Max 100 words."', True),
        ("Same model. Different brief. Different result.", False),
        ("The iterative habit matters more than the perfect prompt.", False),
        ("A rough ask the model can build on beats silence. Read the output, tighten the next ask, repeat.", False),
    ], note="The four elements are a structure, not a script to recite. The durable skill is clarity — giving the model what it needs to not guess. How you deliver that varies by task.")

    add_concept_slide(prs, "When You Are Stuck", [
        ("If you do not know how to frame the prompt, ask the model to interview you.", False),
        ('"I need help with [task]. Ask me the questions you need answered before you start."', True),
        ("The model surfaces what context it is missing.", True),
        ("Answer the questions, then ask it to proceed.", True),
        ("The template trap: four elements are a structure, not a script.", False),
        ("Magic phrases are not the lesson. Clarity is the lesson.", False),
        ("Students who focus on tricks over structure hit a ceiling fast.", False),
    ])

    add_check_on_learning(prs,
        "What specifically changed between a one-line prompt and a structured prompt?\n\n"
        "Which of the four elements made the biggest difference for your task?\n\n"
        "You ask the model to 'write a report.' It produces something generic and too long. "
        "What is the most effective fix?\n\n"
        "A) Switch to a more powerful model\n"
        "B) Ask the same question again and hope for a better result\n"
        "C) Add role, context, an example of good output, and a length/format spec to the prompt\n"
        "D) Break the report into smaller pieces and ask for each separately")

    add_hands_on(prs, "Prompt Engineering", [
        "Pick a real task you need help with right now.",
        "Write a one-line version of the request. Submit it. Save the output.",
        "Add role, context, an example of what good looks like, and an output spec. Submit again.",
        "Compare the two outputs side by side. Which required fewer follow-up exchanges?",
        "Optional: if stuck, type 'I need help with [task]. Ask me the questions you need answered before you start.'",
    ], "[IMAGE: Two mission briefs side by side — one a wall of text, one a clean 5-paragraph OPORD format]")

    add_section_summary(prs, "Prompt Engineering", [
        "Four elements: Role, Context, Example, Output Spec — give the model what it needs to not guess.",
        "A vague prompt is a vague order. Structured prompts reduce the surface area for wrong assumptions.",
        "The iterative habit matters more than the perfect prompt — read, tighten, repeat.",
        "If stuck, ask the model to interview you before it starts.",
    ])

    # =========================================================================
    # SECTION 05 — How AI Is Delivered and Paid For
    # =========================================================================
    add_section_header(prs, "05", "How AI Is\nDelivered and\nPaid For",
        "The same model reaches you through different doors — flat-rate app, pay-per-token API, or your organization's key. "
        "Tokens cost money, and model choice has a price.",
        "[IMAGE: Supply sergeant reviewing a billing ledger against available budget — resource accountability]")

    add_concept_slide(prs, "Three Delivery Models", [
        ("Subscription App — flat monthly fee, model access included.", False),
        ("Example: ChatGPT Plus, Claude.ai Pro", True),
        ("Pay-Per-Token API — you pay per token sent and received.", False),
        ("Example: Anthropic API, OpenAI API", True),
        ("Bring-Your-Own-Key — your organization's API key foots the bill.", False),
        ("Example: Enterprise deployments, Claude Code", True),
        ("Know which model you are on before you start a long or expensive task.", False),
    ])

    add_concept_slide(prs, "Cloud vs. Local Delivery", [
        ("Most models run in the cloud — your input goes to a remote server, the response comes back.", False),
        ("Some smaller models can run entirely on your machine.", False),
        ("Local delivery: no connectivity required, data stays local.", False),
        ("Trade-off: local models are smaller and less capable than frontier cloud models.", False),
        ("The cost ladder — token cost scales with model size and reasoning effort.", False),
        ("Small, fast model: less per token. Large frontier model: more per token.", True),
        ("'Always use the biggest model' is wrong — match the model to the job.", True),
    ], note="Do not teach specific token prices — they change with every major release. Teach the structure: tokens cost money, bigger models cost more, match the model to the mission.")

    add_example_slide(prs, "Example: Match the Model to the Job",
        "Scenario: Selecting the right tier for each task",
        ["WRONG APPROACH:",
         "  Use the most powerful model for everything.",
         "  Like running a diesel generator to charge a phone — the capability is real, the waste is real.",
         "",
         "RIGHT APPROACH — Match the tier to the task:",
         "  Drafting a quick message: small, fast model.",
         "  Analyzing a complex document: larger model.",
         "  Writing and running code in an agentic environment: frontier model with tool use.",
         "",
         "Know your billing model before you start a long task on a pay-per-token plan."],
        "[IMAGE: Army equipment lineup — different tools selected for different mission requirements]")

    add_check_on_learning(prs,
        "You need to run an analysis task with a very long document and several back-and-forth exchanges.\n\n"
        "What factors determine the cost of that task, and how would you reduce it?\n\n"
        "Why is 'always use the most powerful model' not good advice?\n\n"
        "A) Powerful models make more mistakes than smaller ones\n"
        "B) Powerful models are slower and less reliable\n"
        "C) Larger models cost more per token — for simple tasks, the cost is not justified by the capability difference\n"
        "D) You can only access the most powerful models through a paid subscription")

    add_hands_on(prs, "How AI Is Delivered and Paid For", [
        "Open your chatbot's settings or account page. Find which model you are currently using. Note the name.",
        "Open anthropic.com/pricing (or your provider's equivalent). Find two models — one fast/small, one large/powerful. Read the price difference. Notice the structure, not the numbers.",
        "Identify which delivery model you are on: subscription app, pay-per-token API, or organization-managed key.",
        "Optional: paste a paragraph into a tokenizer tool. Count the tokens. Compare to the word count.",
        "Pick one task you might use AI for this week. Decide which model tier is appropriate and write one sentence explaining why.",
    ], "[IMAGE: Soldier cross-referencing a field manual with a digital screen — selecting the right reference for the task]")

    add_section_summary(prs, "How AI Is Delivered and Paid For", [
        "Three delivery models: subscription app, pay-per-token API, organization-managed key.",
        "Most models run in the cloud — local models trade capability for privacy and offline access.",
        "Token cost scales with model size. Bigger is not always better — match the model to the job.",
        "Know your billing model before starting any long or high-volume task.",
    ])

    # =========================================================================
    # SECTION 06 — Data Handling: What Never to Paste
    # =========================================================================
    add_section_header(prs, "06", "Data Handling:\nWhat Never to\nPaste",
        "Once you are pasting real work into a cloud-connected tool, what you paste matters. "
        "This rule does not expire under time pressure.",
        "[IMAGE: Soldier locking a classified document safe before using a civilian computer terminal]")

    add_concept_slide(prs, "Authorization Is a Property of the System", [
        ("Authorization is a property of the system — not the impressiveness of the tool.", False),
        ("A highly capable model on an unauthorized system does not become authorized because it is impressive.", False),
        ("The boundary is set by policy, not capability.", False),
        ("What must never go into an unauthorized system:", False),
        ("Personally identifiable information (PII): names combined with identifiers, SSNs, DOBs, home addresses", True),
        ("Sensitive, controlled, or classified material of any kind", True),
        ("Anything above the system's authorization ceiling", True),
    ], note="This applies to military and civilian audiences alike. For a military audience: state it explicitly, including the 'paraphrase and summarize' loophole — which does not exist. For civilian: medical records, financial data, HR files, client information.")

    add_concept_slide(prs, "The Default Posture and the Classify-Before-You-Paste Habit", [
        ("Default posture: when in doubt, do not paste.", False),
        ("Ask someone who can authorize it before you proceed.", False),
        ("Classify before you paste — before pasting anything into an AI tool, take two seconds:", False),
        ("What is this content?", True),
        ("Is this system authorized for it?", True),
        ("Build that pause until it is automatic.", True),
        ("This bright line does not expire — not at the end of this course, not under time pressure.", False),
    ], note="One careless paste in an unauthorized system is the kind of mistake that has real and lasting consequences. Default to no until you hear yes from someone who can authorize it. Do not soften this.")

    add_check_on_learning(prs,
        "You are under a deadline and want to paste a document into your AI tool to summarize it quickly. "
        "You are not sure whether the system is authorized for that content.\n\n"
        "What do you do?\n\n"
        "You have a borderline-sensitive document and an unauthorized but highly capable AI tool that would save hours of work. "
        "What is the correct call?\n\n"
        "A) Paste it — the efficiency gain justifies a judgment call\n"
        "B) Summarize the key points in your head first, then paste the summary\n"
        "C) Do not paste. Authorization does not change based on capability or time pressure. Ask before proceeding.\n"
        "D) Paste it, but delete the conversation immediately after")

    add_hands_on(prs, "Data Handling: What Never to Paste", [
        "No prompting today. Think about the last three things you pasted into an AI tool.",
        "For each one: was the system authorized for that type of content?",
        "Identify one category of content you work with regularly that you will never paste into an unauthorized tool.",
        "Write that category down. It is your personal bright line.",
        "Keep it. It does not expire.",
    ], "[IMAGE: Soldier reviewing a checklist before using a terminal — deliberate pause before action]")

    add_section_summary(prs, "Data Handling: What Never to Paste", [
        "Authorization is a property of the system — not the tool's capability or your time pressure.",
        "PII, sensitive, controlled, and classified material must never go into an unauthorized system.",
        "Default posture: when in doubt, do not paste. Ask first.",
        "Classify before you paste — build the two-second pause until it is automatic.",
    ])

    # =========================================================================
    # MODULE SUMMARY TABLE
    # =========================================================================
    col_x = [L, Inches(3.5), Inches(8.1)]
    col_w = [Inches(2.9), Inches(4.4), Inches(5.0)]
    add_summary_table(prs,
        "MODULE SUMMARY",
        "AI Literacy — Six Core Concepts",
        [
            ("What an LLM Is",        "Prediction from learned patterns",          "Fluent output is not correct output. The only verb is 'predicts.'"),
            ("Tokens and Context",    "Finite working memory — fill it and it fails", "Start fresh proactively. Instructions buried mid-chat drift."),
            ("How LLMs Fail",         "Four failure modes — all expected behavior",  "Verify anything that matters. Confidence is not accuracy."),
            ("Prompt Engineering",    "Role, Context, Example, Output Spec",         "Structure closes the gap. A vague prompt is a vague order."),
            ("Delivery and Cost",     "Three models, token pricing, model tiers",    "Match the model to the job. Know your billing model."),
            ("Data Handling",         "Authorization is a property of the system",   "When in doubt, do not paste. The bright line does not expire."),
        ],
        ["CONCEPT", "CORE IDEA", "WHY IT MATTERS"],
        col_x, col_w)

    # =========================================================================
    # READINESS CHECK
    # =========================================================================
    add_readiness_check(prs, [
        "I can state in one sentence what an LLM does — using only the verb 'predicts'",
        "I can explain 'trained not programmed' and what it means for accuracy",
        "I can define a token and explain what happens when a context window fills",
        "I can name the four failure modes and have produced a hallucination with my own hands",
        "I can write a deliberate prompt using role, context, example, and output spec",
        "I can name the three delivery models and explain why model choice has a cost",
        "I know what must never go into an unauthorized system and I have identified my personal bright line",
    ])

    # =========================================================================
    # END SLIDE
    # =========================================================================
    add_end_slide(prs, "Bedrock:\nAI Literacy", [
        "Read this module twice — once to understand, once to internalize.",
        "Produce a hallucination with your own hands before moving to the next module.",
        "Write a before/after prompt comparison on a real task you currently have.",
        "Identify your personal bright line for content you handle — and write it down.",
    ])

    out = os.path.join(os.path.dirname(__file__), "bedrock-ai-literacy.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
