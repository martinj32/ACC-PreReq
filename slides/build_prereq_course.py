"""
ACC Prerequisite Course Slide Deck Builder — Army Cyber Dark Theme
Source: docs/prereq-course/course-design.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_summary_table(prs, rows):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE SUMMARY", "Eight Modules — Zero to Agentic-Ready")
    _footer(slide)

    col_x = [L, Inches(3.5), Inches(8.1)]
    col_w = [Inches(2.9), Inches(4.4), Inches(5.0)]
    ht = Inches(1.52)
    rh = Inches(0.42)

    hdr = _rect(slide, L, ht, R, rh, ELECTRIC_BLUE)
    _glow_blue(hdr, 4)
    for i, h in enumerate(["MODULE", "FOCUS", "OUTCOME"]):
        _tb(slide, h, col_x[i] + Inches(0.1), ht + Inches(0.06), col_w[i], rh,
            sz=11, bold=True, color=CYBER_BLACK)

    for ri, (module, focus, outcome) in enumerate(rows):
        top = ht + rh + ri * Inches(0.72)
        _rect(slide, L, top, R, Inches(0.72), DARK_NAVY if ri % 2 == 0 else NAVY_MID)
        bar = _rect(slide, L, top, Pt(3), Inches(0.72), ELECTRIC_BLUE)
        _glow_blue(bar, 3)
        for ci, (val, color) in enumerate([(module, ELECTRIC_BLUE), (focus, WHITE), (outcome, LIGHT_GRAY)]):
            _tb(slide, val, col_x[ci] + Inches(0.12), top + Inches(0.08),
                col_w[ci] - Inches(0.15), Inches(0.65),
                sz=10, bold=(ci == 0), color=color)


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # -----------------------------------------------------------------------
    # TITLE SLIDE
    # -----------------------------------------------------------------------
    add_title_slide(prs,
        "AI and Agentics\nBasics",
        "PREREQUISITE COURSE  |  2-3 DAYS  |  8 MODULES",
        "[IMAGE: Soldier at a workstation with multiple screens displaying code and AI interface]")

    # -----------------------------------------------------------------------
    # COURSE OVERVIEW
    # -----------------------------------------------------------------------
    add_overview_slide(prs,
        "Eight modules take students from zero to agentic-ready. Every module is hands-on first; concepts follow practice.",
        "2-3 days | 8 modules | Hands-on, project-based",
        [
            "Understand what LLMs are, how they work, and how to interact with them",
            "Navigate and manipulate file systems confidently from the command line",
            "Understand and use Git for basic version control workflows",
            "Write clear, structured Markdown for documentation and prompts",
            "Understand fundamental programming concepts (variables, functions, logic flow)",
            "Set up and authenticate development tools",
            "Troubleshoot basic technical problems independently",
        ])

    # -----------------------------------------------------------------------
    # MODULE 1: LLMs AND PROMPTS
    # -----------------------------------------------------------------------
    add_section_header(prs, "01", "Hello, Claude\nLLMs and Prompts",
        "Understand what an LLM is, how to prompt it clearly, and when to verify the output.",
        "[IMAGE: Soldier at a computer terminal composing a message — deliberate, structured communication]")

    add_concept_slide(prs, "What Is a Large Language Model?", [
        ("An LLM is trained on massive text data to predict the next token in a sequence.", False),
        ("It does not look things up — it generates based on statistical patterns.", False),
        ("Key concepts every operator must understand:", False),
        ("Tokens: the chunks the model processes — not words, but word-pieces", True),
        ("Context window: the finite working memory visible to the model at one time", True),
        ("Temperature: a parameter that controls randomness in generation", True),
        ("The model generates plausible text — not verified fact.", False),
    ], note="The model has no truth-checking step. Confident output and correct output are entirely unrelated. Verification is always the operator's responsibility.")

    add_concept_slide(prs, "Prompts and System Prompts", [
        ("A prompt is your input — the instruction or question you send to the model.", False),
        ("A system prompt is a platform-level instruction that shapes all behavior before the conversation begins.", False),
        ("Typing 'You are a...' in the chat is persona injection — not a true system prompt.", False),
        ("Persona injection is visible in the chat history. A system prompt is not.", False),
        ("This distinction matters when you configure harness files in Claude Code.", False),
    ], note="True system prompts are loaded by the platform before any user turn. In Claude Code, your CLAUDE.md and me.md files function as system context — not mid-conversation instructions.")

    add_concept_slide(prs, "The RGCOA Prompt Structure", [
        ("RGCOA: Role, Goal, Context, Output, Asks — a structured prompting framework.", False),
        ("Role — the expert perspective the model should apply", True),
        ("Goal — the specific task, single and concrete", True),
        ("Context — background required to do the task well", True),
        ("Output — format, length, structure, and tone", True),
        ("Asks — what to do when information is missing; permission to surface uncertainty", True),
        ("Specificity is not optional. Vague orders get vague results.", False),
    ], note="The Asks element is where you give the model permission to flag gaps rather than paper over them with a plausible-sounding guess. Always include it.")

    add_check_on_learning(prs,
        "You send two prompts:\n\n"
        "Prompt A: 'Write about networking.'\n\n"
        "Prompt B: 'Write a 200-word guide on how to build professional relationships "
        "for a newly promoted junior officer — BLUF format, plain language, no jargon.'\n\n"
        "Predict the difference in output quality. What specific elements of Prompt B produce a better result?\n\n"
        "Which prompt required fewer follow-up iterations?")

    add_hands_on(prs, "Module 1 — LLMs and Prompts", [
        "Ice-Breaker: Ask 'What is the weather?' then ask the same with your location and season. Compare specificity.",
        "System Prompt Experiment: Interview Claude as a Marine Corps drill sergeant, then as a kindergarten teacher. Same question, different system prompt. Note 3 differences.",
        "Clarity Matters: Rewrite three vague prompts to be specific and actionable using the RGCOA structure.",
        "Context Window Boundary: Paste a 10,000-word document and ask questions. At what point does quality degrade?",
        "Iterative Prompting: Start vague, iterate 3 times. Track improvement each round. Annotate what changed.",
    ], "[IMAGE: Two mission briefs side by side — one vague wall of text, one structured OPORD]")

    add_section_summary(prs, "Module 1 — Hello, Claude", [
        "LLMs generate plausible text — they do not look up facts. Verification is always your responsibility.",
        "Prompts are your input. System prompts are platform-level instructions that shape all behavior.",
        "The RGCOA structure (Role, Goal, Context, Output, Asks) produces better results with fewer iterations.",
        "Specificity is not optional — vague orders produce vague results.",
    ])

    # -----------------------------------------------------------------------
    # DATA HYGIENE AND OPSEC (standalone section between M1 and M2)
    # -----------------------------------------------------------------------
    add_section_header(prs, "1.5", "Data Hygiene\nand OPSEC",
        "The model is not the classification authority. You are. Know what you are doing before you type.",
        "[IMAGE: Soldier reviewing a document with a red 'CLASSIFIED' stamp before deciding what to share]")

    add_concept_slide(prs, "What Never Goes Into Any AI Tool", [
        ("Classified material at any level", False),
        ("CUI — Controlled Unclassified Information (includes all FOUO-marked documents)", False),
        ("PII: real names combined with any other identifier (SSN, DOB, address, phone)", False),
        ("Personnel records: evaluations, medical, legal, disciplinary", False),
        ("Operational planning detail: timelines, routes, objectives, force compositions", False),
        ("Credentials and API keys: passwords, tokens, certificates — never under any circumstances", False),
    ], note="DANGER: If the document carries ANY marking, or describes real people, real units, real operations, or real capabilities — it does not get pasted into a consumer AI tool. The marking is not required for the rule to apply. The content triggers it.")

    add_concept_slide(prs, "Where Your Prompts Go", [
        ("Consumer AI tools (Claude.ai, ChatGPT.com) transmit input to a provider-operated backend.", False),
        ("Depending on service tier and opt-out settings, input may be reviewed by humans or used in training.", False),
        ("'This tool is impressive' is not the same as 'this tool is authorized for work use.'", False),
        ("Authorization is a property of the contract and ATO — not the capability of the tool.", False),
        ("Without an approved system designation and ATO, a consumer tool is not authorized for CUI.", False),
        ("Claude Code writes conversation history to disk in .claude/ — that log requires the same handling as its content.", False),
    ], note="Scrub discipline applies to what you type, not just what you paste. A session log containing CUI is CUI.")

    add_concept_slide(prs, "The Scrub Checklist", [
        ("Remove ALL of the following before pasting any document into an AI tool:", False),
        ("Real names of any person", True),
        ("Real unit designations, base names, grid coordinates, facility names", True),
        ("Dates tied to real operations or events", True),
        ("Any document markings (CUI, FOUO, LES, classification)", True),
        ("Phone numbers, personal email addresses, account identifiers", True),
        ("Credentials, API keys, passwords — never under any circumstances", True),
        ("System names, network names, IP addresses belonging to real infrastructure", True),
    ], note="Bracketed Placeholder Technique: Replace removed specifics with [NCO], [LOCATION], [UNIT], [DATE], [SYSTEM]. The model reasons about structure and process — not identity. Your scrubbed prompt is fully functional.")

    add_hands_on(prs, "Data Hygiene — Scrub Drill", [
        "Receive the instructor-provided one-page fabricated document (contains name, rank, unit, base, grid, date, marking).",
        "Produce a paste-safe version using bracketed placeholders. Preserve structure and substance.",
        "Remove every item on the scrub checklist. Use [NCO], [UNIT], [LOCATION], [DATE], [SYSTEM].",
        "Debrief: explain what you removed and why. Identify the category each item falls into.",
        "Written outcome: scrubbed document + list of what was removed and its category.",
    ], "[IMAGE: Analyst with a document and a red pen, carefully reviewing and redacting specific fields]")

    add_section_summary(prs, "Data Hygiene and OPSEC", [
        "The model is not the classification authority — you are. Know what you are doing before you type.",
        "Consumer AI tools are not authorized for CUI without an approved system designation and ATO.",
        "Scrub checklist: real names, units, locations, dates, markings, credentials — all must go.",
        "Use bracketed placeholders. The model needs structure, not identity.",
    ])

    # -----------------------------------------------------------------------
    # LLM FAILURE MODES (part of M1 content, surfaced as its own section)
    # -----------------------------------------------------------------------
    add_section_header(prs, "1.6", "LLM Failure\nModes",
        "The model generates plausible text — not verified fact. Know these failure modes before you rely on any output.",
        "[IMAGE: Soldier double-checking a report against source documents before signing off]")

    add_concept_slide(prs, "Failure Mode 1: Hallucination", [
        ("The model generates tokens that are statistically likely to follow prior context.", False),
        ("It has no truth-checking mechanism.", False),
        ("It cannot distinguish accurate training data from a plausible completion.", False),
        ("A detailed, well-cited, grammatically perfect answer can be entirely fabricated.", False),
        ("Confident output and correct output are entirely unrelated.", False),
        ("This is expected system behavior — not a rare bug or edge case.", False),
    ], note="The sharp analyst who never says 'I don't know.' Confident, fluent, and occasionally making it up. You check the work. Fluency is not evidence of accuracy.")

    add_concept_slide(prs, "Failure Modes 2 and 3: Nondeterminism and Knowledge Cutoff", [
        ("Nondeterminism — the model uses temperature, introducing randomness by design.", False),
        ("Same prompt, different run, different result. One output is not 'the answer.'", True),
        ("For high-stakes output: run the prompt multiple times and compare. Divergent answers signal verification.", True),
        ("Knowledge Cutoff — three distinct behaviors:", False),
        ("(a) Model flags uncertainty — honest behavior. Do not mistake it for inability.", True),
        ("(b) Model answers confidently using stale data with no uncertainty flag — the dangerous behavior.", True),
        ("(c) Model has a retrieval tool and pulls current data — verify the tool was actually invoked.", True),
    ], note="All three failure modes look identical from the outside. Verify anything time-sensitive regardless of how confidently the model presents it. The only safe assumption: anything that could have changed since training may have changed.")

    add_hands_on(prs, "LLM Failure Modes — Hallucination in the Wild", [
        "Ask the model to produce 5 peer-reviewed academic papers on a narrow, specific topic.",
        "Request full citations: author, title, journal, year, volume, and page numbers.",
        "Select one citation and verify it using Google Scholar, PubMed, or similar.",
        "Debrief: what made the citation look legitimate? What would you have done without verification?",
        "Written outcome: the citation produced, the verification attempt, and a one-paragraph reflection on hallucination risk in your own work.",
    ], "[IMAGE: Analyst searching an academic database — comparing a citation against real journal records]")

    add_section_summary(prs, "LLM Failure Modes", [
        "Hallucination: the model completes plausibly, not accurately. No truth-checking step.",
        "Nondeterminism: same prompt, different run, different result — by design.",
        "Knowledge cutoff: confident answers about events past training may be wrong.",
        "Verification is your job — always, regardless of output confidence.",
    ])

    # -----------------------------------------------------------------------
    # MODULE 2: COMMAND LINE BASICS
    # -----------------------------------------------------------------------
    add_section_header(prs, "02", "The Terminal\nIs Friendly",
        "The command line is where agentic tools live. Navigate it confidently before you trust an agent to do it for you.",
        "[IMAGE: Soldier confidently operating a field communications terminal with multiple inputs]")

    add_concept_slide(prs, "Windows Students: Use WSL", [
        ("On Windows, use WSL (Windows Subsystem for Linux) — not native PowerShell.", False),
        ("All terminal exercises in this course assume a Linux/Unix environment.", False),
        ("Install WSL: open PowerShell and run wsl --install (no admin required on Windows 11).", False),
        ("First launch prompts for a Linux username and password — separate from Windows login.", False),
        ("Your Windows files inside WSL: /mnt/c/Users/YourName/", False),
        ("Path rule: C:\\Users\\Jake\\Documents becomes /mnt/c/Users/Jake/Documents", False),
    ], note="VS Code + WSL Integration: Install the 'WSL' extension. From a WSL terminal, type 'code .' to open the current folder in VS Code with full WSL integration.")

    add_concept_slide(prs, "Core Terminal Concepts", [
        ("The file system: folders, files, hierarchy — a tree rooted at /", False),
        ("Current working directory: pwd shows where you are right now", False),
        ("Listing files: ls and its flags (-l for detail, -a for hidden files)", False),
        ("Paths: absolute (/home/user/docs) vs relative (../docs)", False),
        ("Creating and deleting: mkdir, touch, rm, rmdir", False),
        ("Moving and copying: mv, cp", False),
        ("Help systems: man ls, ls --help — never guess when docs exist", False),
        ("Ctrl+C — the most important shortcut: stops any running command immediately", False),
    ], note="Ctrl+C sends an interrupt signal and stops the current process. Name it explicitly and drill it early — students need this reflex before anything else.")

    add_concept_slide(prs, "Piping, Redirection, and PATH", [
        ("Piping (|): send the output of one command as input to another", False),
        ("ls | grep .txt — list files, filter for .txt", True),
        ("Redirection: > writes output to a file, >> appends, < reads from a file", False),
        ("ls > file-list.txt — saves directory listing to a file", True),
        ("PATH: the list of directories the shell searches for executable commands", False),
        ("'Command not found' almost always means: not on PATH", True),
        ("Environment variables: KEY=VALUE pairs the shell and programs can read", False),
        ("Hidden files and dotfiles: files starting with . (like .env, .gitignore)", False),
    ])

    add_hands_on(prs, "Module 2 — Terminal Basics", [
        "Orientation Walk: Start at home directory, navigate to 5 locations, use pwd after each move, create a file.",
        "File Manipulation Sprint: Build projects/acc-prep/module-2/ — create 3 files, copy, rename, delete, list final structure.",
        "Path Puzzle: Write absolute and relative paths to a given file. Navigate using those paths to verify.",
        "The Help System: Use man ls and ls --help to answer questions. Build a reference card for 6 commands.",
        "Piping and Redirection: Produce 3 piping examples and 2 redirection examples in a terminal transcript.",
        "Real-World Scenario: Find all .txt files in a 50-file folder, sort by date, move to an archive folder.",
    ], "[IMAGE: Soldier navigating a command-line interface with multiple directories open on screen]")

    add_readiness_check(prs, [
        "You can navigate to any directory using absolute and relative paths",
        "You can create, copy, move, and delete files and folders from the terminal",
        "You know how to stop a running command (Ctrl+C)",
        "You can pipe and redirect command output",
        "You can explain what PATH is and why it matters",
    ])

    add_section_summary(prs, "Module 2 — The Terminal", [
        "The terminal is the operating environment for all agentic tools — own it before you delegate to an agent.",
        "Ctrl+C is your emergency brake. Drill it early.",
        "Absolute paths start at /. Relative paths start from where you are. Both matter.",
        "Piping and redirection chain commands — a core skill for automation.",
    ])

    # -----------------------------------------------------------------------
    # MODULE 3: GIT BASICS
    # -----------------------------------------------------------------------
    add_section_header(prs, "03", "Git Basics\nVersion Control",
        "Git is the audit trail for agentic work. Every agent change gets a commit. Every commit is reviewable.",
        "[IMAGE: NCO maintaining a detailed operations log — every action recorded and traceable]")

    add_concept_slide(prs, "What Is Version Control?", [
        ("Git tracks changes to files over time — a full history of every modification.", False),
        ("The three Git states every operator must know:", False),
        ("Working directory — your current files, with unsaved changes", True),
        ("Staging area — changes you have selected for the next commit (git add)", True),
        ("Committed — a permanent, immutable snapshot in the repository history", True),
        ("Local vs remote: your machine vs GitHub (the shared record)", False),
        ("Commits are immutable records of change — they are your audit trail.", False),
    ], note="The PR (pull request) is the quality gate for agentic work: when an agent produces changes on a branch, the PR is how you review what it did before it touches main.")

    add_concept_slide(prs, "Branches and Commit Discipline", [
        ("A branch is an isolated line of work — safe to modify without affecting main.", False),
        ("Workflow: create branch → make changes → commit → merge back to main.", False),
        ("Merge conflicts occur when two branches modify the same line — resolve manually.", False),
        ("Commit messages are documentation. Bad examples: 'fix stuff', 'update', 'work'.", False),
        ("Good messages: 'Add authentication logic to login form', 'Fix off-by-one in pagination'", False),
        ("Convention: present tense, action-focused, informative.", False),
    ], note="WARNING: .gitignore must be created BEFORE the first commit. If a credential file is committed first, adding it to .gitignore later does NOT remove it from history. The file remains recoverable.")

    add_concept_slide(prs, "The .gitignore and Pull Requests", [
        ("The .gitignore file tells Git which files to never track.", False),
        ("Minimum entries for any project:", False),
        (".env (API keys and secrets — never commit)", True),
        ("node_modules/ or .venv/ (dependencies — rebuilt from lock files)", True),
        (".DS_Store, __pycache__/ (OS and language artifacts)", True),
        ("A pull request (PR) is a proposal to merge one branch into another, with a review step.", False),
        ("The PR shows a diff — exactly what changed, line by line. Review it before merging.", False),
        ("This is the human review gate for all agentic output.", False),
    ])

    add_hands_on(prs, "Module 3 — Git Basics", [
        "Your First Repo: init → create file → stage → commit → view log → make change → commit again.",
        "Branching Sprint: Create feature/add-content branch, commit changes, switch to main, verify unchanged, merge, delete branch.",
        "Cloning and Pushing: Clone a public repo, make a change, push to your own GitHub repo successfully.",
        "Merge Conflicts: Create two branches modifying the same line, merge one (succeeds), merge other (conflict), resolve manually.",
        "Commit Messages: Write 5 meaningful commit messages for provided scenarios. Present tense, action-focused.",
        "Pull Request: Create branch, push to GitHub, open PR with gh pr create, review the diff, merge.",
    ], "[IMAGE: Soldier reviewing a git log on screen — tracking exactly what changed and when]")

    add_readiness_check(prs, [
        "You can initialize a repo, stage files, and commit with a meaningful message",
        "You can create a branch, make changes, and merge it back to main",
        "You can clone a repo and push changes to GitHub",
        "Your .gitignore exists before any commits and excludes .env and credential files",
        "You can open a PR and read the diff before merging",
    ])

    add_section_summary(prs, "Module 3 — Git", [
        "Git is the audit trail for agentic work — every agent change gets a commit, every commit is reviewable.",
        ".gitignore must exist before the first commit. Credentials committed once stay in history forever.",
        "Branches isolate work. PRs are the review gate before anything touches main.",
        "Commit messages are documentation — write them for the operator who reviews the diff.",
    ])

    # -----------------------------------------------------------------------
    # MODULE 4: MARKDOWN
    # -----------------------------------------------------------------------
    add_section_header(prs, "04", "Markdown\nDocumenting Like a Pro",
        "Markdown is the formatting language of prompts, documentation, and context files. Broken Markdown breaks your instructions.",
        "[IMAGE: Soldier writing a clear, structured field report with headings, lists, and tables]")

    add_concept_slide(prs, "Markdown Fundamentals", [
        ("Markdown is plain text with lightweight formatting syntax — readable as-is, rendered in tools.", False),
        ("This course teaches GitHub Flavored Markdown (GFM) — renders in VS Code and GitHub.", False),
        ("Core syntax operators must know:", False),
        ("Headers: # H1, ## H2, ### H3 — require a space after the # character", True),
        ("Lists: - or * for unordered, 1. for ordered — require a blank line before the list", True),
        ("Bold: **text**, Italic: *text*, Inline code: `command`", True),
        ("Code blocks: triple backtick with language hint (```python, ```bash)", True),
        ("Tables: require both a header row and a separator row (|---|---|)", True),
    ], note="Critical Spacing Rules — Most Beginner Errors Live Here: space after #, blank line before lists, blank line before code blocks, separator row in tables. Miss these and your instructions render as a wall of text.")

    add_concept_slide(prs, "Markdown for Context Files", [
        ("Markdown is how you write CLAUDE.md and me.md — the system context for Claude Code.", False),
        ("The model reads ## headers as distinct instruction blocks.", False),
        ("A bullet list under a header creates discrete, unambiguous rules.", False),
        ("If the Markdown formatting breaks, the model reads a wall of text instead of structure.", False),
        ("Smart-quote warning: do not paste Markdown from Word or Google Docs.", False),
        ("Smart quotes break code blocks and inline code — write Markdown in VS Code only.", False),
        ("Use Ctrl+Shift+V in VS Code to preview before submitting.", False),
    ], note="YAML frontmatter: optional metadata block at the top of a Markdown file (title, author, date). Written between triple-dashes (---). Treated as metadata, not content.")

    add_hands_on(prs, "Module 4 — Markdown", [
        "Markdown Anatomy: Read a well-formatted file, identify all syntax, recreate from scratch with 10+ elements.",
        "Personal README: Write a markdown file with heading, bullets, link, bold/italic, code block. Commit to Git.",
        "Code Blocks and Syntax: Create examples in Python, Bash, and JavaScript with syntax highlighting.",
        "Tables and Organization: Build two well-formatted tables — one comparing tools, one as a project plan.",
        "Documentation Project: Document a simple process (300-500 words) with headers, lists, code, and a table of contents.",
    ], "[IMAGE: Split screen — raw Markdown text on left, rendered preview on right in VS Code]")

    add_readiness_check(prs, [
        "You can write a properly formatted Markdown document with 5+ syntax elements",
        "You know the critical spacing rules (# space, blank lines before lists and code blocks)",
        "You write Markdown in VS Code, not Word or Google Docs",
        "You can create tables, code blocks, and task lists",
    ])

    add_section_summary(prs, "Module 4 — Markdown", [
        "Markdown is the formatting language of context files, prompts, and all documentation in this course.",
        "Broken formatting means broken instructions — the model reads a wall of text instead of structured rules.",
        "Write Markdown in VS Code. Preview before committing. Never paste from Word.",
        "Master the critical spacing rules — they account for the majority of beginner formatting errors.",
    ])

    # -----------------------------------------------------------------------
    # MODULE 5: PROGRAMMING CONCEPTS
    # -----------------------------------------------------------------------
    add_section_header(prs, "05", "Programming\nConcepts",
        "You do not need to be a developer. You need to read code, understand what it does, and direct an agent that writes it.",
        "[IMAGE: Soldier reading a technical manual and annotating key sections — comprehension before execution]")

    add_concept_slide(prs, "Variables, Types, and Operators", [
        ("A variable is a named container for data. Name it clearly — the name is documentation.", False),
        ("Data types: string ('text'), number (42), boolean (true/false), array ([item1, item2]), object ({key: value})", False),
        ("Type errors occur when an operation expects one type but receives another.", False),
        ("Operators: arithmetic (+, -, *, /), comparison (==, !=, >, <), logical (&&, ||, !)", False),
        ("Write pseudocode first — the concept, not the syntax.", False),
        ("Then implement. Then test.", False),
    ], note="These exercises use pseudocode and one real language (JavaScript). The goal is concepts, not syntax mastery. Understanding what code does is a prerequisite to directing an agent that writes it.")

    add_concept_slide(prs, "Control Flow and Functions", [
        ("Conditionals: if/else — execute different code based on a condition being true or false.", False),
        ("Loops: for (known count), while (until condition met), forEach (each item in a list)", False),
        ("Functions: a named, reusable block of code with inputs (parameters) and outputs (return values)", False),
        ("Scope: variables inside a function are not visible outside it.", False),
        ("Debugging: read the error message first — it tells you the type, file, and line.", False),
        ("Understand: typos, logic errors, boundary conditions — three categories of bugs.", False),
    ])

    add_hands_on(prs, "Module 5 — Programming Concepts", [
        "Variables and Types: Pseudocode + JavaScript — store name, age, boolean; calculate next year's age; combine in a sentence.",
        "Control Flow: A coffee machine with price logic. Pseudocode first, then JavaScript, then test 3 types × 2 sizes.",
        "Loops: Print 1-10, repeat a phrase, sum 1-100, iterate an array. 4 examples with comments.",
        "Functions: greeting, sum, password check, vowel counter — 4 JavaScript functions with clear names.",
        "Debugging Basics: Fix 5 buggy snippets (syntax errors, logic errors, off-by-one). Annotate each fix.",
        "Pseudocode Planning: Grade calculator — write pseudocode first, then implement, then test edge cases.",
    ], "[IMAGE: Soldier annotating a technical flowchart before writing a single line of code]")

    add_readiness_check(prs, [
        "You can write pseudocode before writing real code",
        "You understand variables, conditionals, loops, and functions",
        "You can read a buggy snippet and identify the error type",
        "You can write a working function with a clear name and documented purpose",
    ])

    add_section_summary(prs, "Module 5 — Programming Concepts", [
        "Write pseudocode before code — plan the logic before you worry about syntax.",
        "Variables, conditionals, loops, and functions are the building blocks of every program.",
        "Read the error message first — it tells you the type, location, and often the fix.",
        "Comprehension enables direction. You must understand what an agent writes before you approve it.",
    ])

    # -----------------------------------------------------------------------
    # MODULE 6: DEVELOPER TOOLS SETUP
    # -----------------------------------------------------------------------
    add_section_header(prs, "06", "Developer Tools\nInstallation and Auth",
        "An agent is only as effective as the tools it can reach. Get your environment fully operational before you build.",
        "[IMAGE: Soldier performing pre-mission equipment checks — methodical, complete, nothing skipped]")

    add_concept_slide(prs, "Tools You Will Install", [
        ("Claude Code — the agentic development harness. The primary tool for this course.", False),
        ("GitHub CLI (gh) — command-line interface to GitHub. Enables PR creation, auth, and repo management.", False),
        ("VS Code — the code editor. Integrates with WSL and Git natively.", False),
        ("Node.js or Python — the runtime for your capstone project.", False),
        ("Git — already covered, but verify installation before proceeding.", False),
        ("Environment variables (.env files) — how credentials are stored without appearing in code.", False),
    ], note="Every tool must be installed AND verified. 'I think I installed it' is not verified. Run the version check command for each tool. No exceptions.")

    add_concept_slide(prs, "Authentication and API Keys", [
        ("Authentication proves your identity to a system — GitHub, Anthropic API, etc.", False),
        ("API keys are secret strings that let code authenticate as you.", False),
        ("They must never appear in code, commits, or chat sessions.", False),
        ("HTTPS vs SSH: recommend HTTPS for beginners, SSH for advanced users.", False),
        ("Token-based authentication: GitHub personal access tokens, Anthropic API keys.", False),
        ("PATH and executable discovery: 'Command not found' means the tool is not on PATH.", False),
    ], note="DANGER: .env files must never be committed. Create .gitignore before the first commit and include .env as the very first entry. A key committed once stays in history.")

    add_hands_on(prs, "Module 6 — Developer Tools", [
        "Install Claude Code: follow official instructions for your OS. Verify with claude --version. Run claude to confirm interaction.",
        "GitHub Authentication: install gh, run gh auth login (choose HTTPS), verify with gh auth status, clone a test repo.",
        "VS Code Setup: install VS Code, open a folder, install Markdown Preview + Git Graph + Prettier extensions.",
        "Environment Variables: create a .env with sample variables, create a .gitignore that excludes .env.",
        "Tool Verification Sprint: run claude --version, node --version, git --version, gh --version, code --version. Document any failures.",
    ], "[IMAGE: Soldier running a pre-mission equipment checklist with each item checked off systematically]")

    add_readiness_check(prs, [
        "Claude Code is installed and claude --version passes",
        "GitHub CLI is authenticated (gh auth status passes)",
        "VS Code is installed with your preferred extensions",
        "You know what .env files are and why they must not be committed",
    ])

    add_section_summary(prs, "Module 6 — Developer Tools", [
        "Install and verify every tool — version checks are not optional.",
        "API keys go in .env files. .env goes in .gitignore. .gitignore is created before the first commit.",
        "Authentication failure blocks the entire workflow — resolve it before moving to Module 7.",
        "Claude Code is the primary tool for this course. Confirm it runs before continuing.",
    ])

    # -----------------------------------------------------------------------
    # MODULE 7: CONTEXT AND CONFIGURATION
    # -----------------------------------------------------------------------
    add_section_header(prs, "07", "Context and\nConfiguration",
        "CLAUDE.md and me.md are how you brief the agent before the mission begins. Write them like an OPORD, not a wish list.",
        "[IMAGE: Commander at a map board briefing the team — clear intent, explicit constraints, defined success]")

    add_concept_slide(prs, "Why Context Files Matter", [
        ("Claude Code reads CLAUDE.md and me.md before any conversation begins.", False),
        ("These files function as the system prompt for your project.", False),
        ("They set: what the project is, what is off-limits, what success looks like.", False),
        ("Vague context = vague results. Specific context = predictable, defensible results.", False),
        ("The .claude/ folder: where Claude Code stores context files and session logs.", False),
        ("Context files are written in Markdown — formatting matters. See Module 4.", False),
    ], note="Every element of your CLAUDE.md is Markdown syntax you learned in Module 4. The ## headers create distinct instruction blocks. The bullet lists under Rules create discrete, unambiguous constraints.")

    add_concept_slide(prs, "CLAUDE.md: Project Context", [
        ("CLAUDE.md contains the project brief the agent reads before every session.", False),
        ("Essential sections:", False),
        ("What this project does — the mission in one paragraph", True),
        ("What it is NOT — explicit scope boundaries", True),
        ("Rules — discrete, enforceable constraints (what the model can and cannot do)", True),
        ("Files — what Claude can and cannot modify", True),
        ("Success criteria — what done looks like", True),
        ("Write constraints as discrete bullet points, not paragraphs.", False),
    ], note="Example rule: 'Never include real names in output — use [NAME] as a placeholder.' This is enforceable. 'Be careful with names' is not.")

    add_concept_slide(prs, "me.md: Personal Operating Profile", [
        ("me.md tells Claude how to work with you — your style, preferences, and decision-making approach.", False),
        ("Include: your name and role, decision-making style, communication preferences.", False),
        ("Include: what you value in work (speed, clarity, thoroughness), how you want Claude to interact.", False),
        ("A conflict resolution matrix clarifies what wins when requirements compete:", False),
        ("Speed vs Security — which wins for this project?", True),
        ("Simplicity vs Features — where do you draw the line?", True),
        ("Completeness vs Time-to-delivery — what is your threshold?", True),
        ("Write it once. Update it when your preferences change.", False),
    ])

    add_example_slide(prs, "Example: CLAUDE.md in Practice",
        "Scenario: Daily Brief Generator project",
        [
            "# Project: Daily Brief Generator",
            "  ## What this does",
            "  Pulls the unit's daily log and produces a formatted brief.",
            "  ## Rules",
            "  - Never include real names in the output -- use [NAME] as a placeholder",
            "  - Use the provided template structure, do not invent new sections",
            "  - Flag any entries with incomplete location or time data",
            "  ## Files",
            "  - brief-template.md -- the output template (read-only)",
            "  - input-log.txt -- the daily log to process",
            "",
            "Every section is enforceable. Every constraint is discrete.",
            "The model cannot comply with a vague instruction -- specificity is the gate.",
        ],
        "[IMAGE: Commander's written intent document — clear mission, explicit constraints, defined end state]")

    add_hands_on(prs, "Module 7 — Context and Configuration", [
        "CLAUDE.md Anatomy: Read the example CLAUDE.md. Annotate each section and what behavior it produces.",
        "Write Your me.md: name, role, decision-making style, communication preference, how you want Claude to interact (5-10 lines).",
        "Project CLAUDE.md: Scenario — task manager for Marines. Write CLAUDE.md with what, what NOT, constraints, files, success criteria.",
        "Conflict Resolution Matrix: Rate Speed vs Security, Simplicity vs Features, Completeness vs Time for your scenario project.",
        "Set Up Project Context: create folder, .claude/ subfolder, add CLAUDE.md and me.md, run Claude Code, confirm it reads your context.",
    ], "[IMAGE: Soldier setting up a command post — context boards, reference materials, communication equipment all in place]")

    add_readiness_check(prs, [
        "You have a complete CLAUDE.md that specifies project scope, rules, and constraints",
        "You have a personal me.md with communication preferences and working style",
        "Claude Code reads your context files and references them in conversation",
        "You understand how Markdown formatting affects how the model reads instructions",
    ])

    add_section_summary(prs, "Module 7 — Context and Configuration", [
        "CLAUDE.md is your project brief — what the project is, what is off-limits, what success looks like.",
        "me.md is your operating profile — how you work, what you value, how Claude should interact with you.",
        "Specific, discrete constraints are enforceable. Vague instructions are not.",
        "Formatting matters — broken Markdown means broken instructions.",
    ])

    # -----------------------------------------------------------------------
    # MODULE 8: CAPSTONE PROJECT
    # -----------------------------------------------------------------------
    add_section_header(prs, "08", "Putting It\nAll Together",
        "The capstone integrates all eight modules. Plan, build, commit, push. Own the result.",
        "[IMAGE: Soldiers conducting a full mission rehearsal — all elements integrated, leader at the front]")

    add_concept_slide(prs, "Three Capstone Paths", [
        ("Choose one path. Build something real. Document it. Push it to GitHub.", False),
        ("Path A: CLI Tool — password strength checker, markdown-to-HTML converter, task logger", False),
        ("Takes input from user or file, processes with logic, outputs result, handles errors", True),
        ("Path B: Web App — note-taking app, unit converter, quote generator", False),
        ("HTML + CSS + JavaScript, takes user input, runs locally, documented", True),
        ("Path C: API Integration — CLI or web app calling a free public API", False),
        ("Authenticates with a real API, handles errors gracefully, displays results clearly", True),
        ("All paths: JavaScript (Node.js) or Python. Minimum 3 commits. README required.", False),
    ], note="Technology recommendation: JavaScript (Node.js) or Python. Both are well-supported in Claude Code. Choose the language you are most comfortable reading.")

    add_concept_slide(prs, "The Capstone Workflow", [
        ("Step 1: Plan (20 min) — write a one-paragraph description, list 3-5 features, first commit.", False),
        ("Step 2: Set Up (20 min) — create .claude/ folder, init git, create file structure, commit.", False),
        ("Step 3: Build Phase 1 (60 min) — simplest version first, use Claude Code, commit per feature.", False),
        ("Step 4: Test and Iterate (40 min) — use your own tool, find bugs, ask Claude to help improve.", False),
        ("Step 5: Document (20 min) — write README, add code comments, final commit.", False),
        ("Step 6: Push to GitHub (10 min) — verify all commits local, push, verify on GitHub.com.", False),
    ])

    add_concept_slide(prs, "Capstone Assessment Rubric", [
        ("Does it work? (50%) — runs without errors, does what was promised, handles basic errors", False),
        ("Is it documented? (20%) — README is clear, includes setup/usage, code has comments", False),
        ("Are commits meaningful? (15%) — 5+ commits with clear messages, pushed to GitHub", False),
        ("Did you iterate? (10%) — evidence of feedback loops, bug fixes, improvements from testing", False),
        ("Is it polished? (5%) — readable, organized, no debug statements left behind", False),
        ("Bonus: +5 points — deployment, automated tests, second feature, or other demonstrated ambition", False),
    ], note="Passing: 70+ points. Grades: 70-79 Pass, 80-89 Good, 90+ Excellent. 7/8 modules passed = course complete. 8/8 = honors.")

    add_hands_on(prs, "Module 8 — Capstone Project", [
        "Choose your path (CLI Tool, Web App, or API Integration). Write one-paragraph project description.",
        "Set up: .claude/ folder, CLAUDE.md, me.md, git init, first commit: 'Initial project plan'.",
        "Build core functionality using Claude Code. Commit per feature with meaningful messages.",
        "Test your own tool. Find bugs. Ask Claude to help with error handling and edge cases. Commit fixes.",
        "Write a clear README (what it does, how to run it, any requirements). Final commit: 'Update documentation'.",
        "Push to GitHub. Verify all commits appear. Confirm publicly visible. Submit repo link.",
    ], "[IMAGE: Soldier presenting a completed project to peers — confident delivery, clear documentation in hand]")

    add_readiness_check(prs, [
        "Your capstone project runs without errors",
        "It is documented with a clear README (what it does, how to run it)",
        "You have 5+ commits with meaningful messages",
        "No credentials appear in your git history",
        "Your project is pushed to GitHub and publicly visible",
    ])

    # -----------------------------------------------------------------------
    # SUMMARY TABLE
    # -----------------------------------------------------------------------
    add_summary_table(prs, [
        ("M1: LLMs & Prompts",     "RGCOA prompting, failure modes, context windows", "Clear, verified prompts using structured format"),
        ("M2: Terminal",            "Navigation, file ops, piping, PATH",              "Confident command-line operation from any directory"),
        ("M3: Git",                 "Commits, branches, PRs, .gitignore",              "Version-controlled workflow with traceable changes"),
        ("M4: Markdown",            "Syntax, spacing rules, context files",            "Well-formatted docs and context files for Claude"),
        ("M5: Programming",         "Variables, logic, functions, debugging",          "Read and direct code without syntax dependence"),
        ("M6: Tools Setup",         "Claude Code, GitHub CLI, VS Code, .env",          "Fully operational development environment"),
        ("M7: Context Config",      "CLAUDE.md, me.md, conflict matrix",               "Properly briefed agent with enforceable constraints"),
        ("M8: Capstone",            "Full workflow: plan → build → commit → push",     "Working project on GitHub, documented and defensible"),
    ])

    # -----------------------------------------------------------------------
    # COURSE READINESS CHECK
    # -----------------------------------------------------------------------
    add_readiness_check(prs, [
        "I can write a clear, specific prompt using the RGCOA structure",
        "I know the difference between a prompt and a system prompt",
        "I can identify at least one LLM failure mode and explain how to verify for it",
        "I know what never goes into a consumer AI tool (classified, CUI, PII, credentials)",
        "I can navigate the terminal confidently using absolute and relative paths",
        "I can initialize a Git repo, commit with meaningful messages, and open a PR",
        "I can write well-formatted Markdown including tables, code blocks, and task lists",
        "I understand variables, conditionals, loops, and functions",
        "All developer tools are installed and version-verified",
        "I have a CLAUDE.md and me.md for my capstone project",
        "My capstone project runs, is documented, and is pushed to GitHub",
    ])

    # -----------------------------------------------------------------------
    # COURSE COMPLETION
    # -----------------------------------------------------------------------
    add_end_slide(prs, "AI and Agentics\nBasics — Complete", [
        "You are ready for advanced agentic work. The prerequisites are behind you.",
        "Keep your CLAUDE.md and me.md updated as your projects and preferences evolve.",
        "Apply OPSEC discipline in every session — no credentials, no CUI, no classified material.",
        "Review your capstone commit history — this is the habit of work that defines an operator.",
    ])

    out = os.path.join(os.path.dirname(__file__), "prereq-course.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
