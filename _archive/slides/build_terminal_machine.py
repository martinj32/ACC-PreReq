"""
ACC Terminal Basics — The Machine Slide Deck Builder
Source: docs/terminal-basics/the-machine.md
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree
import os

# ---------------------------------------------------------------------------
# Color palette
# ---------------------------------------------------------------------------
CYBER_BLACK    = RGBColor(0x0A, 0x0C, 0x14)
DARK_NAVY      = RGBColor(0x0F, 0x14, 0x23)
NAVY_MID       = RGBColor(0x15, 0x1E, 0x35)
ELECTRIC_BLUE  = RGBColor(0x00, 0xB4, 0xFF)
CYBER_GOLD     = RGBColor(0xFF, 0xAA, 0x00)
TERMINAL_GREEN = RGBColor(0x00, 0xE5, 0x7A)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY     = RGBColor(0xA8, 0xB4, 0xC8)
DIM_GRAY       = RGBColor(0x3A, 0x44, 0x58)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

# ---------------------------------------------------------------------------
# Glow
# ---------------------------------------------------------------------------

def _glow(shape, color_rgb, radius_pt=4, alpha_pct=25):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    lst = spPr.find(qn('a:effectLst'))
    if lst is None:
        lst = etree.SubElement(spPr, qn('a:effectLst'))
    for old in lst.findall(qn('a:glow')):
        lst.remove(old)
    g = etree.SubElement(lst, qn('a:glow'))
    g.set('rad', str(int(radius_pt * 12700)))
    c = etree.SubElement(g, qn('a:srgbClr'))
    c.set('val', f'{color_rgb[0]:02X}{color_rgb[1]:02X}{color_rgb[2]:02X}')
    a = etree.SubElement(c, qn('a:alpha'))
    a.set('val', str(int(alpha_pct * 1000)))

def _glow_blue(s, r=4):  _glow(s, ELECTRIC_BLUE, r, 25)
def _glow_gold(s, r=3):  _glow(s, CYBER_GOLD, r, 22)
def _glow_green(s, r=3): _glow(s, TERMINAL_GREEN, r, 22)

# ---------------------------------------------------------------------------
# Primitives
# ---------------------------------------------------------------------------

def _bg(slide, color):
    f = slide.background.fill
    f.solid()
    f.fore_color.rgb = color

def _rect(slide, l, t, w, h, fill, line=None, lw=None):
    s = slide.shapes.add_shape(1, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        if lw: s.line.width = Pt(lw)
    else:
        s.line.fill.background()
    return s

def _tb(slide, text, l, t, w, h, sz=15, bold=False, color=WHITE,
        align=PP_ALIGN.LEFT, italic=False, font="Calibri"):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(sz)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return box

def _paras(slide, items, l, t, w, h):
    """items: list of dicts: text, sz, bold, color, space_before"""
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = item.get('align', PP_ALIGN.LEFT)
        p.space_before = Pt(item.get('sb', 5))
        run = p.add_run()
        run.text = item['text']
        run.font.name = item.get('font', 'Calibri')
        run.font.size = Pt(item.get('sz', 14))
        run.font.bold = item.get('bold', False)
        run.font.color.rgb = item.get('color', WHITE)
    return box

def _line(slide, l, t, w, color, thickness_pt=2, glow_r=4, glow_a=25):
    s = _rect(slide, l, t, w, Pt(thickness_pt), color)
    _glow(s, color, glow_r, glow_a)
    return s

def _card(slide, l, t, w, h, accent=None):
    c = _rect(slide, l, t, w, h, DARK_NAVY)
    if accent:
        bar = _rect(slide, l, t, Pt(4), h, accent)
        _glow(bar, accent, 3, 22)
    return c

def _img_placeholder(slide, l, t, w, h, label, accent=ELECTRIC_BLUE):
    c = _rect(slide, l, t, w, h, DARK_NAVY, accent, 1.5)
    _glow(c, accent, 3, 20)
    _tb(slide, "◼  IMAGE", l, t + (h/2) - Inches(0.35), w, Inches(0.3),
        sz=9, bold=True, color=accent, align=PP_ALIGN.CENTER)
    _tb(slide, label, l + Inches(0.2), t + (h/2) - Inches(0.05),
        w - Inches(0.4), Inches(1.0), sz=9, color=LIGHT_GRAY, align=PP_ALIGN.CENTER)
    return c

# ---------------------------------------------------------------------------
# Common chrome
# ---------------------------------------------------------------------------

L = Inches(0.45)   # left margin
R = Inches(12.43)  # content width (L to right edge minus margin)

def _header(slide, title, sub=None, accent=ELECTRIC_BLUE):
    _tb(slide, title, L, Inches(0.2), R, Inches(0.75),
        sz=28, bold=True, color=WHITE)
    line_top = Inches(0.95) if not sub else Inches(1.25)
    if sub:
        _tb(slide, sub, L, Inches(0.95), R, Inches(0.35), sz=12, color=LIGHT_GRAY)
    _line(slide, L, line_top, R, accent, 2, 4, 25)

def _footer(slide, text="ACC  |  AGENTIC COMMANDERS COURSE", accent=ELECTRIC_BLUE):
    _line(slide, Inches(0), Inches(7.22), SLIDE_W, accent, 1.5, 3, 20)
    _tb(slide, text, L, Inches(7.3), Inches(8), Inches(0.2),
        sz=8, color=DIM_GRAY)

def _callout(slide, text, l, t, w, h, label="NOTE", accent=CYBER_GOLD, sz=12):
    _card(slide, l, t, w, h, accent)
    _tb(slide, label, l + Inches(0.2), t + Inches(0.07), Inches(1.5), Inches(0.28),
        sz=9, bold=True, color=accent)
    _tb(slide, text, l + Inches(0.2), t + Inches(0.34), w - Inches(0.4),
        h - Inches(0.42), sz=sz, color=LIGHT_GRAY)

# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------

def add_title_slide(prs, title, subtitle, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    vert = _rect(slide, Inches(0), Inches(0), Pt(10), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "AGENTIC COMMANDERS COURSE", Inches(0.3), Inches(0.3),
        Inches(6), Inches(0.35), sz=10, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, title, Inches(0.3), Inches(1.0), Inches(6.5), Inches(2.8),
        sz=40, bold=True, color=WHITE)
    gl = _line(slide, Inches(0.3), Inches(4.0), Inches(5.8), CYBER_GOLD, 2, 3, 22)
    _tb(slide, subtitle, Inches(0.3), Inches(4.15), Inches(6), Inches(0.5),
        sz=13, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.1), Inches(0.35), Inches(5.9), Inches(6.65), img_label)
    _footer(slide)


def add_overview_slide(prs, bluf, duration, objectives):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "MODULE OVERVIEW", "Situation Brief")
    _footer(slide)

    _card(slide, L, Inches(1.45), R, Inches(1.05), CYBER_GOLD)
    _tb(slide, "BLUF", L + Inches(0.2), Inches(1.5), Inches(1), Inches(0.28),
        sz=9, bold=True, color=CYBER_GOLD)
    _tb(slide, bluf, L + Inches(0.2), Inches(1.77), R - Inches(0.4), Inches(0.65),
        sz=14, color=WHITE)

    _tb(slide, f"⏱  {duration}", L, Inches(2.6), Inches(5), Inches(0.35),
        sz=12, bold=True, color=CYBER_GOLD)

    _card(slide, L, Inches(3.05), R, Inches(3.95), ELECTRIC_BLUE)
    _tb(slide, "LEARNING OBJECTIVES", L + Inches(0.2), Inches(3.12),
        Inches(6), Inches(0.28), sz=9, bold=True, color=ELECTRIC_BLUE)

    items = [{'text': f"  ›  {o}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 4}
             for o in objectives]
    _paras(slide, items, L + Inches(0.2), Inches(3.45), R - Inches(0.4), Inches(3.4))


def add_section_header(prs, num, title, one_liner, img_label):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)

    # Decorative large dim number
    _tb(slide, num, Inches(-0.15), Inches(0.4), Inches(4), Inches(5),
        sz=190, bold=True, color=NAVY_MID)

    _tb(slide, f"SECTION  {num}", L, Inches(0.95), Inches(4), Inches(0.38),
        sz=11, bold=True, color=CYBER_GOLD)
    _line(slide, L, Inches(1.38), Inches(5.8), CYBER_GOLD, 1.5, 3, 22)
    _tb(slide, title, L, Inches(1.55), Inches(6.5), Inches(3.0),
        sz=44, bold=True, color=WHITE)
    _line(slide, L, Inches(4.7), Inches(5.8), ELECTRIC_BLUE, 2, 4, 25)
    _tb(slide, one_liner, L, Inches(4.85), Inches(6.5), Inches(1.2),
        sz=14, color=LIGHT_GRAY)
    _img_placeholder(slide, Inches(7.3), Inches(0.6), Inches(5.7), Inches(6.4), img_label)
    _footer(slide)


def add_concept_slide(prs, title, bullets, note=None, accent=ELECTRIC_BLUE):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, accent=accent)
    _footer(slide)

    ct = Inches(1.5)
    ch = Inches(4.85) if not note else Inches(3.45)
    _card(slide, L, ct, R, ch, accent)

    items = []
    for text, is_sub in bullets:
        items.append({
            'text': f"      ◦  {text}" if is_sub else f"  ›  {text}",
            'sz': 13 if is_sub else 15,
            'color': LIGHT_GRAY if is_sub else WHITE,
            'sb': 3 if is_sub else 7,
        })
    _paras(slide, items, L + Inches(0.2), ct + Inches(0.18), R - Inches(0.4), ch - Inches(0.3))

    if note:
        _callout(slide, note, L, Inches(5.1), R, Inches(1.75))


def add_example_slide(prs, title, scenario, lines, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, scenario)
    _footer(slide)

    cw = Inches(7.1) if img_label else R
    _card(slide, L, Inches(1.55), cw, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': ln, 'sz': 12,
              'color': WHITE if not ln.startswith('  ') else LIGHT_GRAY,
              'sb': 3 if ln else 8}
             for ln in lines]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), cw - Inches(0.3), Inches(5.1))

    if img_label:
        _img_placeholder(slide, Inches(7.75), Inches(1.55), Inches(5.13), Inches(5.5), img_label)


def add_check_on_learning(prs, question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), ELECTRIC_BLUE)
    _glow_blue(bar, 4)
    _tb(slide, "CHECK ON LEARNING", Inches(0.45), Inches(0.1), Inches(8), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)

    _card(slide, L, Inches(0.78), R, Inches(6.27), ELECTRIC_BLUE)
    _tb(slide, "Before You Continue", L + Inches(0.2), Inches(0.9),
        Inches(10), Inches(0.52), sz=18, bold=True, color=CYBER_GOLD)
    _line(slide, L + Inches(0.2), Inches(1.47), Inches(12.0), ELECTRIC_BLUE, 1.5, 3, 20)
    _tb(slide, question, L + Inches(0.2), Inches(1.65), R - Inches(0.4), Inches(5.2),
        sz=15, color=WHITE)


def add_hands_on(prs, section_title, steps, img_label=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide, accent=TERMINAL_GREEN)

    bar = _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.6), TERMINAL_GREEN)
    _glow_green(bar, 4)
    _tb(slide, "HANDS-ON", Inches(0.45), Inches(0.1), Inches(2.5), Inches(0.4),
        sz=14, bold=True, color=CYBER_BLACK)
    _tb(slide, section_title, Inches(2.2), Inches(0.1), Inches(9), Inches(0.4),
        sz=13, color=CYBER_BLACK)

    cw = Inches(7.5) if img_label else R
    _card(slide, L, Inches(0.78), cw, Inches(6.27), TERMINAL_GREEN)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': WHITE, 'sb': 10}
             for i, s in enumerate(steps)]
    _paras(slide, items, L + Inches(0.2), Inches(0.98), cw - Inches(0.3), Inches(5.85))

    if img_label:
        _img_placeholder(slide, Inches(8.15), Inches(0.78), Inches(4.73), Inches(6.27),
                         img_label, accent=TERMINAL_GREEN)


def add_section_summary(prs, section_title, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "SECTION SUMMARY", section_title)
    _footer(slide)

    _card(slide, L, Inches(1.55), R, Inches(5.5), ELECTRIC_BLUE)
    items = [{'text': f"  ›  {b}", 'sz': 16, 'color': WHITE, 'sb': 14}
             for b in bullets]
    _paras(slide, items, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_summary_table(prs, title, subtitle, rows, col_headers):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, title, subtitle)
    _footer(slide)

    col_x = [L, Inches(3.5), Inches(8.1)]
    col_w = [Inches(2.9), Inches(4.4), Inches(5.0)]
    ht = Inches(1.52)
    rh = Inches(0.42)

    hdr = _rect(slide, L, ht, R, rh, ELECTRIC_BLUE)
    _glow_blue(hdr, 4)
    for i, h in enumerate(col_headers):
        _tb(slide, h, col_x[i] + Inches(0.1), ht + Inches(0.06), col_w[i], rh,
            sz=11, bold=True, color=CYBER_BLACK)

    for ri, row in enumerate(rows):
        top = ht + rh + ri * Inches(0.83)
        _rect(slide, L, top, R, Inches(0.83), DARK_NAVY if ri % 2 == 0 else NAVY_MID)
        bar = _rect(slide, L, top, Pt(3), Inches(0.83), ELECTRIC_BLUE)
        _glow_blue(bar, 3)
        colors = [ELECTRIC_BLUE, WHITE, LIGHT_GRAY]
        for ci, (val, color) in enumerate(zip(row, colors)):
            _tb(slide, val, col_x[ci] + Inches(0.12), top + Inches(0.08),
                col_w[ci] - Inches(0.15), Inches(0.8),
                sz=11, bold=(ci == 0), color=color)


def add_readiness_check(prs, items):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _header(slide, "READINESS CHECK", "Confirm before moving on", accent=CYBER_GOLD)
    _footer(slide, accent=CYBER_GOLD)

    _card(slide, L, Inches(1.55), R, Inches(5.5), CYBER_GOLD)
    paras = [{'text': f"  ☐   {item}", 'sz': 13, 'color': LIGHT_GRAY, 'sb': 8}
             for item in items]
    _paras(slide, paras, L + Inches(0.2), Inches(1.73), R - Inches(0.4), Inches(5.1))


def add_end_slide(prs, module_title, next_steps):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _bg(slide, CYBER_BLACK)
    _footer(slide)

    vert = _rect(slide, Inches(0), Inches(0), Pt(18), SLIDE_H, ELECTRIC_BLUE)
    _glow_blue(vert, 5)

    _tb(slide, "END OF MODULE", L, Inches(0.6), Inches(8), Inches(0.38),
        sz=11, bold=True, color=ELECTRIC_BLUE)
    _tb(slide, module_title, L, Inches(1.1), Inches(8), Inches(2.0),
        sz=40, bold=True, color=WHITE)
    _line(slide, L, Inches(3.25), Inches(6), CYBER_GOLD, 2, 3, 22)
    _tb(slide, "NEXT STEPS", L, Inches(3.45), Inches(4), Inches(0.35),
        sz=11, bold=True, color=CYBER_GOLD)
    items = [{'text': f"  {i+1}.   {s}", 'sz': 14, 'color': LIGHT_GRAY, 'sb': 10}
             for i, s in enumerate(next_steps)]
    _paras(slide, items, L, Inches(3.9), Inches(7.5), Inches(2.8))


# ---------------------------------------------------------------------------
# Build
# ---------------------------------------------------------------------------

def build():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    # --- TITLE ---
    add_title_slide(prs,
        "Terminal Basics:\nThe Machine",
        "PREREQUISITE MODULE  |  TERMINAL BASICS",
        "[IMAGE: Soldier at a command post terminal reviewing folder structures on dual monitors]")

    # --- OVERVIEW ---
    add_overview_slide(prs,
        "Everything on a computer lives in a file, files live in folders, and understanding that structure — "
        "plus the difference between plaintext and rich text — is the ground everything else stands on.",
        "Self-paced (approximately 45 minutes)",
        [
            "Navigate the filesystem without using search",
            "Read a file path as a route through the folder tree",
            "Explain the difference between plaintext and rich text and why it matters",
            "Open and view files in a code editor without changing anything",
            "Enable file extensions and understand why they are critical",
        ])

    # =========================================================================
    # SECTION 1 — Files, Folders, and the Tree
    # =========================================================================
    add_section_header(prs, "01", "Files, Folders,\nand the Tree",
        "Every file has exactly one location in the tree. A path is the route to reach it — apps have been hiding this from you.",
        "[IMAGE: Soldier reading a detailed map of an area of operations — every grid square has an address]")

    add_concept_slide(prs, "The Structure Under the Apps", [
        ("A file holds data: a document, an image, a script.", False),
        ("A folder (directory) holds files and other folders.", False),
        ("Folders nest inside folders, forming a tree.", False),
        ("Every file on the machine has exactly one location in that tree.", False),
        ("Most apps hide this — when you save in Word, Word decides where it goes.", False),
        ("The tree exists whether you see it or not.", False),
    ], note="The terminal modules all depend on being able to think about where things are. A student who cannot read a path gets lost the moment they type their first command.")

    add_concept_slide(prs, "What Is a Path?", [
        ("The path is the route through the tree to reach a file.", False),
        ("Windows example: C:\\Users\\YourName\\Documents\\report.txt", False),
        ("Mac/Linux example: /home/yourname/documents/report.txt", False),
        ("Read it left to right — each backslash or forward slash is one step down the tree.", False),
        ("The address bar in File Explorer or Finder shows the path — read it.", False),
    ], note="Walking a path by clicking in File Explorer is the same moves you will type in the terminal. This is not a new concept — it is the same concept made visible.")

    add_example_slide(prs, "Example: Finding a File You Saved Last Week",
        "Scenario: Navigating to a downloaded file without using search",
        [
            "Open File Explorer (Windows) or Finder (Mac).",
            "Navigate to your Downloads folder — click, do not search.",
            "Find a file you downloaded recently.",
            "",
            "Look at the address bar at the top.",
            "  That string of words separated by slashes is the path.",
            "  It is the route through the tree to reach that exact file.",
            "",
            "You have been walking this tree every time you opened a folder.",
            "  The terminal is the same moves — typed instead of clicked.",
        ],
        "[IMAGE: Soldier navigating a grid map by coordinates — each step a deliberate move through known terrain]")

    add_concept_slide(prs, "Turn On File Extensions Now", [
        ("By default, Windows hides file extensions (.docx, .txt, .pdf after the name).", False),
        ("You need to see extensions throughout this course.", False),
        ("Turn them on: File Explorer → View → Show → File name extensions.", False),
        ("Why it matters — the extension tells you what the file is:", False),
        ("report.txt — plaintext. report.docx — rich text. script.py — Python.", True),
        ("Two files named report with different extensions are fundamentally different things.", True),
    ], note="Turning extensions on here prevents confusion in the next section where the difference between .txt and .docx is the whole lesson. Do not skip this step.")

    add_hands_on(prs, "Files, Folders, and the Tree", [
        "Open File Explorer (Windows) or Finder (Mac).",
        "Navigate from your home folder down to Downloads by clicking — do not use search.",
        "Find a file. Read its full path in the address bar.",
        "Create a new folder on your Desktop. Name it 'practice'.",
        "Move any file into that folder by dragging.",
    ], "[IMAGE: Soldier navigating a tactical map to a specific grid — deliberate movement, no shortcuts]")

    add_check_on_learning(prs,
        "You just moved a file by dragging it into the 'practice' folder on your Desktop.\n\n"
        "Write out the full path to where it now lives.\n\n"
        "Could you give someone else those directions — as a path string — and have them find it without "
        "opening File Explorer? If not, what part of the path are you missing?")

    add_section_summary(prs, "Files, Folders, and the Tree", [
        "Every file lives in exactly one location in the folder tree.",
        "A path is the route through that tree — read it left to right, one slash at a time.",
        "Apps hide the tree. The terminal exposes it. Both move through the same structure.",
        "File extensions must be visible. Enable them now — every module depends on this.",
    ])

    add_readiness_check(prs, [
        "I can describe a file, a folder, and how they nest into a tree",
        "I can read a path as a route through that tree",
        "I have navigated to a specific file without using search",
        "File extensions are now visible on my machine",
    ])

    # =========================================================================
    # SECTION 2 — Plaintext vs Rich Text
    # =========================================================================
    add_section_header(prs, "02", "Plaintext vs\nRich Text",
        "A plaintext file is just characters — no hidden formatting. That is exactly what models and tools want.",
        "[IMAGE: Soldier comparing a clean printed report to a marked-up draft with sticky notes and margin edits]")

    add_concept_slide(prs, "What Is Plaintext?", [
        ("Plaintext is exactly what it says: text, nothing else.", False),
        ("Open a .txt file in any program on any machine — you see every character it contains.", False),
        ("No hidden formatting. No metadata. What you see is what the machine reads.", False),
        ("Common plaintext extensions: .txt, .md, .csv, .json, .py", False),
        ("The terminal, AI tools, and version control all work natively with plaintext.", False),
    ], note="Plaintext is honest — there are no hidden characters. Rich text hides things. That is the entire lesson at this level.")

    add_concept_slide(prs, "What Is Rich Text?", [
        ("Rich text files store formatting instructions alongside the actual text.", False),
        ("Bold, font size, margins, revision history — all packed in alongside the words.", False),
        ("Common rich text extensions: .docx, .xlsx, .pdf", False),
        ("When you paste from Word into a terminal or markdown document:", False),
        ("Smart quotes, em-dashes, and non-breaking spaces travel with the content.", True),
        ("Those hidden characters come along silently — the file looks fine but behaves wrong.", True),
        ("Hidden characters are the first thing to check when output looks correct but breaks.", False),
    ], note="Do not go into encodings (UTF-8, ASCII) at this level. 'Plaintext is honest, rich text hides things' is the entire lesson.")

    add_example_slide(prs, "Example: The Word-to-Markdown Trap",
        "Scenario: Pasting formatted content from Word into a Markdown file",
        [
            "Type a bullet point in Word using a hyphen.",
            "  Word automatically converts your straight quote to a 'smart quote'.",
            "  Word converts your double hyphen to an em-dash.",
            "",
            "Paste that into a Markdown file.",
            "  The Markdown parser reads the curly apostrophe as a plain character.",
            "  The em-dash breaks the bullet.",
            "  What looked fine in Word is broken in plaintext.",
            "",
            "The fix: write Markdown in a code editor, not Word.",
            "  Code editors work in plaintext — no automatic formatting changes.",
        ],
        "[IMAGE: Soldier comparing two versions of the same report — clean original vs. garbled retransmission]")

    add_hands_on(prs, "Plaintext vs Rich Text", [
        "Open Notepad (Windows) and type a few sentences. Save as 'test.txt'.",
        "Open the same file in VS Code. Compare how it looks.",
        "Now open any .docx file in VS Code (not Word). Look at what is actually inside it.",
        "You are seeing the difference between 'what a human reads' and 'what a machine reads.'",
    ], "[IMAGE: Soldier reading a signals intercept — what the radio received vs. what was intended]")

    add_check_on_learning(prs,
        "You paste a block of text from a Word document into your AI chatbot and the formatting looks off.\n\n"
        "The text appears correct on your screen, but the model is misreading parts of it.\n\n"
        "What is the most likely cause?\n\n"
        "What would you do instead of pasting from Word?")

    add_section_summary(prs, "Plaintext vs Rich Text", [
        "Plaintext: just characters — no hidden formatting. What you see is what the machine reads.",
        "Rich text: formatting instructions packed alongside the words. Hidden characters travel on paste.",
        "Models, terminals, and version control all expect plaintext.",
        "If something looks correct but behaves wrong — hidden characters are the first thing to check.",
    ])

    add_readiness_check(prs, [
        "I can distinguish plaintext from rich text and give an example of each",
        "I understand why pasting from Word can silently break things",
        "I know the common plaintext extensions (.txt, .md, .csv, .json, .py)",
        "I have seen what a .docx file actually looks like underneath in VS Code",
    ])

    # =========================================================================
    # SECTION 3 — Looking at Files in a Real Editor
    # =========================================================================
    add_section_header(prs, "03", "Looking at Files\nin a Real Editor",
        "Install a code editor, open a folder, read a file, close it without fear. That is the whole job today.",
        "[IMAGE: Soldier setting up a field workstation — equipment staged, ready for the mission, nothing extra]")

    add_concept_slide(prs, "What Is a Code Editor?", [
        ("A code editor is not an IDE — it is a text viewer with good defaults.", False),
        ("Syntax highlighting, a visible file tree, and preview mode for formatted files.", False),
        ("VS Code is the standard choice: free, runs on Windows and Mac, handles every file type in this course.", False),
        ("Three things to know today:", False),
        ("File tree panel (left side) — shows the folder structure, the same tree now labeled.", True),
        ("Editor panel (right side) — shows file content.", True),
        ("Markdown preview — click the split-square icon to render a .md file's formatting.", True),
    ], note="Today's scope: open, view, close. Do not install extensions, change settings, or start editing files. The goal is comfort at the front door, not a full tour.")

    add_concept_slide(prs, "Why Use a Code Editor?", [
        ("Code editors work in plaintext — no automatic formatting changes.", False),
        ("No smart quotes, no em-dashes, no hidden characters added on save.", False),
        ("The file tree panel gives you a visual map of the same folder structure the terminal navigates.", False),
        ("Markdown preview lets you see both the raw plaintext and the rendered output side by side.", False),
        ("The editor and the terminal are looking at the same files.", False),
        ("Learning to read files here makes every terminal command less abstract.", False),
    ])

    add_example_slide(prs, "Example: Rendering a Markdown File",
        "Scenario: Seeing plaintext and rendered output side by side",
        [
            "Open a .md file in VS Code.",
            "In the top-right corner, click the preview icon (or press Ctrl+Shift+V).",
            "",
            "Left pane: raw Markdown — exactly what the machine reads.",
            "  # Header becomes text with a pound sign.",
            "  **bold** appears as asterisks around a word.",
            "",
            "Right pane: rendered output — what a human reads.",
            "  The pound sign becomes a large heading.",
            "  The asterisks disappear and the word appears bold.",
            "",
            "Both panes are the same file. One view is for the machine. One is for you.",
        ],
        "[IMAGE: Soldier viewing the same intelligence product in two formats — raw signal and decoded message]")

    add_hands_on(prs, "Looking at Files in a Real Editor", [
        "Download and install VS Code if you have not already.",
        "Go to File → Open Folder. Open the 'practice' folder from Section 01.",
        "Click the .txt file you created. Read it in the editor panel.",
        "Find a .md file anywhere on your machine. Open it. Toggle the preview with Ctrl+Shift+V.",
    ], "[IMAGE: Soldier at a workstation opening a field report — reading it, not editing, just oriented]")

    add_check_on_learning(prs,
        "You are looking at the same Markdown file in two panes — raw on the left, rendered on the right.\n\n"
        "Which version does the machine work with?\n\n"
        "Which version does the human read?\n\n"
        "When you paste content from the rendered view into a terminal command — what are you actually sending?")

    add_section_summary(prs, "Looking at Files in a Real Editor", [
        "A code editor is a plaintext viewer — no hidden formatting changes on save or paste.",
        "The file tree panel shows the same folder structure you will navigate in the terminal.",
        "Markdown preview: raw plaintext on the left, rendered output on the right — same file.",
        "Today's scope: open, view, close. Nothing more.",
    ])

    add_readiness_check(prs, [
        "VS Code (or equivalent) is installed and opens without errors",
        "I can open a folder and navigate its file tree in the editor panel",
        "I have opened a Markdown file and toggled the rendered preview",
        "I did not change any settings or install any extensions",
    ])

    # =========================================================================
    # MODULE SUMMARY TABLE
    # =========================================================================
    add_summary_table(prs,
        "MODULE SUMMARY", "Terminal Basics — The Machine",
        [
            ("Files and Folders",  "Tree structure with paths",              "Every file has one address. A path is the route to reach it."),
            ("File Extensions",    "Type identifier after the dot",          "Must be visible. The extension tells you what the file is."),
            ("Plaintext",          "Characters only — no hidden formatting", "What you see is what the machine reads. Tools expect this."),
            ("Rich Text",          "Formatting packed alongside content",    "Pasting from Word brings hidden characters. Use a code editor."),
            ("Code Editor",        "Plaintext viewer with good defaults",    "File tree + editor + preview. Open, view, close — nothing more."),
        ],
        ["CONCEPT", "CORE IDEA", "WHY IT MATTERS"])

    # =========================================================================
    # FINAL READINESS CHECK
    # =========================================================================
    add_readiness_check(prs, [
        "I can describe a file, a folder, and the folder tree",
        "I can read a file path as a route through the tree",
        "I have navigated to a file without using search",
        "File extensions are visible on my machine",
        "I can explain the difference between plaintext and rich text",
        "I know which file types are plaintext and which are rich text",
        "VS Code is installed and I have opened a file and toggled Markdown preview",
    ])

    # =========================================================================
    # END SLIDE
    # =========================================================================
    add_end_slide(prs, "Terminal Basics:\nThe Machine", [
        "Navigate your filesystem by path — not by search — for one full week.",
        "Save every file you create in this course as plaintext (.txt, .md, .py) not .docx.",
        "When something looks right but behaves wrong — check for hidden characters first.",
        "Next module: The Terminal — same tree, typed instead of clicked.",
    ])

    out = os.path.join(os.path.dirname(__file__), "terminal-machine.pptx")
    prs.save(out)
    print(f"Saved: {out}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    build()
